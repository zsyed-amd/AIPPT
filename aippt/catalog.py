"""SQLite slide catalog with content hashing and versioning."""
import csv
import hashlib
import logging
import os
import re
import sqlite3
import tempfile
import threading
from datetime import datetime
from pathlib import Path
from typing import TYPE_CHECKING, List, Dict, Optional, Tuple

from pptx import Presentation

from aippt.reverse import extract_text_from_shape
from aippt.sections import read_sections

if TYPE_CHECKING:  # pragma: no cover - typing only
    from aippt.storage import Storage

logger = logging.getLogger(__name__)

# Object-storage key for the catalog snapshot (see snapshot_db/restore_db).
CATALOG_SNAPSHOT_KEY = "catalog/slides.db"

SCHEMA_PATH = Path(__file__).parent / "schema.sql"

_UUID_PREFIX_RE = re.compile(r'^[0-9a-f]{8}(?:[0-9a-f]{24})?_')


def display_name(name: str) -> str:
    """Strip a leading 8- or 32-char hex UUID prefix from a deck name.

    Uploaded decks get a 32-char UUID prefix and generated decks get an
    8-char short ID for collision safety.  This function returns the
    human-friendly name without that prefix.  Names that don't match
    the pattern are returned unchanged.
    """
    return _UUID_PREFIX_RE.sub('', name)


class _CatalogConnection(sqlite3.Connection):
    """A catalog connection whose successful commits trigger a debounced
    object-storage snapshot. The trigger is a no-op when no snapshot scheduler
    is installed (the default, filesystem mode), so committing behaves exactly
    as before unless an object-storage backend is active.
    """

    snapshot_on_commit = False

    def commit(self):
        super().commit()
        if self.snapshot_on_commit:
            request_snapshot()


def get_db(db_path: str = "slides.db") -> sqlite3.Connection:
    """Open database connection, create schema if needed."""
    conn = sqlite3.connect(db_path, timeout=30, factory=_CatalogConnection)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode = WAL")
    conn.execute("PRAGMA foreign_keys = ON")
    with open(SCHEMA_PATH, encoding="utf-8") as f:
        conn.executescript(f.read())

    # Idempotent migrations: add new columns to pre-existing databases.
    # PRAGMA table_info returns one row per column; we collect existing names
    # and only issue ALTER TABLE when the column is absent.
    existing_deck_cols = {
        row[1]
        for row in conn.execute("PRAGMA table_info(decks)").fetchall()
    }
    for col_ddl, col_name in (
        ("author TEXT NOT NULL DEFAULT ''", "author"),
        ("created_date TEXT DEFAULT NULL", "created_date"),
        ("modified_date TEXT DEFAULT NULL", "modified_date"),
        ("subject TEXT NOT NULL DEFAULT ''", "subject"),
        ("description TEXT NOT NULL DEFAULT ''", "description"),
        ("source_script_path TEXT DEFAULT NULL", "source_script_path"),
        ("source_engine TEXT DEFAULT NULL", "source_engine"),
        ("source_theme TEXT DEFAULT NULL", "source_theme"),
        ("outline_path TEXT DEFAULT NULL", "outline_path"),
        ("source_generated_at TEXT DEFAULT NULL", "source_generated_at"),
    ):
        if col_name not in existing_deck_cols:
            conn.execute(f"ALTER TABLE decks ADD COLUMN {col_ddl}")
            logger.debug("Migration: added decks.%s", col_name)

    existing_slide_cols = {
        row[1]
        for row in conn.execute("PRAGMA table_info(slides)").fetchall()
    }
    for col_ddl, col_name in (
        ("author TEXT NOT NULL DEFAULT ''", "author"),
        ("slide_created_date TEXT DEFAULT NULL", "slide_created_date"),
        ("layout_type TEXT DEFAULT NULL", "layout_type"),
    ):
        if col_name not in existing_slide_cols:
            conn.execute(f"ALTER TABLE slides ADD COLUMN {col_ddl}")
            logger.debug("Migration: added slides.%s", col_name)

    conn.commit()
    # Schema setup is done committing; arm snapshot-on-commit so only real
    # catalog mutations from here on trigger a debounced snapshot.
    conn.snapshot_on_commit = True
    return conn


# ---------------------------------------------------------------------------
# Object-storage snapshot / restore for the SQLite catalog
#
# SQLite cannot run directly off object storage, so the catalog lives on a
# local (ephemeral) volume and is snapshotted to / restored from object storage
# as a single consistent file. Snapshots use SQLite's online backup API, which
# yields a WAL-safe image even with writes in flight (superior to copying the
# .db file, which can capture a torn WAL state). Validated end-to-end in the
# 2026-06-11 spike (upload -> cold-restore -> PRAGMA integrity_check = ok).
# ---------------------------------------------------------------------------


def snapshot_db(local_path: str, storage: "Storage", key: str = CATALOG_SNAPSHOT_KEY) -> None:
    """Upload a consistent snapshot of the catalog at *local_path* to *storage*.

    Uses the online backup API so the snapshot is consistent even if WAL writes
    are in flight. No-op-safe to call repeatedly (overwrites the same key).
    """
    src = sqlite3.connect(local_path, timeout=30)
    try:
        with tempfile.TemporaryDirectory() as td:
            snap = os.path.join(td, "snapshot.db")
            dst = sqlite3.connect(snap)
            try:
                with dst:
                    src.backup(dst)
            finally:
                dst.close()
            with open(snap, "rb") as fh:
                storage.put(key, fh, content_type="application/octet-stream")
    finally:
        src.close()
    logger.info("Catalog snapshot uploaded to %s", key)


def restore_db(local_path: str, storage: "Storage", key: str = CATALOG_SNAPSHOT_KEY) -> bool:
    """Restore the catalog snapshot from *storage* into *local_path*.

    Returns True if a snapshot was found and restored, False if no snapshot
    exists yet (cold start with an empty store -- caller starts fresh). Any
    stale WAL/SHM sidecars next to *local_path* are removed first so a restored
    image is never mixed with a previous incarnation's journal.
    """
    if not storage.exists(key):
        logger.info("No catalog snapshot at %s; starting empty", key)
        return False
    os.makedirs(os.path.dirname(os.path.abspath(local_path)), exist_ok=True)
    for sidecar in (local_path + "-wal", local_path + "-shm"):
        try:
            os.remove(sidecar)
        except FileNotFoundError:
            pass
    data = storage.get(key)
    with open(local_path, "wb") as fh:
        fh.write(data)
    logger.info("Catalog restored from %s (%d bytes)", key, len(data))
    return True


class SnapshotScheduler:
    """Debounce catalog snapshots so a burst of writes coalesces into one push.

    Single-writer by design: the production deployment runs ``replicas: 1``, so
    there is no multi-writer race. Each ``request()`` (re)arms a timer; the
    snapshot fires once the writes go quiet for ``debounce_seconds``. ``flush``
    forces an immediate snapshot if one is pending (used on shutdown and in
    tests); ``shutdown`` cancels any pending timer.
    """

    def __init__(
        self,
        db_path: str,
        storage: "Storage",
        key: str = CATALOG_SNAPSHOT_KEY,
        debounce_seconds: float = 5.0,
    ):
        self.db_path = db_path
        self.storage = storage
        self.key = key
        self.debounce_seconds = debounce_seconds
        self._lock = threading.Lock()
        self._timer: Optional[threading.Timer] = None
        self._pending = False

    def request(self) -> None:
        with self._lock:
            self._pending = True
            if self._timer is not None:
                self._timer.cancel()
            self._timer = threading.Timer(self.debounce_seconds, self._fire)
            self._timer.daemon = True
            self._timer.start()

    def _fire(self) -> None:
        with self._lock:
            self._timer = None
            if not self._pending:
                return
            self._pending = False
        try:
            snapshot_db(self.db_path, self.storage, self.key)
        except Exception:  # pragma: no cover - background thread guard
            logger.exception("Debounced catalog snapshot failed")

    def flush(self) -> None:
        """Snapshot immediately if a snapshot is pending."""
        with self._lock:
            if self._timer is not None:
                self._timer.cancel()
                self._timer = None
            pending = self._pending
            self._pending = False
        if pending:
            snapshot_db(self.db_path, self.storage, self.key)

    def shutdown(self) -> None:
        with self._lock:
            if self._timer is not None:
                self._timer.cancel()
                self._timer = None
            self._pending = False


# Process-global scheduler. ``None`` (the default, filesystem mode) makes
# ``request_snapshot`` inert, so the catalog write paths stay byte-for-byte
# identical to historical behavior unless an object-storage backend installs a
# scheduler at startup.
_snapshot_scheduler: Optional[SnapshotScheduler] = None


def set_snapshot_scheduler(scheduler: Optional[SnapshotScheduler]) -> None:
    """Install (or clear with ``None``) the process-global snapshot scheduler."""
    global _snapshot_scheduler
    _snapshot_scheduler = scheduler


def request_snapshot() -> None:
    """Request a debounced catalog snapshot. No-op when no scheduler is installed."""
    scheduler = _snapshot_scheduler
    if scheduler is not None:
        scheduler.request()


def content_hash(title: str, text: str) -> str:
    """SHA-256 hash of normalized title + content text."""
    normalized = f"{title.strip().lower()}\n{text.strip().lower()}"
    return hashlib.sha256(normalized.encode()).hexdigest()


def _resolve_slide_title(slide):
    """Return ``(title, fallback_shape)`` for a slide.

    Standard title placeholders win. When that's empty -- common for decks
    produced by PptxGenJS, which lays the title out as a regular text box
    rather than the title placeholder -- fall back to the first short
    (<=80 char) text shape's first line. The returned ``fallback_shape`` is
    the shape the title came from (or ``None`` if the title placeholder was
    used), so the caller can skip it when building ``content_text`` and
    avoid duplicating the title into the body.
    """
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip(), None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        first_line = shape.text_frame.paragraphs[0].text.strip()
        if first_line and len(first_line) <= 80:
            return first_line, shape

    return "", None


def file_hash(file_path: str) -> str:
    """SHA-256 hash of file contents."""
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def detect_source_engine(script_path: str) -> Optional[str]:
    """Auto-detect engine from script content.

    Scans the first 50 lines for pptxgenjs or python-pptx imports.

    Returns:
        'pptxgenjs', 'python-pptx', or None if undetectable.
    """
    try:
        with open(script_path, encoding="utf-8") as f:
            for i, line in enumerate(f):
                if i >= 50:
                    break
                if "require('pptxgenjs')" in line or 'require("pptxgenjs")' in line:
                    return "pptxgenjs"
                if "from pptx import" in line or "import pptx" in line:
                    return "python-pptx"
    except (OSError, UnicodeDecodeError):
        pass
    return None


def detect_source_theme(script_path: str) -> Optional[str]:
    """Auto-detect theme name from script content.

    Looks for references to ``themes/<name>.yaml``.

    Returns:
        Theme name (e.g. 'amd', 'default') or None.
    """
    theme_re = re.compile(r"themes/(\w+)\.ya?ml")
    try:
        with open(script_path, encoding="utf-8") as f:
            for line in f:
                m = theme_re.search(line)
                if m:
                    return m.group(1)
    except (OSError, UnicodeDecodeError):
        pass
    return None


def catalog_deck(
    deck_path: str,
    db_path: str = "slides.db",
    images_dir: Optional[str] = None,
    base_dir: Optional[str] = None,
    source_script_path: Optional[str] = None,
    source_engine: Optional[str] = None,
    source_theme: Optional[str] = None,
    outline_path: Optional[str] = None,
) -> int:
    """Catalog a PowerPoint deck into the database.

    Args:
        deck_path: Path to the PowerPoint file
        db_path: Path to the SQLite database
        images_dir: Optional directory containing slide images
        base_dir: Base directory for relative path storage (default: cwd)
        source_script_path: Path to the generating script (JS/Python)
        source_engine: Engine used ('pptxgenjs' or 'python-pptx')
        source_theme: Theme name ('amd', 'default', etc.)
        outline_path: Path to the originating markdown outline

    Returns:
        The deck ID
    """
    conn = get_db(db_path)
    base = base_dir or os.getcwd()
    deck_path_abs = os.path.abspath(deck_path)
    deck_path = os.path.relpath(deck_path_abs, base)
    deck_name = os.path.splitext(os.path.basename(deck_path))[0]
    fhash = file_hash(deck_path_abs)

    # Check if this exact file is already cataloged
    existing = conn.execute(
        "SELECT id FROM decks WHERE file_path = ? AND file_hash = ?",
        (deck_path, fhash),
    ).fetchone()
    if existing:
        logger.info(f"Deck already cataloged with same hash: {deck_name}")
        conn.close()
        return existing["id"]

    prs = Presentation(deck_path_abs)

    # Extract core_properties metadata
    cp = prs.core_properties
    deck_author: str = (cp.author or "").strip()
    deck_subject: str = (cp.subject or "").strip()
    # python-pptx exposes 'comments' rather than 'description';
    # fall back gracefully if a future version adds 'description'.
    deck_description: str = (
        getattr(cp, "description", None) or cp.comments or ""
    ).strip()

    _cp_created = cp.created
    if _cp_created is not None:
        deck_created_date: Optional[str] = _cp_created.isoformat()
    else:
        # Fall back to file modification time; if that also fails, use None
        try:
            deck_created_date = datetime.fromtimestamp(
                os.path.getmtime(deck_path_abs)
            ).isoformat()
        except OSError:
            deck_created_date = None

    _cp_modified = cp.modified
    deck_modified_date: Optional[str] = (
        _cp_modified.isoformat() if _cp_modified is not None else None
    )

    # Upsert deck record
    existing_deck = conn.execute(
        "SELECT id FROM decks WHERE file_path = ?", (deck_path,)
    ).fetchone()

    now_iso = datetime.utcnow().isoformat()

    if existing_deck:
        deck_id = existing_deck["id"]
        # Build UPDATE: always update core fields; source fields only when
        # explicitly provided (preserve existing values on re-catalog)
        update_sql = """UPDATE decks
               SET file_hash = ?, slide_count = ?, author = ?, modified_date = ?,
                   subject = ?, description = ?,
                   source_generated_at = ?,
                   updated_at = datetime('now')"""
        params = [fhash, len(prs.slides), deck_author, deck_modified_date,
                  deck_subject, deck_description, now_iso]
        if source_script_path is not None:
            update_sql += ", source_script_path = ?"
            params.append(source_script_path)
        if source_engine is not None:
            update_sql += ", source_engine = ?"
            params.append(source_engine)
        if source_theme is not None:
            update_sql += ", source_theme = ?"
            params.append(source_theme)
        if outline_path is not None:
            update_sql += ", outline_path = ?"
            params.append(outline_path)
        update_sql += " WHERE id = ?"
        params.append(deck_id)
        conn.execute(update_sql, params)
        # Remove old slides for re-catalog
        conn.execute("DELETE FROM slides WHERE deck_id = ?", (deck_id,))
    else:
        cur = conn.execute(
            """INSERT INTO decks (name, file_path, file_hash, slide_count,
                                  author, created_date, modified_date,
                                  subject, description,
                                  source_script_path, source_engine,
                                  source_theme, outline_path,
                                  source_generated_at)
               VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
            (deck_name, deck_path, fhash, len(prs.slides),
             deck_author, deck_created_date, deck_modified_date,
             deck_subject, deck_description,
             source_script_path, source_engine, source_theme,
             outline_path, now_iso if source_script_path else None),
        )
        deck_id = cur.lastrowid

    # Catalog each slide
    for i, slide in enumerate(prs.slides, 1):
        title, title_fallback_shape = _resolve_slide_title(slide)

        # Extract text from all non-title shapes. Skip the fallback shape
        # too (when used) so the title isn't duplicated into content_text.
        texts = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if title_fallback_shape is not None and shape == title_fallback_shape:
                continue
            text = extract_text_from_shape(shape)
            if text:
                texts.append(text)
        content_text = "\n".join(texts)

        chash = content_hash(title, content_text)

        # Notes
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        # Image path
        image_path = None
        if images_dir:
            for ext in (".png", ".PNG", ".jpg", ".jpeg"):
                candidate = os.path.join(images_dir, f"Slide{i}{ext}")
                if os.path.exists(candidate):
                    image_path = os.path.relpath(os.path.abspath(candidate), base)
                    break

        conn.execute(
            """INSERT INTO slides (deck_id, position, title, content_text, content_hash,
               notes, image_path, author, slide_created_date)
               VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)""",
            (deck_id, i, title, content_text, chash, notes, image_path,
             deck_author, deck_created_date),
        )

    # Read and catalog sections from PowerPoint
    ppt_sections = read_sections(prs)

    # Clear old section assignments for this deck
    conn.execute("DELETE FROM slide_sections WHERE slide_id IN "
                 "(SELECT id FROM slides WHERE deck_id = ?)", (deck_id,))
    conn.execute("DELETE FROM sections WHERE deck_id = ?", (deck_id,))

    # Insert sections and map to slides
    seen_names: dict[str, int] = {}
    for pos, section in enumerate(ppt_sections, 1):
        # Deduplicate section names (PPTX can have multiple sections with the same name)
        name = section.name
        if name in seen_names:
            seen_names[name] += 1
            name = f"{name} ({seen_names[name]})"
        else:
            seen_names[name] = 1
        sec_result = conn.execute(
            "INSERT INTO sections (deck_id, name, position) VALUES (?, ?, ?)",
            (deck_id, name, pos)
        )
        section_id = sec_result.lastrowid

        # Map section slide IDs to database slide IDs via position
        for pptx_slide_id in section.slide_ids:
            # Find slide in database by matching to presentation slide order
            for i, prs_slide in enumerate(prs.slides, 1):
                if prs_slide.slide_id == pptx_slide_id:
                    db_slide = conn.execute(
                        "SELECT id FROM slides WHERE deck_id = ? AND position = ?",
                        (deck_id, i)
                    ).fetchone()
                    if db_slide:
                        conn.execute(
                            "INSERT INTO slide_sections (slide_id, section_id) VALUES (?, ?)",
                            (db_slide["id"], section_id)
                        )
                    break

    if ppt_sections:
        logger.info(f"Cataloged {len(ppt_sections)} sections from {deck_name}")

    conn.commit()
    logger.info(f"Cataloged {len(prs.slides)} slides from {deck_name}")
    conn.close()
    return deck_id


def migrate_paths(
    db_path: str = "slides.db",
    base_dir: Optional[str] = None,
) -> Dict[str, int]:
    """Convert absolute paths in the database to relative paths.

    Rewrites ``decks.file_path`` and ``slides.image_path`` from absolute to
    relative (relative to *base_dir*). Already-relative paths are skipped.
    The migration is idempotent — running it twice is safe.

    Args:
        db_path: Path to the SQLite database
        base_dir: Base directory for computing relative paths (default: cwd)

    Returns:
        Dict with counts: ``deck_paths``, ``image_paths``, ``skipped``
    """
    base = base_dir or os.getcwd()
    conn = get_db(db_path)

    deck_count = 0
    image_count = 0
    skipped = 0

    # Migrate decks.file_path
    rows = conn.execute("SELECT id, file_path FROM decks").fetchall()
    for row in rows:
        fp = row["file_path"]
        if not fp or not os.path.isabs(fp):
            skipped += 1
            continue
        rel = os.path.relpath(fp, base)
        conn.execute("UPDATE decks SET file_path = ? WHERE id = ?", (rel, row["id"]))
        deck_count += 1

    # Migrate slides.image_path
    rows = conn.execute("SELECT id, image_path FROM slides WHERE image_path IS NOT NULL").fetchall()
    for row in rows:
        ip = row["image_path"]
        if not ip or not os.path.isabs(ip):
            skipped += 1
            continue
        rel = os.path.relpath(ip, base)
        conn.execute("UPDATE slides SET image_path = ? WHERE id = ?", (rel, row["id"]))
        image_count += 1

    conn.commit()
    conn.close()
    return {"deck_paths": deck_count, "image_paths": image_count, "skipped": skipped}


def search_slides(
    db_path: str = "slides.db",
    tags: Optional[List[str]] = None,
    title_contains: Optional[str] = None,
    section: Optional[str] = None,
) -> List[Dict]:
    """Search slides by tags, title substring, and/or section.

    Args:
        db_path: Path to the SQLite database
        tags: List of tag names to filter by (AND logic)
        title_contains: Substring to search in titles
        section: Section name to filter by (substring match)

    Returns:
        List of slide dictionaries
    """
    conn = get_db(db_path)
    query = """
        SELECT s.id, s.position, s.title, s.content_hash, s.image_path,
               d.name as deck_name, d.file_path as deck_path,
               s.updated_at
        FROM slides s
        JOIN decks d ON s.deck_id = d.id
        WHERE 1=1
    """
    params = []

    if title_contains:
        query += " AND s.title LIKE ?"
        params.append(f"%{title_contains}%")

    if tags:
        placeholders = ",".join("?" * len(tags))
        query += f"""
            AND s.id IN (
                SELECT st.slide_id FROM slide_tags st
                JOIN tags t ON st.tag_id = t.id
                WHERE t.name IN ({placeholders})
            )
        """
        params.extend(tags)

    if section:
        query += """ AND s.id IN (
            SELECT ss.slide_id FROM slide_sections ss
            JOIN sections sec ON ss.section_id = sec.id
            WHERE sec.name LIKE ?
        )"""
        params.append(f"%{section}%")

    query += " ORDER BY d.name, s.position"
    rows = conn.execute(query, params).fetchall()
    conn.close()
    return [dict(r) for r in rows]


def check_newer_versions(
    slides: List[Dict],
    db_path: str = "slides.db",
) -> List[Dict]:
    """Check if any slides have newer versions in other decks.

    Args:
        slides: List of slide dictionaries with 'title', 'content_hash', 'updated_at'
        db_path: Path to the SQLite database

    Returns:
        List of warning dictionaries
    """
    conn = get_db(db_path)
    warnings = []

    for slide in slides:
        # Find slides with same title but different hash, newer timestamp
        rows = conn.execute(
            """SELECT s.title, s.content_hash, s.updated_at, d.name as deck_name
               FROM slides s JOIN decks d ON s.deck_id = d.id
               WHERE s.title = ? AND s.content_hash != ? AND s.updated_at > ?
               ORDER BY s.updated_at DESC LIMIT 1""",
            (slide.get("title", ""), slide.get("content_hash", ""), slide.get("updated_at", "")),
        ).fetchall()

        for row in rows:
            warnings.append({
                "slide_title": slide.get("title", ""),
                "current_deck": slide.get("deck_name", ""),
                "newer_deck": row["deck_name"],
                "newer_updated": row["updated_at"],
            })

    conn.close()
    return warnings


def get_all_tags(db_path: str = "slides.db") -> List[Dict]:
    """Return all tags that are applied to at least one slide, with counts and categories.

    Each dict has keys: name, category, count.
    Category comes from the taxonomy table (empty string if not in taxonomy).
    Tags with zero slide associations are excluded.

    Args:
        db_path: Path to the SQLite database

    Returns:
        List of dicts sorted by category then name
    """
    conn = get_db(db_path)
    rows = conn.execute(
        """SELECT t.name, COALESCE(tx.category, '') AS category, COUNT(st.slide_id) AS count
           FROM tags t
           JOIN slide_tags st ON t.id = st.tag_id
           LEFT JOIN taxonomy tx ON t.name = tx.name
           GROUP BY t.id
           ORDER BY category, t.name"""
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


def add_tags(
    slide_id: int,
    tag_names: List[str],
    source: str = "ai",
    db_path: str = "slides.db",
):
    """Add tags to a slide.

    Args:
        slide_id: The slide ID
        tag_names: List of tag names to add
        source: Tag source ('ai', 'taxonomy', 'manual')
        db_path: Path to the SQLite database
    """
    conn = get_db(db_path)
    for name in tag_names:
        name = name.strip().lower()
        if not name:
            continue
        # Upsert tag
        conn.execute(
            "INSERT OR IGNORE INTO tags (name, source) VALUES (?, ?)",
            (name, source),
        )
        tag_row = conn.execute("SELECT id FROM tags WHERE name = ?", (name,)).fetchone()
        if tag_row:
            conn.execute(
                "INSERT OR IGNORE INTO slide_tags (slide_id, tag_id) VALUES (?, ?)",
                (slide_id, tag_row["id"]),
            )
    conn.commit()
    conn.close()


def get_slide_tags(slide_id: int, db_path: str = "slides.db") -> List[str]:
    """Get all tags for a slide.

    Args:
        slide_id: The slide ID
        db_path: Path to the SQLite database

    Returns:
        List of tag names
    """
    conn = get_db(db_path)
    rows = conn.execute(
        """SELECT t.name FROM tags t
           JOIN slide_tags st ON t.id = st.tag_id
           WHERE st.slide_id = ?
           ORDER BY t.name""",
        (slide_id,),
    ).fetchall()
    conn.close()
    return [r["name"] for r in rows]


def get_deck_by_id(deck_id: int, db_path: str = "slides.db") -> Optional[Dict]:
    """Get deck information by ID.

    Args:
        deck_id: The deck ID
        db_path: Path to the SQLite database

    Returns:
        Deck dictionary or None
    """
    conn = get_db(db_path)
    row = conn.execute(
        "SELECT id, name, file_path, file_hash, slide_count, author, cataloged_at, updated_at, subject, description, source_script_path, source_engine, source_theme, outline_path, source_generated_at FROM decks WHERE id = ?",
        (deck_id,),
    ).fetchone()
    conn.close()
    return dict(row) if row else None


def get_deck_by_file_hash(sha256: str, db_path: str = "slides.db") -> Optional[Dict]:
    """Return deck metadata for the deck whose ``file_hash`` matches ``sha256``.

    Used by the SPA upload-pre-check endpoint to detect duplicates *before*
    the user re-uploads a deck that's already cataloged. ``decks.file_hash``
    has an index (see ``schema.sql``) so this is O(log n).

    Args:
        sha256: Lowercase 64-character hex SHA-256 of the file's raw bytes.
        db_path: Path to the SQLite database.

    Returns:
        A dict of deck metadata, or None if no deck has that hash. The
        returned shape mirrors ``get_deck_by_id``.
    """
    conn = get_db(db_path)
    row = conn.execute(
        "SELECT id, name, file_path, file_hash, slide_count, author, "
        "cataloged_at, updated_at, subject, description, source_script_path, "
        "source_engine, source_theme, outline_path, source_generated_at "
        "FROM decks WHERE file_hash = ? LIMIT 1",
        (sha256,),
    ).fetchone()
    conn.close()
    return dict(row) if row else None


def resolve_deck(identifier: str, db_path: str = "slides.db"):
    """Look up a deck by ID (integer) or name (case-insensitive partial match).

    Args:
        identifier: Integer ID string or name substring
        db_path: Path to the SQLite database

    Returns:
        - dict: Single matching deck
        - list: Multiple matching decks (ambiguous)
        - None: No match found
    """
    conn = get_db(db_path)
    _deck_cols = "id, name, file_path, file_hash, slide_count, author, subject, description, cataloged_at, updated_at, source_script_path, source_engine, source_theme, outline_path, source_generated_at"
    # Try integer ID first
    try:
        deck_id = int(identifier)
        row = conn.execute(
            f"SELECT {_deck_cols} FROM decks WHERE id = ?",
            (deck_id,),
        ).fetchone()
        conn.close()
        return dict(row) if row else None
    except ValueError:
        pass

    # Partial name match (case-insensitive)
    rows = conn.execute(
        f"SELECT {_deck_cols} FROM decks WHERE name LIKE ? COLLATE NOCASE",
        (f"%{identifier}%",),
    ).fetchall()
    conn.close()

    if len(rows) == 1:
        return dict(rows[0])
    if len(rows) > 1:
        return [dict(r) for r in rows]
    return None


def delete_deck(deck_id: int, db_path: str = "slides.db") -> Optional[Dict]:
    """Delete a deck and all associated data (cascade).

    Returns:
        Dict with deleted deck info (name, slide_count, tag_count, section_count),
        or None if deck not found
    """
    conn = get_db(db_path)
    row = conn.execute(
        "SELECT id, name, slide_count FROM decks WHERE id = ?", (deck_id,)
    ).fetchone()
    if not row:
        conn.close()
        return None

    info = dict(row)
    info["tag_count"] = conn.execute(
        "SELECT COUNT(*) FROM slide_tags WHERE slide_id IN (SELECT id FROM slides WHERE deck_id = ?)",
        (deck_id,),
    ).fetchone()[0]
    info["section_count"] = conn.execute(
        "SELECT COUNT(*) FROM sections WHERE deck_id = ?", (deck_id,),
    ).fetchone()[0]

    conn.execute("DELETE FROM decks WHERE id = ?", (deck_id,))
    conn.commit()
    conn.close()
    return info


def rename_deck(deck_id: int, new_name: str, db_path: str = "slides.db") -> Optional[str]:
    """Rename a deck. Returns old name if renamed, None if not found."""
    conn = get_db(db_path)
    row = conn.execute("SELECT name FROM decks WHERE id = ?", (deck_id,)).fetchone()
    if not row:
        conn.close()
        return None
    old_name = row["name"]
    conn.execute("UPDATE decks SET name = ? WHERE id = ?", (new_name, deck_id))
    conn.commit()
    conn.close()
    return old_name


def get_deck_tag_count(deck_id: int, db_path: str = "slides.db") -> int:
    """Count distinct tags across all slides in a deck."""
    conn = get_db(db_path)
    row = conn.execute(
        "SELECT COUNT(DISTINCT st.tag_id) FROM slide_tags st JOIN slides s ON st.slide_id = s.id WHERE s.deck_id = ?",
        (deck_id,),
    ).fetchone()
    conn.close()
    return row[0]


def get_deck_top_tags(deck_id: int, db_path: str = "slides.db", limit: int = 10) -> List[Tuple[str, int]]:
    """Get the most common tags in a deck, ordered by frequency."""
    conn = get_db(db_path)
    rows = conn.execute(
        """SELECT t.name, COUNT(st.slide_id) as cnt
           FROM slide_tags st
           JOIN tags t ON st.tag_id = t.id
           JOIN slides s ON st.slide_id = s.id
           WHERE s.deck_id = ?
           GROUP BY t.id
           ORDER BY cnt DESC, t.name
           LIMIT ?""",
        (deck_id, limit),
    ).fetchall()
    conn.close()
    return [(r["name"], r["cnt"]) for r in rows]


def get_deck_slides(deck_id: int, db_path: str = "slides.db") -> List[Dict]:
    """Get all slides for a deck.

    Args:
        deck_id: The deck ID
        db_path: Path to the SQLite database

    Returns:
        List of slide dictionaries
    """
    conn = get_db(db_path)
    rows = conn.execute(
        """SELECT id, position, title, content_text, content_hash, notes, image_path, created_at, updated_at, layout_type
           FROM slides WHERE deck_id = ? ORDER BY position""",
        (deck_id,),
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


def list_decks(db_path: str = "slides.db") -> List[Dict]:
    """List all cataloged decks.

    Args:
        db_path: Path to the SQLite database

    Returns:
        List of deck dictionaries
    """
    conn = get_db(db_path)
    rows = conn.execute(
        "SELECT id, name, file_path, slide_count, author, cataloged_at, updated_at, subject, description FROM decks ORDER BY updated_at DESC"
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


# ---------------------------------------------------------------------------
# Tag removal
# ---------------------------------------------------------------------------


def get_deck_origin(deck_id: int, db_path: str = "slides.db") -> dict:
    """Return the derived origin block for a deck.

    Reads the 5 origin columns and computes ``kind`` from their state:

    * ``outline_path`` set, ``source_script_path`` NULL  →  ``"outline"``
    * ``source_script_path`` set (regardless of outline)  →  ``"script"``
    * Both NULL  →  ``"upload"``

    Args:
        deck_id: The deck ID.
        db_path: Path to the SQLite database.

    Returns:
        Dict with keys: ``kind``, ``outline_path``, ``source_script_path``,
        ``engine``, ``theme``, ``generated_at``.  Returns ``None`` if deck
        not found.
    """
    conn = get_db(db_path)
    row = conn.execute(
        "SELECT outline_path, source_script_path, source_engine, "
        "source_theme, source_generated_at FROM decks WHERE id = ?",
        (deck_id,),
    ).fetchone()
    conn.close()
    if row is None:
        return None

    outline_path = row["outline_path"]
    script_path = row["source_script_path"]

    if script_path:
        kind = "script"
    elif outline_path:
        kind = "outline"
    else:
        kind = "upload"

    return {
        "kind": kind,
        "outline_path": outline_path,
        "source_script_path": script_path,
        "engine": row["source_engine"],
        "theme": row["source_theme"],
        "generated_at": row["source_generated_at"],
    }


def remove_slide_tag(
    slide_id: int,
    tag_name: str,
    db_path: str = "slides.db",
) -> bool:
    """Remove a specific tag from a slide.

    Args:
        slide_id: The slide ID
        tag_name: Tag name to remove
        db_path: Path to the SQLite database

    Returns:
        True if a tag was removed, False if it wasn't found
    """
    conn = get_db(db_path)
    tag_name = tag_name.strip().lower()
    tag_row = conn.execute("SELECT id FROM tags WHERE name = ?", (tag_name,)).fetchone()
    if not tag_row:
        conn.close()
        return False
    cur = conn.execute(
        "DELETE FROM slide_tags WHERE slide_id = ? AND tag_id = ?",
        (slide_id, tag_row["id"]),
    )
    conn.commit()
    removed = cur.rowcount > 0
    conn.close()
    return removed


def remove_all_slide_tags(
    slide_id: int,
    db_path: str = "slides.db",
) -> int:
    """Remove all tags from a slide.

    Args:
        slide_id: The slide ID
        db_path: Path to the SQLite database

    Returns:
        Number of tags removed
    """
    conn = get_db(db_path)
    cur = conn.execute("DELETE FROM slide_tags WHERE slide_id = ?", (slide_id,))
    conn.commit()
    count = cur.rowcount
    conn.close()
    return count


def rename_tag(
    old_name: str,
    new_name: str,
    db_path: str = "slides.db",
) -> int:
    """Rename a tag everywhere (tags table + taxonomy table).

    Args:
        old_name: Current tag name
        new_name: New tag name
        db_path: Path to the SQLite database

    Returns:
        Number of slide associations that were affected
    """
    conn = get_db(db_path)
    old_name = old_name.strip().lower()
    new_name = new_name.strip().lower()

    # Count slide associations before rename
    tag_row = conn.execute("SELECT id FROM tags WHERE name = ?", (old_name,)).fetchone()
    assoc_count = 0
    if tag_row:
        assoc_count = conn.execute(
            "SELECT COUNT(*) as cnt FROM slide_tags WHERE tag_id = ?", (tag_row["id"],)
        ).fetchone()["cnt"]
        conn.execute("UPDATE tags SET name = ? WHERE id = ?", (new_name, tag_row["id"]))

    # Also rename in taxonomy if present
    conn.execute("UPDATE taxonomy SET name = ? WHERE name = ?", (new_name, old_name))

    conn.commit()
    conn.close()
    return assoc_count


# ---------------------------------------------------------------------------
# Taxonomy management
# ---------------------------------------------------------------------------


def list_taxonomy(db_path: str = "slides.db") -> List[Dict]:
    """List all taxonomy tags grouped by category.

    Args:
        db_path: Path to the SQLite database

    Returns:
        List of dicts with 'name' and 'category' keys
    """
    conn = get_db(db_path)
    rows = conn.execute(
        "SELECT name, category FROM taxonomy ORDER BY category, name"
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


def add_taxonomy_tags(
    tags: List[Dict],
    db_path: str = "slides.db",
) -> Tuple[int, int]:
    """Add tags to the taxonomy (upsert: insert new, update category on existing).

    Args:
        tags: List of dicts with 'name' and optional 'category' keys
        db_path: Path to the SQLite database

    Returns:
        Tuple of (new_count, updated_count)
    """
    conn = get_db(db_path)
    new_count = 0
    updated_count = 0
    for tag in tags:
        name = tag["name"].strip().lower()
        category = tag.get("category", "").strip()
        if not name:
            continue
        existing = conn.execute("SELECT id FROM taxonomy WHERE name = ?", (name,)).fetchone()
        if existing:
            conn.execute("UPDATE taxonomy SET category = ? WHERE id = ?", (category, existing["id"]))
            updated_count += 1
        else:
            conn.execute(
                "INSERT INTO taxonomy (name, category) VALUES (?, ?)",
                (name, category),
            )
            new_count += 1
    conn.commit()
    conn.close()
    return new_count, updated_count


def remove_taxonomy_tag(
    tag_name: str,
    db_path: str = "slides.db",
) -> bool:
    """Remove a tag from the taxonomy.

    Does NOT remove the tag from slides -- taxonomy removal only means
    'stop suggesting this tag'.

    Args:
        tag_name: Tag name to remove
        db_path: Path to the SQLite database

    Returns:
        True if removed, False if not found
    """
    conn = get_db(db_path)
    tag_name = tag_name.strip().lower()
    cur = conn.execute("DELETE FROM taxonomy WHERE name = ?", (tag_name,))
    conn.commit()
    removed = cur.rowcount > 0
    conn.close()
    return removed


def import_taxonomy_csv(
    csv_path: str,
    db_path: str = "slides.db",
) -> Tuple[int, int]:
    """Import taxonomy tags from a CSV file.

    Expects columns 'name' and optionally 'category'.
    Uses upsert semantics: new tags are inserted, existing tags get their
    category updated.

    Args:
        csv_path: Path to the CSV file
        db_path: Path to the SQLite database

    Returns:
        Tuple of (new_count, updated_count)
    """
    tags = []
    with open(csv_path, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        fieldnames = reader.fieldnames or []
        name_field = "name" if "name" in fieldnames else (fieldnames[0] if fieldnames else "")
        cat_field = "category" if "category" in fieldnames else None
        for row in reader:
            name = row.get(name_field, "").strip()
            if not name:
                continue
            category = row.get(cat_field, "").strip() if cat_field else ""
            tags.append({"name": name, "category": category})
    return add_taxonomy_tags(tags, db_path)


def export_taxonomy_csv(
    csv_path: str,
    db_path: str = "slides.db",
) -> int:
    """Export taxonomy tags to a CSV file.

    Args:
        csv_path: Output CSV file path
        db_path: Path to the SQLite database

    Returns:
        Number of tags exported
    """
    tags = list_taxonomy(db_path)
    with open(csv_path, "w", encoding="utf-8", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=["name", "category"])
        writer.writeheader()
        for tag in tags:
            writer.writerow(tag)
    return len(tags)


def get_taxonomy_names(db_path: str = "slides.db") -> List[str]:
    """Get flat list of taxonomy tag names (for LLM constrained tagging).

    Args:
        db_path: Path to the SQLite database

    Returns:
        List of tag name strings
    """
    conn = get_db(db_path)
    rows = conn.execute("SELECT name FROM taxonomy ORDER BY name").fetchall()
    conn.close()
    return [r["name"] for r in rows]


# ---------------------------------------------------------------------------
# Section management
# ---------------------------------------------------------------------------


def set_slide_section(
    slide_id: int,
    section_name: str,
    deck_id: int,
    db_path: str = "slides.db"
) -> None:
    """Assign a slide to a section.

    Args:
        slide_id: The slide ID
        section_name: Section name to assign
        deck_id: The deck ID (sections are deck-scoped)
        db_path: Path to the SQLite database
    """
    conn = get_db(db_path)
    name = section_name.strip().lower()
    if not name:
        conn.close()
        return

    # Get or create section
    conn.execute(
        "INSERT OR IGNORE INTO sections (deck_id, name, position) "
        "VALUES (?, ?, (SELECT COALESCE(MAX(position), 0) + 1 FROM sections WHERE deck_id = ?))",
        (deck_id, name, deck_id)
    )
    section_row = conn.execute(
        "SELECT id FROM sections WHERE deck_id = ? AND name = ?",
        (deck_id, name)
    ).fetchone()

    if section_row:
        # Upsert slide association
        conn.execute(
            "INSERT OR REPLACE INTO slide_sections (slide_id, section_id) VALUES (?, ?)",
            (slide_id, section_row["id"])
        )

    conn.commit()
    conn.close()


def get_slide_section(slide_id: int, db_path: str = "slides.db") -> Optional[str]:
    """Get the section name for a slide.

    Args:
        slide_id: The slide ID
        db_path: Path to the SQLite database

    Returns:
        Section name or None if no section assigned
    """
    conn = get_db(db_path)
    row = conn.execute(
        """SELECT sec.name FROM sections sec
           JOIN slide_sections ss ON sec.id = ss.section_id
           WHERE ss.slide_id = ?""",
        (slide_id,)
    ).fetchone()
    conn.close()
    return row["name"] if row else None


def get_deck_sections(deck_id: int, db_path: str = "slides.db") -> List[Dict]:
    """Get all sections for a deck with slide counts.

    Args:
        deck_id: The deck ID
        db_path: Path to the SQLite database

    Returns:
        List of dicts with 'id', 'name', 'position', 'slide_count' keys
    """
    conn = get_db(db_path)
    rows = conn.execute(
        """SELECT sec.id, sec.name, sec.position,
                  COUNT(ss.slide_id) as slide_count
           FROM sections sec
           LEFT JOIN slide_sections ss ON sec.id = ss.section_id
           WHERE sec.deck_id = ?
           GROUP BY sec.id
           ORDER BY sec.position""",
        (deck_id,)
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


def remove_slide_section(slide_id: int, db_path: str = "slides.db") -> bool:
    """Remove a slide's section assignment.

    Args:
        slide_id: The slide ID
        db_path: Path to the SQLite database

    Returns:
        True if a section was removed, False otherwise
    """
    conn = get_db(db_path)
    cur = conn.execute(
        "DELETE FROM slide_sections WHERE slide_id = ?",
        (slide_id,)
    )
    conn.commit()
    removed = cur.rowcount > 0
    conn.close()
    return removed


def rename_section(
    deck_id: int,
    old_name: str,
    new_name: str,
    db_path: str = "slides.db"
) -> int:
    """Rename a section within a deck.

    Args:
        deck_id: The deck ID
        old_name: Current section name
        new_name: New section name
        db_path: Path to the SQLite database

    Returns:
        Number of affected slide associations
    """
    conn = get_db(db_path)
    old_name = old_name.strip().lower()
    new_name = new_name.strip().lower()

    # Get section and count associations
    section_row = conn.execute(
        "SELECT id FROM sections WHERE deck_id = ? AND name = ?",
        (deck_id, old_name)
    ).fetchone()

    assoc_count = 0
    if section_row:
        assoc_count = conn.execute(
            "SELECT COUNT(*) as cnt FROM slide_sections WHERE section_id = ?",
            (section_row["id"],)
        ).fetchone()["cnt"]

        conn.execute(
            "UPDATE sections SET name = ? WHERE id = ?",
            (new_name, section_row["id"])
        )

    conn.commit()
    conn.close()
    return assoc_count


_EDITABLE_FIELDS = {"notes"}


def record_edit(
    slide_id: int,
    field: str,
    new_value: str,
    *,
    source: str = "web",
    db_path: str = "slides.db",
) -> bool:
    """Update a slide field and record the change in edit_history.

    Returns True if the field was changed, False if value was identical.
    Skips history write and field update when old == new.
    """
    if field not in _EDITABLE_FIELDS:
        raise ValueError(f"Field {field!r} is not editable via record_edit()")
    conn = get_db(db_path)
    try:
        row = conn.execute(
            f"SELECT {field} FROM slides WHERE id = ?", (slide_id,)
        ).fetchone()
        if row is None:
            raise ValueError(f"Slide {slide_id} not found")

        old_value = row[field] or ""
        if old_value == new_value:
            return False

        conn.execute(
            "INSERT INTO edit_history (slide_id, field, old_value, new_value, source) "
            "VALUES (?, ?, ?, ?, ?)",
            (slide_id, field, old_value, new_value, source),
        )
        conn.execute(
            f"UPDATE slides SET {field} = ?, updated_at = datetime('now') WHERE id = ?",
            (new_value, slide_id),
        )
        conn.commit()
        return True
    finally:
        conn.close()

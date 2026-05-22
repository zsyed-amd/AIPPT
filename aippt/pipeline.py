"""Presentation creation pipeline — shared by CLI and web UI."""

import os
from dataclasses import dataclass, field
from typing import Callable, Optional
import logging

logger = logging.getLogger(__name__)


@dataclass
class PipelineConfig:
    """Configuration for the presentation creation pipeline."""

    # Required
    outline_text: str
    template_path: str
    output_path: str

    # Enhancement
    enhance: bool = False
    model: Optional[str] = None
    audience: Optional[str] = None
    show_plan: bool = False
    no_plan: bool = False

    # LLM connection
    gateway_config: Optional[str] = None
    api_key: Optional[str] = None
    api_base: Optional[str] = None

    # Image generation
    image_gen: str = "none"
    mcp_config: str = "mcp_servers.json"
    mcp_server: str = "txt2img"
    classification: str = "internal"

    # File context
    outline_path: Optional[str] = None

    # Source / origin tracking
    source_engine: Optional[str] = "python-pptx"
    source_theme: Optional[str] = None
    source_kind: Optional[str] = "outline"

    # Corporate template merge
    corp_template: Optional[str] = None

    # Callbacks
    progress_callback: Optional[Callable] = field(default=None, repr=False)


@dataclass
class PipelineResult:
    """Result from a completed pipeline run."""

    output_path: str
    slide_count: int
    title: str
    engine: Optional[str] = None
    theme: Optional[str] = None
    source_kind: Optional[str] = None


def run_pipeline(config: PipelineConfig) -> PipelineResult:
    """Create a PPTX presentation from a PipelineConfig.

    This is the canonical orchestrator extracted from cli.create_deck().
    All parameters are sourced from `config`; caller constructs the
    PipelineConfig dataclass rather than passing individual keyword args.

    Args:
        config: Fully-populated PipelineConfig.

    Returns:
        PipelineResult with output_path, slide_count, and title.

    Raises:
        FileNotFoundError: If config.template_path does not exist.
        RuntimeError: If generation fails fatally (e.g. final save error).
    """
    from pptx import Presentation

    from aippt.parser import parse_outline
    from aippt.llm import LLMClient, load_gateway_config
    from aippt.enhancer import enhance_with_llm
    from aippt.images import setup_image_directory
    from aippt.layouts import remove_all_slides
    from aippt.builder import BuildContext, build_slide

    def _notify(step, detail=""):
        logger.info(detail or step)
        if config.progress_callback:
            config.progress_callback(step, detail)

    # Validate template
    if not os.path.exists(config.template_path):
        raise FileNotFoundError(f"Template file not found: {config.template_path}")

    # Extract frontmatter metadata (audience, goal, tone)
    from aippt.parser import parse_frontmatter
    outline_text = config.outline_text
    frontmatter, outline_text = parse_frontmatter(outline_text)

    # Resolve audience: config arg > frontmatter > default
    audience = config.audience
    if audience is not None:
        audience_source = "cli"
    else:
        fm_audience = frontmatter.get('audience', '').lower()
        valid_audiences = {'engineers', 'executives', 'product', 'mixed'}
        if fm_audience in valid_audiences:
            audience = fm_audience
            audience_source = "frontmatter"
        else:
            audience = "mixed"
            audience_source = "default"
    logger.info(f"Target audience: {audience} (source: {audience_source})")

    # Parse the outline text
    _notify("parse", f"Parsing outline ({len(outline_text)} chars)")
    parsed = parse_outline(outline_text)
    slides = parsed['slides']
    sections = parsed['sections']
    total_slides = len(slides)
    logger.info(f"Loaded outline with {total_slides} slides")

    # Resolve IMAGE: directive paths relative to the outline file
    has_images = any('image' in s for s in slides)
    if has_images and not config.outline_path:
        logger.warning(
            "IMAGE: directives found but no outline file path provided; "
            "image resolution skipped"
        )
        for slide in slides:
            slide.pop('image', None)
    elif config.outline_path:
        from aippt.parser import resolve_image_path
        outline_dir = os.path.dirname(os.path.abspath(config.outline_path))
        for slide in slides:
            if 'image' in slide:
                resolved = resolve_image_path(slide['image'], outline_dir)
                if resolved:
                    slide['image'] = resolved
                else:
                    del slide['image']

    if sections:
        logger.info(f"Found {len(sections)} sections in outline")

    _notify("parse", f"Parsed outline with {total_slides} slides")

    # Load template
    prs = Presentation(config.template_path)
    slide_count_before = len(prs.slides)
    remove_all_slides(prs)
    if slide_count_before > 0:
        logger.info(f"Removed {slide_count_before} template placeholder slide(s)")
    logger.info(f"Loaded template: {config.template_path}")

    # Create image directory if using image generation
    image_dir = None
    if config.image_gen != 'none':
        image_dir = setup_image_directory(config.output_path)
        logger.info(f"Created image directory: {image_dir}")

    # Create MCP manager for image generation
    mcp_manager = None
    if config.image_gen == 'mcp':
        from aippt.mcp import MCPManager
        try:
            mcp_manager = MCPManager(config.mcp_config)
            if config.mcp_server not in mcp_manager.servers:
                logger.warning(
                    f"MCP server '{config.mcp_server}' not found in config; "
                    "image generation disabled"
                )
                mcp_manager = None
            else:
                logger.info(f"MCP image generation enabled via server: {config.mcp_server}")
        except Exception as e:
            logger.warning(f"Failed to initialize MCP manager: {e}; image generation disabled")
            mcp_manager = None

    # Setup LLM client if needed
    client = None
    resolved_model = config.model
    if config.enhance or config.image_gen != 'none':
        from aippt.config import get_model_default, ConfigError
        try:
            resolved_model = config.model or get_model_default("enhance")
        except ConfigError as exc:
            raise RuntimeError(str(exc)) from exc

        gateway = None
        if config.gateway_config and os.path.exists(config.gateway_config):
            gateway = load_gateway_config(config.gateway_config)
            if gateway:
                logger.info(f"Using gateway config: {config.gateway_config}")

        try:
            client = LLMClient(
                model=resolved_model,
                api_key=config.api_key,
                api_base=config.api_base,
                gateway=gateway,
            )
        except (ConfigError, ValueError) as exc:
            raise RuntimeError(str(exc)) from exc
        logger.info(f"Using model: {resolved_model} via {client.model_config.provider} API")

    # Deck-level narrative planning (before per-slide enhancement)
    if config.show_plan and config.no_plan:
        logger.warning("--show-plan ignored because --no-plan was specified")
    deck_plan = None
    if config.enhance and not config.no_plan:
        from aippt.enhancer import plan_deck
        _notify("plan", "Planning deck narrative structure...")
        deck_plan = plan_deck(slides, client, audience=audience, image_gen=config.image_gen)
        if deck_plan['slides']:
            logger.info(
                f"Deck plan: {deck_plan['narrative_arc']} arc, "
                f"{len(deck_plan['slides'])} slides planned"
            )
        else:
            logger.warning("Deck planning returned empty plan; enhancing without deck context")
        if config.show_plan and deck_plan['slides']:
            print("\n=== Deck Narrative Plan ===")
            print(f"Narrative arc: {deck_plan['narrative_arc']}")
            print(f"Assessment: {deck_plan['arc_assessment']}")
            print()
            for entry in deck_plan['slides']:
                print(
                    f"  Slide {entry['index'] + 1}: [{entry['role']}] "
                    f"{entry['title']} -> {entry['suggested_layout']}"
                )
                if entry.get('context_hint'):
                    print(f"    Context: {entry['context_hint']}")
                if entry.get('transition_to_next'):
                    print(f"    Transition: {entry['transition_to_next']}")
            print("===========================\n")

    # Enhance slides with LLM if requested
    if config.enhance:
        for i, slide in enumerate(slides, 1):
            _notify("enhance", f"Enhancing slide {i}/{len(slides)}: {slide['title']}")

            # Skip enhancement for functional slides (title slides, section
            # dividers) — those with sparse content and no explicit LAYOUT
            # directive.  Enhancing these risks hallucinated content.
            has_layout_directive = 'layout' in slide
            bullet_count = sum(1 for line in slide.get('content', []) if line.strip())
            if bullet_count <= 2 and not has_layout_directive:
                logger.info(
                    f"Skipping enhancement for functional slide: "
                    f"'{slide['title']}' ({bullet_count} bullet(s))"
                )
                slide['original_content'] = list(slide['content'])
                # Still attach deck plan context for notes if available
                if deck_plan and deck_plan.get('slides'):
                    plan_entries = deck_plan['slides']
                    if i - 1 < len(plan_entries):
                        entry = plan_entries[i - 1]
                        slide['_deck_context'] = {
                            'role': entry.get('role', ''),
                            'suggested_layout': entry.get('suggested_layout', ''),
                            'transition_to_next': entry.get('transition_to_next', ''),
                            'context_hint': entry.get('context_hint', ''),
                        }
                    if deck_plan:
                        slide['_narrative_arc'] = deck_plan.get('narrative_arc', '')
                continue

            try:
                slide['original_content'] = list(slide['content'])
                # Look up this slide's deck context from the plan
                slide_deck_context = None
                if deck_plan and deck_plan.get('slides'):
                    plan_entries = deck_plan['slides']
                    if i - 1 < len(plan_entries):
                        entry = plan_entries[i - 1]
                        slide_deck_context = {
                            'role': entry.get('role', ''),
                            'suggested_layout': entry.get('suggested_layout', ''),
                            'transition_to_next': entry.get('transition_to_next', ''),
                            'context_hint': entry.get('context_hint', ''),
                        }
                enhanced_content = enhance_with_llm(
                    slide, client, image_gen=config.image_gen,
                    has_image='image' in slide,
                    audience=audience,
                    deck_context=slide_deck_context,
                )
                slide['content'] = enhanced_content.split('\n')
                if slide_deck_context:
                    slide['_deck_context'] = slide_deck_context
                if deck_plan:
                    slide['_narrative_arc'] = deck_plan.get('narrative_arc', '')
            except Exception as e:
                logger.error(f"Error enhancing slide {i}: {str(e)}")
                logger.info("Continuing with original content for this slide")

            # Save after each enhancement
            try:
                prs.save(config.output_path)
                logger.info(f"Progress saved after enhancing slide {i}")
            except Exception as e:
                logger.error(f"Error saving progress after slide {i}: {str(e)}")
        _notify("enhance", f"All {len(slides)} slides enhanced")

    # Build context shared across all slides
    build_ctx = BuildContext(
        client=client,
        image_gen=config.image_gen,
        image_dir=image_dir,
        model=resolved_model if config.enhance else None,
        mcp_manager=mcp_manager,
        mcp_server=config.mcp_server,
        classification=config.classification,
        audience=audience,
        audience_source=audience_source,
    )

    # Create slides
    layout_counts = {}
    for i, slide_item in enumerate(slides, 1):
        _notify("build", f"Creating slide {i}/{len(slides)}: {slide_item['title']}")
        try:
            slide_item['slide_num'] = i
            layout_type = build_slide(prs, slide_item, build_ctx)
            if layout_type:
                layout_counts[layout_type] = layout_counts.get(layout_type, 0) + 1

            # Save after each slide creation
            try:
                prs.save(config.output_path)
                logger.info(f"Progress saved after creating slide {i}")
            except Exception as e:
                logger.error(f"Error saving progress after slide {i}: {str(e)}")

        except Exception as e:
            logger.error(f"Error creating slide {i}: {str(e)}")
            try:
                prs.save(config.output_path)
                logger.info("Progress saved despite error")
            except Exception as save_error:
                logger.error(f"Error saving progress: {str(save_error)}")
            continue

    _notify("build", f"Built {len(prs.slides)} slides")

    # Log layout distribution summary
    if layout_counts:
        summary = ", ".join(f"{count} {ltype}" for ltype, count in sorted(layout_counts.items()))
        logger.info(f"Layout mix: {summary}")

    # Add image directory info to presentation notes
    if image_dir and len(prs.slides) > 0:
        try:
            notes_slide = prs.slides[0].notes_slide
            current_notes = notes_slide.notes_text_frame.text
            notes_slide.notes_text_frame.text = f"Image Directory: {image_dir}\n\n{current_notes}"
        except Exception as e:
            logger.error(f"Error adding image directory to notes: {str(e)}")

    # Apply sections from outline structure
    if sections:
        try:
            from aippt.sections import write_sections, Section
            sections_to_write = []
            for section_data in sections:
                slide_ids = [prs.slides[i].slide_id for i in section_data["slide_indices"]]
                sections_to_write.append(Section(name=section_data["name"], slide_ids=slide_ids))
            write_sections(prs, sections_to_write)
            logger.info(f"Applied {len(sections_to_write)} sections to presentation")
        except Exception as e:
            logger.error(f"Error applying sections: {str(e)}")

    # Final save
    try:
        prs.save(config.output_path)
        logger.info(f"PowerPoint presentation completed: {config.output_path}")
    except Exception as e:
        raise RuntimeError(f"Error saving final presentation: {str(e)}") from e

    # Corporate template merge post-processing
    if config.corp_template:
        from aippt.template_merge import merge_with_template
        pre_merge = config.output_path + ".pre-merge.pptx"
        os.rename(config.output_path, pre_merge)
        try:
            merge_result = merge_with_template(pre_merge, config.corp_template, config.output_path)
            logger.info(
                f"Corporate template merge: {merge_result['slide_count']} slides, "
                f"{len(set(a['target_layout'] for a in merge_result['layout_assignments']))} layout types used"
            )
        finally:
            if os.path.exists(pre_merge):
                os.unlink(pre_merge)

    # Embed [AIPPT-META] lineage on slide 1 of the generated PPTX
    if config.source_kind and len(prs.slides) > 0:
        from datetime import datetime, timezone as _tz
        from aippt.metadata import write_deck_lineage
        engine_label = config.source_engine or "python-pptx"
        source_str = f"{config.source_kind} -> {engine_label}"
        try:
            write_deck_lineage(
                pptx_path=config.output_path,
                source=source_str,
                engine=engine_label,
                theme=config.source_theme,
                generated_at=datetime.now(_tz.utc).isoformat(),
            )
        except Exception as _exc:
            logger.warning("write_deck_lineage failed (non-fatal): %s", _exc)

    title = slides[0]['title'] if slides else ""
    return PipelineResult(
        output_path=config.output_path,
        slide_count=len(prs.slides),
        title=title,
        engine=config.source_engine,
        theme=config.source_theme,
        source_kind=config.source_kind,
    )

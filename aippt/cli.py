"""CLI entry point with subcommands."""
import argparse
import logging
import os
import subprocess
import sys

from aippt import graph, render
from aippt.config import load_sharepoint_config

logger = logging.getLogger(__name__)


def cmd_create(args):
    """Create a presentation from a markdown outline."""
    from aippt.layouts import inspect_template
    from aippt.pipeline import PipelineConfig, run_pipeline

    # Validate input files
    if not os.path.exists(args.outline):
        logger.error(f"Outline file not found: {args.outline}")
        return 1

    if not os.path.exists(args.template):
        logger.error(f"Template file not found: {args.template}")
        return 1

    # Analyze template if requested
    if args.analyze_template:
        try:
            template_info = inspect_template(args.template)
            logger.info("Template analysis:")
            logger.info(f"Available layouts: {[l['name'] for l in template_info['layouts']]}")
            logger.info(f"Slide size: {template_info['slide_size']}")
        except Exception as e:
            logger.error(f"Error analyzing template: {str(e)}")

    # Read outline
    with open(args.outline, 'r', encoding='utf-8') as file:
        outline_text = file.read()

    # Apply --test slicing at the CLI layer by rebuilding a trimmed outline
    if args.test:
        from aippt.parser import parse_outline
        parsed = parse_outline(outline_text)
        all_slides = parsed['slides']
        total_slides = len(all_slides)
        test_slides = min(args.test, total_slides)
        logger.info(f"Test mode: Processing first {test_slides} of {total_slides} slides")

        sliced = all_slides[:test_slides]
        lines = []
        for slide in sliced:
            lines.append(f"## {slide['title']}")
            if 'layout' in slide:
                lines.append(f"LAYOUT: {slide['layout']}")
            if 'image' in slide:
                lines.append(f"IMAGE: {slide['image']}")
            for line in slide['content']:
                lines.append(line)
            lines.append("")
        outline_text = "\n".join(lines)

    try:
        config = PipelineConfig(
            outline_text=outline_text,
            template_path=args.template,
            output_path=args.output,
            enhance=args.enhance,
            model=args.model,
            gateway_config=args.gateway_config,
            api_key=args.api_key,
            api_base=args.api_base,
            image_gen=args.image_gen,
            outline_path=args.outline,
            mcp_config=args.mcp_config,
            classification=args.classification,
            mcp_server=args.mcp_server,
            audience=getattr(args, 'audience', None),
            show_plan=getattr(args, 'show_plan', False),
            no_plan=getattr(args, 'no_plan', False),
            corp_template=getattr(args, 'corp_template', None),
        )
        run_pipeline(config)
    except FileNotFoundError as e:
        logger.error(str(e))
        return 1
    except RuntimeError as e:
        logger.error(str(e))
        return 1

    logger.info(f"PowerPoint presentation completed: {args.output}")
    return 0


def cmd_reverse(args):
    """Convert a PowerPoint back to markdown outline."""
    from aippt.reverse import convert_pptx_to_outline

    if not os.path.exists(args.input):
        logger.error(f"File not found: {args.input}")
        return 1

    output = args.output or os.path.splitext(args.input)[0] + '.md'
    strip_notes = getattr(args, 'strip_notes', False)
    include_notes = not getattr(args, 'no_notes', False) and not strip_notes
    enhance = getattr(args, 'enhance', False)

    # Enhanced mode: set up LLM client
    llm_client = None
    if enhance:
        from aippt.llm import LLMClient, load_gateway_config
        from aippt.config import get_model_default, ConfigError

        try:
            model = getattr(args, 'model', None) or get_model_default("reverse")
        except (ConfigError, KeyError, ValueError) as exc:
            logger.error(str(exc))
            return 1

        gateway = None
        gateway_config_path = getattr(args, 'gateway_config', None)
        if gateway_config_path and os.path.exists(gateway_config_path):
            gateway = load_gateway_config(gateway_config_path)

        try:
            llm_client = LLMClient(model=model, gateway=gateway)
        except (ConfigError, ValueError) as exc:
            logger.error(str(exc))
            return 1

        print(f"Enhanced reverse using model: {model}")

    images_dir = getattr(args, 'images_dir', None)

    success = convert_pptx_to_outline(
        args.input, output, include_notes,
        enhance=enhance, llm_client=llm_client, images_dir=images_dir,
    )

    if success:
        print(f"Converted to: {output}")
    return 0 if success else 1


def cmd_catalog(args):
    """Catalog a deck into the slide database."""
    from aippt.catalog import catalog_deck
    from aippt.config import load_dirs_config

    if not os.path.exists(args.deck):
        logger.error(f"File not found: {args.deck}")
        return 1

    dirs = load_dirs_config()
    images_dir = getattr(args, 'images_dir', None)
    db_path = args.db if args.db != "slides.db" else dirs["directories"]["db"]
    deck_id = catalog_deck(args.deck, db_path=db_path, images_dir=images_dir, base_dir=dirs["base_dir"])
    print(f"Cataloged deck (id={deck_id}): {args.deck}")
    return 0


def cmd_search(args):
    """Search cataloged slides by tags, title, or section."""
    from aippt.catalog import search_slides
    from aippt.remix import generate_manifest

    tags = [t.strip() for t in args.tags.split(",")] if args.tags else None
    results = search_slides(
        db_path=args.db,
        tags=tags,
        title_contains=args.title_contains,
        section=args.section
    )

    if not results:
        print("No matching slides found.")
        return 0

    print(f"Found {len(results)} matching slide(s):\n")
    for r in results:
        print(f"  [{r['deck_name']}] Slide {r['position']}: {r['title']}")

    if args.export_manifest:
        manifest_yaml = generate_manifest(results, db_path=args.db)
        with open(args.export_manifest, "w", encoding="utf-8") as f:
            f.write(manifest_yaml)
        print(f"\nManifest written to {args.export_manifest}")

    return 0


def cmd_remix(args):
    """Assemble a new deck from a manifest file."""
    from aippt.remix import assemble_deck

    if not os.path.exists(args.manifest):
        logger.error(f"Manifest file not found: {args.manifest}")
        return 1

    count = assemble_deck(args.manifest, args.output, db_path=args.db)
    print(f"Remixed deck saved to {args.output} ({count} slides)")
    return 0


def cmd_analyze(args):
    """Run multimodal analysis on slide images."""
    from pptx import Presentation

    from aippt.llm import LLMClient, load_gateway_config
    from aippt.analyze import analyze_slide, parse_tags_response, load_taxonomy
    from aippt.catalog import get_db, add_tags, catalog_deck

    if not os.path.exists(args.deck):
        logger.error(f"File not found: {args.deck}")
        return 1

    # Resolve model: CLI flag > models.yaml per-mode default
    from aippt.config import get_model_default, ConfigError
    # Map analyze modes to config operations
    mode_to_operation = {"feedback": "feedback", "notes": "notes", "tags": "tags", "improvements": "feedback"}
    operation = mode_to_operation.get(args.mode, "feedback")
    try:
        model = args.model or get_model_default(operation)
    except ConfigError as exc:
        logger.error(str(exc))
        return 1

    # Setup LLM client
    gateway = None
    if args.gateway_config and os.path.exists(args.gateway_config):
        gateway = load_gateway_config(args.gateway_config)

    try:
        client = LLMClient(
            model=model,
            api_key=getattr(args, 'api_key', None),
            gateway=gateway,
        )
    except (ConfigError, ValueError) as exc:
        logger.error(str(exc))
        return 1

    # Determine images directory
    images_dir = getattr(args, 'images_dir', None)
    if not images_dir:
        deck_name = os.path.splitext(os.path.basename(args.deck))[0]
        images_dir = os.path.join("images", deck_name)

    if not os.path.isdir(images_dir):
        logger.warning(
            f"Images directory not found: {images_dir} — will use text-only analysis for all slides."
        )
        images_dir = None

    # Ensure deck is cataloged
    deck_id = catalog_deck(args.deck, db_path=args.db, images_dir=images_dir)

    # Load taxonomy: CLI CSV flag > DB taxonomy table > freeform
    taxonomy = None
    if getattr(args, 'taxonomy', None):
        taxonomy = load_taxonomy(args.taxonomy)
        logger.info(f"Loaded {len(taxonomy)} tags from taxonomy CSV")
    elif args.mode == "tags":
        from aippt.catalog import get_taxonomy_names
        db_taxonomy = get_taxonomy_names(args.db)
        if db_taxonomy:
            taxonomy = db_taxonomy
            logger.info(f"Using {len(taxonomy)} tags from database taxonomy")

    # Get slides from catalog then close the connection immediately so that
    # add_tags (which opens its own connection) does not contend with an
    # open handle on the same database file.
    conn = get_db(args.db)
    slides = conn.execute(
        "SELECT id, position, title, image_path FROM slides WHERE deck_id = ? ORDER BY position",
        (deck_id,),
    ).fetchall()
    conn.close()

    prs = Presentation(args.deck)

    for slide_row in slides:
        image_path = slide_row["image_path"]
        if not image_path and images_dir:
            # Try to find image by position
            for ext in (".png", ".PNG", ".jpg", ".jpeg"):
                candidate = os.path.join(images_dir, f"Slide{slide_row['position']}{ext}")
                if os.path.exists(candidate):
                    image_path = candidate
                    break

        # Extract text content from PPTX shapes for text-only fallback
        pptx_slide = prs.slides[slide_row["position"] - 1]
        content_text = _extract_slide_text(pptx_slide, slide_row["title"])

        if image_path and os.path.exists(image_path):
            logger.info(f"Analyzing slide {slide_row['position']} (image): {slide_row['title']}")
        else:
            logger.info(
                f"Analyzing slide {slide_row['position']} (text-only): {slide_row['title']}"
            )
            image_path = None  # normalise so analyze_slide receives None

        try:
            result = analyze_slide(
                client=client,
                image_path=image_path,
                mode=args.mode,
                title=slide_row["title"],
                taxonomy=taxonomy,
                content_text=content_text,
            )

            if args.mode == "feedback":
                print(f"\n--- Slide {slide_row['position']}: {slide_row['title']} ---")
                print(result)

            elif args.mode == "notes":
                # Write notes back into PPTX (pptx_slide already fetched above)
                notes_slide = pptx_slide.notes_slide
                notes_slide.notes_text_frame.text = result
                # Also update notes in the database
                from aippt.catalog import record_edit
                record_edit(
                    slide_row["id"], "notes", result,
                    source="ai-notes", db_path=args.db,
                )
                print(f"Slide {slide_row['position']}: notes generated")

            elif args.mode == "tags":
                tags = parse_tags_response(result)
                source = "taxonomy" if taxonomy else "ai"
                add_tags(slide_row["id"], tags, source=source, db_path=args.db)
                print(f"Slide {slide_row['position']}: tagged with {tags}")

            elif args.mode == "improvements":
                print(f"\n--- Slide {slide_row['position']}: {slide_row['title']} ---\n")
                print(result)

        except Exception as e:
            logger.error(f"Error analyzing slide {slide_row['position']}: {e}")
            continue

    if args.mode == "notes":
        prs.save(args.deck)
        logger.info(f"Notes saved to {args.deck}")

    return 0


def cmd_improve(args):
    """Improve an existing presentation using LLM analysis and rewrite."""
    from aippt.improve import improve_deck
    from aippt.llm import LLMClient, load_gateway_config
    from aippt.config import get_model_default, ConfigError

    if not os.path.exists(args.deck):
        logger.error(f"File not found: {args.deck}")
        return 1

    # Resolve model — try 'improve' default, fall back to 'enhance' for older configs
    try:
        model = args.model or get_model_default("improve")
    except (ConfigError, KeyError, ValueError):
        try:
            model = args.model or get_model_default("enhance")
        except ConfigError as exc:
            logger.error(str(exc))
            return 1

    # Setup LLM client
    gateway = None
    if args.gateway_config and os.path.exists(args.gateway_config):
        gateway = load_gateway_config(args.gateway_config)
        if gateway:
            logger.info(f"Using gateway config: {args.gateway_config}")

    try:
        client = LLMClient(model=model, api_key=args.api_key, gateway=gateway)
    except (ConfigError, ValueError) as exc:
        logger.error(str(exc))
        return 1

    logger.info(f"Using model: {model} via {client.model_config.provider} API")

    # Auto-detect images directory
    images_dir = getattr(args, 'images_dir', None)
    if not images_dir:
        deck_name = os.path.splitext(os.path.basename(args.deck))[0]
        candidate = os.path.join('images', deck_name)
        if os.path.isdir(candidate):
            images_dir = candidate
            logger.info(f"Auto-detected images directory: {images_dir}")

    # Parse slides filter
    slides_filter = None
    if args.slides:
        slides_filter = [int(s.strip()) for s in args.slides.split(',')]

    audience = getattr(args, 'audience', None) or 'mixed'

    results = improve_deck(
        pptx_path=args.deck,
        output_path=args.output,
        images_dir=images_dir,
        slides_filter=slides_filter,
        passes=args.passes,
        dry_run=args.dry_run,
        client=client,
        focus=args.focus,
        audience=audience,
        keep_titles=getattr(args, 'keep_titles', False),
        max_retries=getattr(args, 'max_retries', 2),
        no_validate=getattr(args, 'no_validate', False),
    )

    # Print summary
    from collections import Counter
    counts = Counter(r.get('status', 'unknown') for r in results)

    if args.dry_run:
        print(f"\n[DRY RUN] Would improve {counts.get('dry_run', 0)} slide(s), "
              f"skipped {counts.get('no_content', 0)} (no content)")
        for r in results:
            if r.get('improved'):
                print(f"\n--- Slide {r.get('slide_num', '?')}: {r['title']} ---")
                if r.get('title_rewritten') and r.get('new_title'):
                    print(f"Title: '{r['original_title']}' → '{r['new_title']}'")
                print(f"Original:\n{r['original']}\n")
                print(f"Improved:\n{r['improved']}")
    else:
        print(f"\nImproved {counts.get('applied', 0)} slide(s), "
              f"no content: {counts.get('no_content', 0)}, "
              f"no placeholder: {counts.get('no_placeholder', 0)}, "
              f"errors: {counts.get('error', 0)}")

    return 0


def cmd_export(args):
    """Export slide metadata to CSV."""
    from aippt.export import export_csv

    deck_path = getattr(args, 'deck', None)
    export_all = getattr(args, 'all', False)

    if not deck_path and not export_all:
        logger.error("Specify a deck path or use --all")
        return 1

    count = export_csv(
        output_path=args.output,
        db_path=args.db,
        deck_path=deck_path,
        export_all=export_all,
    )
    print(f"Exported {count} slides to {args.output}")
    return 0


def cmd_serve(args):
    """Launch the web UI."""
    import uvicorn
    from aippt.web.app import create_app
    from aippt.config import load_dirs_config, resolve_path

    dirs = load_dirs_config()
    base = dirs["base_dir"]
    d = dirs["directories"]

    gateway_config = getattr(args, 'gateway_config', None)
    uploads_dir = getattr(args, 'uploads_dir', None) or resolve_path(d["uploads"], base)
    db_path = args.db if args.db != "slides.db" else resolve_path(d["db"], base)
    view_only = True if getattr(args, 'view_only', False) else None
    app = create_app(db_path=db_path, gateway_config=gateway_config, uploads_dir=uploads_dir, project_root=base, view_only=view_only)
    mode = "view-only" if app.state.view_only else "full"
    print(f"Starting AIPPT web UI on http://{args.host}:{args.port} ({mode} mode)")
    uvicorn.run(app, host=args.host, port=args.port)
    return 0


def cmd_tags(args):
    """Manage the taxonomy of predefined tags."""
    from aippt.catalog import (
        list_taxonomy,
        add_taxonomy_tags,
        remove_taxonomy_tag,
        import_taxonomy_csv,
        export_taxonomy_csv,
        rename_tag,
    )

    action = getattr(args, "tags_action", None)

    if action == "add":
        category = getattr(args, "category", "") or ""
        new, updated = add_taxonomy_tags(
            [{"name": args.tag, "category": category}], db_path=args.db
        )
        if new:
            cat_msg = f" (category: {category})" if category else ""
            print(f"Added '{args.tag}' to taxonomy{cat_msg}")
        else:
            print(f"Updated '{args.tag}' in taxonomy")
        return 0

    if action == "remove":
        removed = remove_taxonomy_tag(args.tag, db_path=args.db)
        if removed:
            print(f"Removed '{args.tag}' from taxonomy")
        else:
            print(f"Tag '{args.tag}' not found in taxonomy")
        return 0

    if action == "import":
        new, updated = import_taxonomy_csv(args.csv_file, db_path=args.db)
        total = new + updated
        print(f"Imported {total} tags ({new} new, {updated} existing)")
        return 0

    if action == "export":
        count = export_taxonomy_csv(args.csv_file, db_path=args.db)
        print(f"Exported {count} tags to {args.csv_file}")
        return 0

    if action == "rename":
        assoc_count = rename_tag(args.old_name, args.new_name, db_path=args.db)
        print(f"Renamed '{args.old_name}' -> '{args.new_name}' (updated {assoc_count} slide associations)")
        return 0

    # Default: list taxonomy tags grouped by category
    tags = list_taxonomy(db_path=args.db)
    if not tags:
        print("Taxonomy is empty. Use 'tags add' or 'tags import' to populate it.")
        return 0

    # Group by category
    categories = {}
    for t in tags:
        cat = t["category"] or "(uncategorized)"
        categories.setdefault(cat, []).append(t["name"])

    for cat, names in sorted(categories.items()):
        print(f"  Category: {cat}")
        print(f"    {', '.join(names)}")

    return 0


def cmd_decks(args):
    """Manage cataloged decks."""
    import json as json_mod
    from aippt.catalog import (
        list_decks,
        get_deck_slides,
        get_deck_sections,
        resolve_deck,
        delete_deck,
        rename_deck,
        get_deck_tag_count,
        get_deck_top_tags,
        display_name,
    )

    action = getattr(args, "decks_action", None)

    if action == "list":
        decks = list_decks(db_path=args.db)
        if not decks:
            print("No decks cataloged.")
            return 0

        if getattr(args, "json", False):
            for d in decks:
                d["display_name"] = display_name(d["name"])
                d["tag_count"] = get_deck_tag_count(d["id"], args.db)
            print(json_mod.dumps(decks, indent=2))
            return 0

        # Table output
        total_slides = 0
        rows = []
        for d in decks:
            dname = display_name(d["name"])
            tag_count = get_deck_tag_count(d["id"], args.db)
            cataloged = (d.get("cataloged_at") or "")[:10]
            author = d.get("author", "") or "(none)"
            file_path = os.path.basename(d.get("file_path", ""))
            rows.append((d["id"], dname, d["slide_count"], author, cataloged, tag_count, file_path))
            total_slides += d["slide_count"]

        # Column headers and widths
        headers = ("ID", "Name", "Slides", "Author", "Cataloged", "Tags", "File")
        widths = [max(len(str(r[i])) for r in rows + [headers]) for i in range(len(headers))]
        fmt = "  ".join(f"{{:<{w}}}" if i != 2 and i != 5 else f"{{:>{w}}}" for i, w in enumerate(widths))

        print(fmt.format(*headers))
        print(fmt.format(*(chr(0x2500) * w for w in widths)))
        for row in rows:
            print(fmt.format(*[str(v) for v in row]))
        print(f"\n{len(rows)} decks, {total_slides} slides total")
        return 0

    if action == "info":
        result = resolve_deck(args.deck, db_path=args.db)
        if result is None:
            print(f"No deck found matching '{args.deck}'")
            return 1
        if isinstance(result, list):
            print(f"Multiple decks match '{args.deck}':")
            for d in result:
                print(f"  ID {d['id']}: {display_name(d['name'])}")
            print("Use a more specific name or the deck ID.")
            return 1

        deck = result
        slides = get_deck_slides(deck["id"], db_path=args.db)
        sections = get_deck_sections(deck["id"], db_path=args.db)
        tag_count = get_deck_tag_count(deck["id"], args.db)
        top_tags = get_deck_top_tags(deck["id"], args.db)

        if getattr(args, "json", False):
            data = dict(deck)
            data["display_name"] = display_name(deck["name"])
            data["slides"] = [{"position": s["position"], "title": s["title"]} for s in slides]
            data["sections"] = sections
            data["tag_count"] = tag_count
            data["top_tags"] = [{"name": t[0], "count": t[1]} for t in top_tags]
            print(json_mod.dumps(data, indent=2))
            return 0

        dname = display_name(deck["name"])
        print(f"Deck: {dname} (ID: {deck['id']})")
        print(f"File: {deck.get('file_path', '')}")
        if deck.get("author"):
            print(f"Author: {deck['author']}")
        if deck.get("subject"):
            print(f"Subject: {deck['subject']}")
        print(f"Slides: {deck['slide_count']}")
        print(f"Cataloged: {deck.get('cataloged_at', '')}")
        if deck.get("updated_at"):
            print(f"Updated: {deck['updated_at']}")

        # Source tracking metadata
        if deck.get("source_script_path"):
            print(f"\nSource:")
            print(f"  Script: {deck['source_script_path']}")
            if deck.get("source_engine"):
                print(f"  Engine: {deck['source_engine']}")
            if deck.get("source_theme"):
                print(f"  Theme: {deck['source_theme']}")
            if deck.get("outline_path"):
                print(f"  Outline: {deck['outline_path']}")
            if deck.get("source_generated_at"):
                print(f"  Generated: {deck['source_generated_at']}")

        if slides:
            print(f"\nSlides:")
            for s in slides:
                title = s["title"] or "(untitled)"
                print(f"  {s['position']}. {title}")

        if top_tags:
            tag_str = ", ".join(f"{name} ({count})" for name, count in top_tags)
            print(f"\nTags ({tag_count} total):")
            print(f"  {tag_str}")

        return 0

    if action == "rename":
        result = resolve_deck(args.deck, db_path=args.db)
        if result is None:
            print(f"No deck found matching '{args.deck}'")
            return 1
        if isinstance(result, list):
            print(f"Multiple decks match '{args.deck}':")
            for d in result:
                print(f"  ID {d['id']}: {display_name(d['name'])}")
            print("Use a more specific name or the deck ID.")
            return 1

        deck = result
        old_name = rename_deck(deck["id"], args.new_name, db_path=args.db)
        dname = display_name(old_name)
        print(f'Renamed deck {deck["id"]}: "{dname}" -> "{args.new_name}"')
        return 0

    if action == "delete":
        result = resolve_deck(args.deck, db_path=args.db)
        if result is None:
            print(f"No deck found matching '{args.deck}'")
            return 1
        if isinstance(result, list):
            print(f"Multiple decks match '{args.deck}':")
            for d in result:
                print(f"  ID {d['id']}: {display_name(d['name'])}")
            print("Use a more specific name or the deck ID.")
            return 1

        deck = result
        dname = display_name(deck["name"])
        tag_count = get_deck_tag_count(deck["id"], args.db)
        sections = get_deck_sections(deck["id"], db_path=args.db)

        # Derive image directory from slides
        image_dir = None
        if getattr(args, "purge_images", False):
            slides = get_deck_slides(deck["id"], db_path=args.db)
            for s in slides:
                if s.get("image_path"):
                    image_dir = os.path.dirname(s["image_path"])
                    break

        if not getattr(args, "force", False):
            parts = [f"{deck['slide_count']} slides"]
            if tag_count:
                parts.append(f"{tag_count} tags")
            if sections:
                parts.append(f"{len(sections)} sections")
            print(f'Delete deck {deck["id"]} "{dname}"?')
            print(f"  {', '.join(parts)} will be removed.")
            if image_dir:
                print(f"  Image directory will be deleted: {image_dir}")
            try:
                answer = input("  Type 'yes' to confirm: ")
            except (EOFError, KeyboardInterrupt):
                print("\nAborted.")
                return 1
            if answer.strip().lower() != "yes":
                print("Aborted.")
                return 1

        info = delete_deck(deck["id"], db_path=args.db)
        print(f'Deleted deck {deck["id"]} "{dname}" and all associated data.')

        if image_dir and os.path.isdir(image_dir):
            import shutil
            shutil.rmtree(image_dir)
            print(f"Removed image directory: {image_dir}")

        return 0

    if action == "source":
        result = resolve_deck(args.deck, db_path=args.db)
        if result is None:
            print(f"No deck found matching '{args.deck}'")
            return 1
        if isinstance(result, list):
            print(f"Multiple decks match '{args.deck}':")
            for d in result:
                print(f"  ID {d['id']}: {display_name(d['name'])}")
            print("Use a more specific name or the deck ID.")
            return 1

        deck = result
        script_path = deck.get("source_script_path")
        if not script_path:
            dname = display_name(deck["name"])
            print(f'Deck "{dname}" has no source script tracked.')
            print("Provide a direct script path to /edit-deck instead.")
            return 1

        if getattr(args, "cat", False):
            if not os.path.exists(script_path):
                print(f"Source script not found: {script_path}")
                return 1
            with open(script_path, encoding="utf-8") as f:
                print(f.read())
        else:
            print(script_path)
        return 0

    # Default: same as list
    decks = list_decks(db_path=args.db)
    if not decks:
        print("No decks cataloged. Use 'aippt ingest' or 'aippt catalog' to add decks.")
        return 0

    total_slides = 0
    rows = []
    for d in decks:
        dname = display_name(d["name"])
        tag_count = get_deck_tag_count(d["id"], args.db)
        cataloged = (d.get("cataloged_at") or "")[:10]
        author = d.get("author", "") or "(none)"
        file_path = os.path.basename(d.get("file_path", ""))
        rows.append((d["id"], dname, d["slide_count"], author, cataloged, tag_count, file_path))
        total_slides += d["slide_count"]

    headers = ("ID", "Name", "Slides", "Author", "Cataloged", "Tags", "File")
    widths = [max(len(str(r[i])) for r in rows + [headers]) for i in range(len(headers))]
    fmt = "  ".join(f"{{:<{w}}}" if i != 2 and i != 5 else f"{{:>{w}}}" for i, w in enumerate(widths))

    print(fmt.format(*headers))
    print(fmt.format(*(chr(0x2500) * w for w in widths)))
    for row in rows:
        print(fmt.format(*[str(v) for v in row]))
    print(f"\n{len(rows)} decks, {total_slides} slides total")
    return 0


def cmd_tag(args):
    """Add tags to a slide."""
    from aippt.catalog import add_tags, get_slide_tags

    tag_names = [t.strip() for t in args.tags.split(",") if t.strip()]
    if not tag_names:
        print("No tags provided.")
        return 1

    add_tags(args.slide_id, tag_names, source="manual", db_path=args.db)
    print(f"Tagged slide {args.slide_id}: {', '.join(tag_names)}")
    return 0


def cmd_untag(args):
    """Remove tags from a slide."""
    from aippt.catalog import remove_slide_tag, remove_all_slide_tags

    if getattr(args, "all", False):
        count = remove_all_slide_tags(args.slide_id, db_path=args.db)
        print(f"Removed all tags from slide {args.slide_id} ({count} removed)")
        return 0

    tag_names = [t.strip() for t in args.tags.split(",") if t.strip()]
    if not tag_names:
        print("No tags provided. Use --all to remove all tags.")
        return 1

    removed = []
    for name in tag_names:
        if remove_slide_tag(args.slide_id, name, db_path=args.db):
            removed.append(name)

    if removed:
        print(f"Untagged slide {args.slide_id}: {', '.join(removed)}")
    else:
        print(f"No matching tags found on slide {args.slide_id}")
    return 0


def _extract_slide_text(pptx_slide, slide_title: str = "") -> str:
    """Extract all visible text from a PPTX slide object, excluding the title shape.

    The title text is already passed separately to ``analyze_slide``; this
    helper collects bullet-point and body text from the remaining shapes.

    Args:
        pptx_slide: A ``pptx.slide.Slide`` object
        slide_title: The slide title string (used to skip the title shape)

    Returns:
        Newline-joined body text, or an empty string if no body text found
    """
    lines = []
    for shape in pptx_slide.shapes:
        if not shape.has_text_frame:
            continue
        # Skip the title shape (it duplicates what we already pass as `title`)
        shape_text = shape.text_frame.text.strip()
        if shape_text and shape_text == slide_title.strip():
            continue
        for para in shape.text_frame.paragraphs:
            para_text = para.text.strip()
            if para_text:
                lines.append(para_text)
    return "\n".join(lines)


def _find_powershell():
    """Find available PowerShell executable, preferring pwsh.exe (PowerShell 7)."""
    import shutil
    for ps_cmd in ("pwsh.exe", "pwsh", "powershell.exe", "powershell"):
        path = shutil.which(ps_cmd)
        if path:
            return path
    return None


# Cache for WSL detection result -- computed once per process.
_wsl_detected = None


def _is_wsl():
    """Return True if the current process is running inside WSL."""
    global _wsl_detected
    if _wsl_detected is not None:
        return _wsl_detected
    try:
        with open("/proc/version", "r") as fh:
            content = fh.read().lower()
        _wsl_detected = "microsoft" in content or "wsl" in content
    except OSError:
        _wsl_detected = False
    return _wsl_detected


def _wsl_to_windows_path(path):
    """Convert a Linux path to a Windows UNC path using wslpath.

    Returns the converted path on success, or the original path on failure
    so that callers do not need to handle errors specially.
    """
    import subprocess
    try:
        result = subprocess.run(
            ["wslpath", "-w", path],
            capture_output=True,
            text=True,
            timeout=10,
        )
        if result.returncode == 0:
            return result.stdout.strip()
    except (OSError, FileNotFoundError):
        pass
    return path


def _is_windows_powershell(ps_exe):
    """Return True when ps_exe is a Windows-side binary (needs Windows paths)."""
    # On WSL, Windows executables are mounted under /mnt/ or end with .exe.
    if ps_exe is None:
        return False
    lower = ps_exe.lower()
    return lower.endswith(".exe") or lower.startswith("/mnt/")


def cmd_ingest(args):
    """Ingest a deck: export images, catalog, and optionally tag."""
    from aippt.ingest import ingest_deck
    from aippt.config import load_dirs_config, resolve_path

    dirs = load_dirs_config()
    base = dirs["base_dir"]
    d = dirs["directories"]

    db_path = args.db if args.db != "slides.db" else resolve_path(d["db"], base)
    images_dir = args.images_dir  # ingest_deck derives default from deck name

    total_steps = 3 if args.tags else 2
    step = 0

    def progress(stage, detail):
        nonlocal step
        if stage in ("export_images", "catalog", "tags"):
            step += 1
            print(f"\n[{step}/{total_steps}] {detail}")
        elif stage.endswith("_done") or stage.endswith("_skipped"):
            print(f"  {detail}")

    try:
        result = ingest_deck(
            deck_path=args.deck,
            db_path=db_path,
            images_dir=images_dir,
            generate_tags=args.tags,
            taxonomy=args.taxonomy,
            model=args.model,
            gateway_config=args.gateway_config,
            api_key=getattr(args, 'api_key', None),
            width=args.width,
            height=args.height,
            progress_callback=progress,
            source_script_path=getattr(args, 'source', None),
            source_theme=getattr(args, 'theme', None),
        )
    except FileNotFoundError as exc:
        logger.error(str(exc))
        return 1
    except RuntimeError as exc:
        logger.error(str(exc))
        return 1

    print(f"\n{'=' * 50}")
    print(f"INGEST COMPLETE")
    print(f"{'=' * 50}")
    print(f"  Deck: {args.deck}")
    print(f"  Deck ID: {result['deck_id']}")
    print(f"  Images: {result['images_dir']}")
    print(f"  Database: {args.db}")
    if result['tags_generated']:
        print(f"  Tags: generated")
    if result.get('source_tracked'):
        print(f"  Source: tracked")
    print(f"{'=' * 50}\n")

    return 0


def _export_images_linux(args, out_dir: str) -> int:
    """Linux branch: render via Microsoft Graph (PPTX → SP → PDF → pdftoppm)."""
    gateway_path = getattr(args, "gateway_config", None) or "gateway.yaml"
    sp_config = load_sharepoint_config(gateway_path)
    if sp_config is None:
        logger.error(
            "SharePoint render config missing. Add a 'sharepoint' block to "
            "%s with render_site_id and render_drive_id (see "
            "gateway.yaml.example).", gateway_path,
        )
        return 1

    token = (
        getattr(args, "ms_token", None)
        or graph.get_token_from_env()
    )
    if not token:
        logger.error(
            "Microsoft sign-in required. Set MS_ACCESS_TOKEN or pass "
            "--ms-token to use the Linux render pipeline.",
        )
        return 1

    ntid = (
        os.environ.get("AIPPT_USER_NTID", "").strip()
        or os.environ.get("USER", "").strip()
        or "anonymous"
    )

    try:
        render.render_pptx_to_pngs(
            pptx_path=os.path.abspath(args.deck),
            out_dir=out_dir,
            token=token,
            ntid=ntid,
            site_id=sp_config.site_id,
            drive_id=sp_config.drive_id,
            root_path=sp_config.root_path,
        )
    except graph.GraphError as exc:
        logger.error("Microsoft Graph error: %s", exc)
        return 1
    except FileNotFoundError as exc:
        logger.error("%s", exc)
        return 1
    except subprocess.CalledProcessError as exc:
        logger.error("pdftoppm failed: %s", exc)
        return 1
    return 0


def cmd_export_images(args):
    """Export slides to PNG images.

    On Linux: renders via Microsoft Graph (PPTX → SharePoint → PDF →
    pdftoppm). On Windows: drives PowerPoint via the bundled PowerShell
    script (COM automation).
    """
    from aippt.config import load_dirs_config, resolve_path

    if not os.path.exists(args.deck):
        logger.error(f"File not found: {args.deck}")
        return 1

    dirs = load_dirs_config()
    images_base = resolve_path(dirs["directories"]["images"], dirs["base_dir"])

    out_dir = args.out_dir
    if not out_dir:
        deck_name = os.path.splitext(os.path.basename(args.deck))[0]
        out_dir = os.path.join(images_base, deck_name)

    if sys.platform.startswith("linux"):
        return _export_images_linux(args, out_dir)

    # Windows / WSL+Windows-PowerShell branch
    script_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "scripts")
    ps_script = os.path.join(script_dir, "Export-SlidesToImages.ps1")

    if not os.path.exists(ps_script):
        logger.error(f"PowerShell script not found: {ps_script}")
        logger.info("Expected at: scripts/Export-SlidesToImages.ps1")
        return 1

    # Find PowerShell executable
    ps_exe = _find_powershell()
    if not ps_exe:
        logger.error("PowerShell not found. Install PowerShell 7 (pwsh) or run from Windows.")
        return 1

    deck_path = os.path.abspath(args.deck)
    out_dir_abs = os.path.abspath(out_dir)
    script_path = ps_script

    # When running under WSL with a Windows-side PowerShell binary, Linux paths
    # are not accessible to the Windows process.  Convert all three paths to
    # Windows UNC format via wslpath before passing them to PowerShell.
    if _is_wsl() and _is_windows_powershell(ps_exe):
        # Ensure the output directory exists before wslpath tries to resolve it.
        os.makedirs(out_dir_abs, exist_ok=True)

        win_deck = _wsl_to_windows_path(deck_path)
        win_out_dir = _wsl_to_windows_path(out_dir_abs)
        win_script = _wsl_to_windows_path(script_path)

        logger.debug(f"WSL path conversion: {deck_path} -> {win_deck}")
        logger.debug(f"WSL path conversion: {out_dir_abs} -> {win_out_dir}")
        logger.debug(f"WSL path conversion: {script_path} -> {win_script}")

        deck_path = win_deck
        out_dir_abs = win_out_dir
        script_path = win_script

    cmd = [
        ps_exe, "-ExecutionPolicy", "Bypass", "-File", script_path,
        "-PptxPath", deck_path,
        "-OutDir", out_dir_abs,
        "-Width", str(args.width),
        "-Height", str(args.height),
    ]

    logger.info(f"Exporting slides from {args.deck} to {out_dir}")
    result = subprocess.run(cmd, capture_output=False)
    return result.returncode


def cmd_models(args):
    """View and manage default model configuration."""
    from aippt.config import (
        load_model_config,
        save_model_config,
        get_model_registry,
        init_model_config,
        VALID_OPERATIONS,
        ConfigError,
    )

    action = getattr(args, "models_action", None)

    if action == "init":
        dest = getattr(args, "config_path", None)
        try:
            path = init_model_config(dest)
            print(f"Created models.yaml from models.yaml.example at: {path}")
        except FileExistsError:
            print("models.yaml already exists. Delete it first if you want to recreate it.")
            return 1
        except ConfigError as exc:
            print(f"Error: {exc}")
            return 1
        return 0

    if action == "set":
        if args.operation not in VALID_OPERATIONS:
            print(f"Unknown operation: {args.operation}")
            print(f"Valid operations: {', '.join(sorted(VALID_OPERATIONS))}")
            return 1
        try:
            config = load_model_config()
            registry = config["registry"]
            if args.model_name not in registry:
                print(f"Error: model '{args.model_name}' is not in the registry.")
                print("Add it to the 'registry' section of models.yaml first.")
                return 1
            config["defaults"][args.operation] = args.model_name
            save_model_config(config["defaults"])
            print(f"Set {args.operation} default to: {args.model_name}")
        except ConfigError as exc:
            print(f"Error: {exc}")
            return 1
        return 0

    if action == "reset":
        print("Error: reset is no longer supported. models.yaml is required.")
        print("Edit models.yaml directly, or delete it and run 'aippt models init' to recreate from example.")
        return 1

    if action == "list-available":
        try:
            registry = get_model_registry()
        except ConfigError as exc:
            print(f"Error: {exc}")
            return 1
        print(f"{'Provider':<12} {'Model':<30} {'Vision':<8} {'Images':<8} {'Context'}")
        for name, cfg in registry.items():
            if cfg.max_input_tokens == 0:
                ctx = "n/a"
            elif cfg.max_input_tokens < 1000000:
                ctx = f"{cfg.max_input_tokens // 1000}k"
            else:
                ctx = f"{cfg.max_input_tokens // 1000000}M"
            print(f"{cfg.provider:<12} {name:<30} {'yes' if cfg.supports_vision else 'no':<8} {'yes' if cfg.supports_images else 'no':<8} {ctx}")
        return 0

    # Default: show current config
    try:
        config = load_model_config()
    except ConfigError as exc:
        print(f"Error: {exc}")
        return 1
    max_key = max(len(k) for k in config["defaults"])
    for op, model in sorted(config["defaults"].items()):
        print(f"  {op:<{max_key + 2}} {model}")
    print(f"  Source: {config['source']}")
    return 0


def cmd_write_notes(args):
    """Write DB notes back to a PPTX file."""
    from aippt.writeback import write_notes_to_pptx, create_backup

    deck_path = args.deck
    db_path = getattr(args, "db", "slides.db")

    try:
        backup_path = create_backup(deck_path)
        print(f"Backup created: {backup_path}")
    except FileNotFoundError:
        print(f"Error: file not found: {deck_path}", file=sys.stderr)
        return 1

    try:
        result = write_notes_to_pptx(deck_path, db_path=db_path)
    except (FileNotFoundError, ValueError) as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1

    print(f"Wrote notes to {result.slides_written} of {result.slides_total} slides "
          f"({result.slides_skipped} skipped — no notes in DB)")
    for w in result.warnings:
        print(f"  Warning: {w}")
    return 0


def cmd_migrate_paths(args):
    """Convert absolute DB paths to relative."""
    from aippt.catalog import migrate_paths
    from aippt.config import load_dirs_config, resolve_path

    dirs = load_dirs_config()
    base_dir = getattr(args, "base_dir", None) or dirs["base_dir"]
    db_path = args.db if args.db != "slides.db" else resolve_path(dirs["directories"]["db"], dirs["base_dir"])
    result = migrate_paths(db_path=db_path, base_dir=base_dir)
    print(f"Migrated {result['deck_paths']} deck path(s), "
          f"{result['image_paths']} image path(s), "
          f"skipped {result['skipped']} (already relative)")
    return 0


def cmd_db_info(args):
    """Dump database schema, statistics, and content."""
    from aippt.dbinfo import dump_db_info

    result = dump_db_info(db_path=args.db, as_json=args.json)

    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(result)
        print(f"Database info written to {args.output}")
    else:
        print(result)
    return 0


def cmd_mcp(args):
    """Manage MCP server connections."""
    import asyncio

    action = getattr(args, "mcp_action", None)

    # Default action is list
    if action == "list" or action is None:
        return asyncio.run(_cmd_mcp_list(args))

    print(f"Unknown mcp action: {action}")
    return 1


async def _cmd_mcp_list(args):
    """List configured MCP servers and their available tools."""
    import json as json_mod

    from aippt.mcp import MCPManager

    config_path = getattr(args, "config", "mcp_servers.json")
    mgr = MCPManager(config_path)

    if not mgr.servers:
        print("No MCP servers configured.")
        if config_path == "mcp_servers.json":
            print("Create mcp_servers.json or copy from mcp_servers.example.json")
        return 0

    use_json = getattr(args, "json", False)
    results = {}

    for name, cfg in mgr.servers.items():
        entry = {
            "url": cfg.display_url,
            "transport": cfg.transport_type,
        }
        if cfg.auth:
            entry["auth"] = cfg.auth

        try:
            tools = await mgr.list_tools(name)
            entry["tools"] = [
                {"name": t.name, "description": t.description} for t in tools
            ]
        except Exception as e:
            entry["error"] = str(e)

        results[name] = entry

    if use_json:
        print(json_mod.dumps(results, indent=2))
        return 0

    # Human-readable output
    print("MCP Servers:\n")
    for name, info in results.items():
        transport = info["transport"]
        auth_str = f", {info['auth']}" if info.get("auth") else ""
        print(f"  {name} ({info['url']}) [{transport}{auth_str}]")

        if "error" in info:
            print(f"    Error: {info['error']}")
        elif "tools" in info:
            if info["tools"]:
                print("    Tools:")
                for tool in info["tools"]:
                    print(f"      - {tool['name']}: {tool['description']}")
            else:
                print("    No tools available")
        print()

    return 0


def cmd_merge(args):
    """Merge multiple PPTX section files into a single deck."""
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'lib'))
    from merge import merge_decks

    chunk_paths = args.chunks
    output_path = args.output
    renumber = args.renumber

    try:
        result = merge_decks(chunk_paths, output_path, renumber=renumber)
        print(f"Merged {len(chunk_paths)} files → {result['output_path']} "
              f"({result['slide_count']} slides)")
        return 0
    except (ValueError, FileNotFoundError) as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1
    except Exception as e:
        logger.error(f"Merge failed: {e}")
        return 1


def cmd_merge_template(args):
    """Merge a generated deck into a corporate template."""
    import json as json_mod
    from aippt.template_merge import merge_with_template

    if not os.path.exists(args.generated_pptx):
        logger.error(f"Generated deck not found: {args.generated_pptx}")
        return 1

    if not os.path.exists(args.corp_template):
        logger.error(f"Corporate template not found: {args.corp_template}")
        return 1

    layout_map = None
    if args.layout_map:
        if not os.path.exists(args.layout_map):
            logger.error(f"Layout map file not found: {args.layout_map}")
            return 1
        with open(args.layout_map, encoding="utf-8") as f:
            layout_map = json_mod.load(f)

    if args.dry_run:
        from pptx import Presentation as PptxPresentation
        from aippt.template_merge import _get_layout_for_slide, CORP_LAYOUT_MAP, FALLBACK_LAYOUT

        effective_map = layout_map if layout_map is not None else CORP_LAYOUT_MAP
        src_prs = PptxPresentation(args.generated_pptx)
        tgt_prs = PptxPresentation(args.corp_template)

        print(f"Dry run: {len(src_prs.slides)} slides from {args.generated_pptx}")
        print(f"Template: {args.corp_template}")
        print()
        for i, slide in enumerate(src_prs.slides, 1):
            _, source_type, target_name = _get_layout_for_slide(
                slide, effective_map, FALLBACK_LAYOUT, tgt_prs
            )
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text.strip()
                    break
            print(f"  Slide {i}: '{title}' [{source_type or 'none'}] -> {target_name}")
        return 0

    try:
        result = merge_with_template(
            args.generated_pptx, args.corp_template, args.output,
            layout_map=layout_map,
        )
        print(f"Merged {result['slide_count']} slides -> {result['output_path']}")
        for a in result["layout_assignments"]:
            print(f"  Slide {a['slide_num']}: '{a['title']}' [{a['source_layout'] or 'none'}] -> {a['target_layout']}")
        return 0
    except (FileNotFoundError, ValueError) as e:
        logger.error(str(e))
        return 1


def cmd_metadata(args):
    """Dump embedded AI metadata from a PPTX file as JSON."""
    import json
    from pptx import Presentation
    from aippt.metadata import extract_metadata

    deck_path = args.deck
    if not os.path.exists(deck_path):
        print(f"File not found: {deck_path}")
        return 1

    prs = Presentation(deck_path)
    slide_filter = getattr(args, 'slide', None)
    results = []

    for i, slide in enumerate(prs.slides, 1):
        if slide_filter and i != slide_filter:
            continue
        entries = extract_metadata(slide)
        if entries or slide_filter:
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text
            results.append({
                "slide": i,
                "title": title,
                "operations": entries,
            })

    print(json.dumps(results, indent=2))
    return 0


def build_parser():
    """Build the argument parser with all subcommands."""
    parser = argparse.ArgumentParser(
        prog="aippt",
        description="Convert markdown outlines to PowerPoint presentations.",
    )
    parser.add_argument("--debug", action="store_true", help="Enable debug logging")
    sub = parser.add_subparsers(dest="command")

    # create
    p_create = sub.add_parser("create", help="Create presentation from outline")
    p_create.add_argument("outline", help="Markdown outline file")
    p_create.add_argument("template", help="PowerPoint template file")
    p_create.add_argument("output", help="Output .pptx file")
    p_create.add_argument("--enhance", action="store_true", help="Use LLM to enhance slides")
    p_create.add_argument("--api-key", help="API key")
    p_create.add_argument("--model", default=None, help="Model to use (overrides models.yaml default)")
    p_create.add_argument("--api-base", help="Base URL for API endpoint")
    p_create.add_argument("--api-provider", choices=["anthropic", "openai", "compatible"])
    p_create.add_argument("--gateway-config", default="gateway.yaml", help="Gateway YAML config path")
    p_create.add_argument("--test", type=int, help="Process only first N slides")
    p_create.add_argument("--analyze-template", action="store_true", help="Analyze template layouts")
    p_create.add_argument("--image-gen", choices=["claude", "dalle", "openai", "mcp", "none"], default="none",
                          help="Enable diagram generation")
    p_create.add_argument("--classification", default="internal",
                          help="Content classification for MCP image generation (default: internal)")
    p_create.add_argument("--mcp-server", default="txt2img",
                          help="MCP server name for image generation (default: txt2img)")
    p_create.add_argument("--mcp-config", default="mcp_servers.json",
                          help="Path to MCP server config file (default: mcp_servers.json)")
    p_create.add_argument("--audience",
                          choices=["engineers", "executives", "product", "mixed"],
                          default=None,
                          help="Target audience (adapts content depth and language; default: mixed)")
    p_create.add_argument("--show-plan", action="store_true",
                          help="Print the deck narrative plan before enhancing (requires --enhance)")
    p_create.add_argument("--no-plan", action="store_true",
                          help="Skip deck-level narrative planning (per-slide enhancement only)")
    p_create.add_argument("--corp-template", default=None,
                          help="Path to corporate template PPTX — merge generated slides into this template as post-processing")

    # reverse
    p_reverse = sub.add_parser("reverse", help="Convert PowerPoint to markdown")
    p_reverse.add_argument("input", help="Input .pptx file")
    p_reverse.add_argument("output", nargs="?", help="Output .md file")
    p_reverse.add_argument("--no-notes", action="store_true", help="Exclude slide notes")
    p_reverse.add_argument("--strip-notes", action="store_true",
                           help="Omit speaker notes entirely from output")
    p_reverse.add_argument("--enhance", action="store_true",
                           help="Use LLM to generate high-quality outline (multimodal)")
    p_reverse.add_argument("--model", default=None,
                           help="Model to use for enhancement (overrides models.yaml)")
    p_reverse.add_argument("--gateway-config", default="gateway.yaml",
                           help="Path to gateway YAML config")
    p_reverse.add_argument("--images-dir",
                           help="Directory with pre-exported slide images (Slide1.PNG, ...)")

    # catalog
    p_catalog = sub.add_parser("catalog", help="Catalog a deck")
    p_catalog.add_argument("deck", help="PowerPoint file to catalog")
    p_catalog.add_argument("--images-dir", help="Directory with exported slide images")
    p_catalog.add_argument("--db", default="slides.db", help="Database file path")

    # search
    p_search = sub.add_parser("search", help="Search cataloged slides")
    p_search.add_argument("--tags", help="Comma-separated tags to filter by")
    p_search.add_argument("--title-contains", help="Filter by title substring")
    p_search.add_argument("--section", help="Filter by section name (substring match)")
    p_search.add_argument("--export-manifest", help="Export results as remix manifest YAML")
    p_search.add_argument("--db", default="slides.db")

    # remix
    p_remix = sub.add_parser("remix", help="Assemble deck from manifest")
    p_remix.add_argument("manifest", help="YAML manifest file")
    p_remix.add_argument("output", help="Output .pptx file")
    p_remix.add_argument("--db", default="slides.db")

    # analyze
    p_analyze = sub.add_parser("analyze", help="AI analysis of slides")
    p_analyze.add_argument("deck", help="PowerPoint file")
    p_analyze.add_argument("--mode", choices=["feedback", "notes", "tags", "improvements"], required=True)
    p_analyze.add_argument("--taxonomy", help="CSV file with pre-defined tags")
    p_analyze.add_argument("--model", default=None, help="Model to use (overrides models.yaml default)")
    p_analyze.add_argument("--api-key")
    p_analyze.add_argument("--gateway-config", default="gateway.yaml")
    p_analyze.add_argument("--images-dir", help="Directory with slide images")
    p_analyze.add_argument("--db", default="slides.db")

    # export
    p_export = sub.add_parser("export", help="Export metadata to CSV")
    p_export.add_argument("deck", nargs="?", help="Specific deck to export")
    p_export.add_argument("--all", action="store_true", help="Export all cataloged decks")
    p_export.add_argument("--output", default="slides.csv")
    p_export.add_argument("--db", default="slides.db")

    # serve
    p_serve = sub.add_parser("serve", help="Launch web UI")
    p_serve.add_argument("--host", default="127.0.0.1", help="Bind address (default: 127.0.0.1)")
    p_serve.add_argument("--port", type=int, default=8000)
    p_serve.add_argument("--db", default="slides.db")
    p_serve.add_argument("--gateway-config", default="gateway.yaml", help="Gateway YAML config path for LLM access")
    p_serve.add_argument("--uploads-dir", default="uploads", help="Directory for uploaded files (default: uploads)")
    p_serve.add_argument("--view-only", action="store_true", help="Disable LLM features (also settable via AIPPT_VIEW_ONLY env var; auto-detected when no gateway/API keys)")

    # tags (taxonomy management)
    p_tags = sub.add_parser("tags", help="Manage taxonomy of predefined tags")
    p_tags.add_argument("--db", default="slides.db")
    tags_sub = p_tags.add_subparsers(dest="tags_action")

    p_tags_add = tags_sub.add_parser("add", help="Add a tag to the taxonomy")
    p_tags_add.add_argument("tag", help="Tag name")
    p_tags_add.add_argument("--category", default="", help="Category for the tag")
    p_tags_add.add_argument("--db", default="slides.db")

    p_tags_remove = tags_sub.add_parser("remove", help="Remove a tag from the taxonomy")
    p_tags_remove.add_argument("tag", help="Tag name to remove")
    p_tags_remove.add_argument("--db", default="slides.db")

    p_tags_import = tags_sub.add_parser("import", help="Import taxonomy from CSV")
    p_tags_import.add_argument("csv_file", help="CSV file to import")
    p_tags_import.add_argument("--db", default="slides.db")

    p_tags_export = tags_sub.add_parser("export", help="Export taxonomy to CSV")
    p_tags_export.add_argument("csv_file", help="Output CSV file")
    p_tags_export.add_argument("--db", default="slides.db")

    p_tags_rename = tags_sub.add_parser("rename", help="Rename a taxonomy tag")
    p_tags_rename.add_argument("old_name", help="Current tag name")
    p_tags_rename.add_argument("new_name", help="New tag name")
    p_tags_rename.add_argument("--db", default="slides.db")

    # tag (add tags to a slide)
    p_tag = sub.add_parser("tag", help="Add tags to a slide")
    p_tag.add_argument("slide_id", type=int, help="Slide ID")
    p_tag.add_argument("tags", help="Comma-separated tag names")
    p_tag.add_argument("--db", default="slides.db")

    # untag (remove tags from a slide)
    p_untag = sub.add_parser("untag", help="Remove tags from a slide")
    p_untag.add_argument("slide_id", type=int, help="Slide ID")
    p_untag.add_argument("tags", nargs="?", default="", help="Comma-separated tag names")
    p_untag.add_argument("--all", action="store_true", help="Remove all tags from the slide")
    p_untag.add_argument("--db", default="slides.db")

    # decks (deck management)
    p_decks = sub.add_parser("decks", help="Manage cataloged decks")
    p_decks.add_argument("--db", default="slides.db")
    decks_sub = p_decks.add_subparsers(dest="decks_action")

    p_decks_list = decks_sub.add_parser("list", help="List all cataloged decks")
    p_decks_list.add_argument("--db", default="slides.db")
    p_decks_list.add_argument("--json", action="store_true", help="Output as JSON")

    p_decks_info = decks_sub.add_parser("info", help="Show detailed deck information")
    p_decks_info.add_argument("deck", help="Deck ID or name substring")
    p_decks_info.add_argument("--db", default="slides.db")
    p_decks_info.add_argument("--json", action="store_true", help="Output as JSON")

    p_decks_rename = decks_sub.add_parser("rename", help="Rename a deck")
    p_decks_rename.add_argument("deck", help="Deck ID or name substring")
    p_decks_rename.add_argument("new_name", help="New display name")
    p_decks_rename.add_argument("--db", default="slides.db")

    p_decks_delete = decks_sub.add_parser("delete", help="Delete a deck and all associated data")
    p_decks_delete.add_argument("deck", help="Deck ID or name substring")
    p_decks_delete.add_argument("--db", default="slides.db")
    p_decks_delete.add_argument("--force", action="store_true", help="Skip confirmation prompt")
    p_decks_delete.add_argument("--purge-images", action="store_true", help="Also delete the image directory")

    p_decks_source = decks_sub.add_parser("source", help="Show source script path for a deck")
    p_decks_source.add_argument("deck", help="Deck ID or name substring")
    p_decks_source.add_argument("--db", default="slides.db")
    p_decks_source.add_argument("--cat", action="store_true", help="Print the script contents")

    # ingest
    p_ingest = sub.add_parser("ingest", help="Ingest a deck: export images, catalog, and optionally tag")
    p_ingest.add_argument("deck", help="PowerPoint file to ingest")
    p_ingest.add_argument("--images-dir", default=None,
                          help="Output directory for slide images (default: images/<deck-name>/)")
    p_ingest.add_argument("--db", default="slides.db", help="Database file path")
    p_ingest.add_argument("--tags", action="store_true", help="Generate AI tags after cataloging")
    p_ingest.add_argument("--taxonomy", help="CSV file for taxonomy-constrained tagging")
    p_ingest.add_argument("--model", default=None, help="Model to use for tag generation")
    p_ingest.add_argument("--gateway-config", default="gateway.yaml", help="Gateway YAML config path")
    p_ingest.add_argument("--api-key", default=None, help="API key for LLM provider")
    p_ingest.add_argument("--width", type=int, default=1920, help="Image export width (default: 1920)")
    p_ingest.add_argument("--height", type=int, default=1080, help="Image export height (default: 1080)")
    p_ingest.add_argument("--source", default=None, help="Path to generating script (JS/Python) for source tracking")
    p_ingest.add_argument("--theme", default=None, help="Theme name (overrides auto-detection from script)")

    # export-images
    p_eimg = sub.add_parser(
        "export-images",
        help="Export slides to PNG images (PowerPoint COM on Windows, Microsoft Graph on Linux)",
    )
    p_eimg.add_argument("deck", help="PowerPoint file to export")
    p_eimg.add_argument("out_dir", nargs="?", default=None,
                        help="Output directory (default: images/<deck-name>/)")
    p_eimg.add_argument("--width", type=int, default=1920, help="Image width in pixels (default: 1920, Windows only)")
    p_eimg.add_argument("--height", type=int, default=1080, help="Image height in pixels (default: 1080, Windows only)")
    p_eimg.add_argument("--ms-token", default=None,
                        help="Microsoft Graph access token (Linux render path; falls back to MS_ACCESS_TOKEN env)")
    p_eimg.add_argument("--gateway-config", default="gateway.yaml",
                        help="Path to gateway.yaml for sharepoint config (Linux render path)")

    # improve
    p_improve = sub.add_parser("improve", help="Improve slides using LLM analysis and rewrite")
    p_improve.add_argument("deck", help="PowerPoint file to improve")
    p_improve.add_argument("--output", default=None, help="Save to different file (default: overwrite)")
    p_improve.add_argument("--dry-run", action="store_true", help="Show changes without modifying")
    p_improve.add_argument("--slides", default=None, help="Comma-separated slide numbers to improve")
    p_improve.add_argument("--passes", type=int, default=1, help="Number of improvement passes")
    p_improve.add_argument("--focus", choices=["accuracy", "detail", "brevity", "structure", "general"],
                           default="general", help="Focus area for improvements (default: general)")
    p_improve.add_argument("--images-dir", default=None, help="Slide images directory")
    p_improve.add_argument("--model", default=None, help="Model for rewrite")
    p_improve.add_argument("--gateway-config", default="gateway.yaml", help="Gateway config path")
    p_improve.add_argument("--api-key", default=None, help="API key")
    p_improve.add_argument("--db", default="slides.db", help="Database path")
    p_improve.add_argument("--audience",
                           choices=["engineers", "executives", "product", "mixed"],
                           default=None,
                           help="Target audience (adapts rewrite prompts; default: mixed)")
    p_improve.add_argument("--keep-titles", action="store_true",
                           help="Skip title rewriting; improve body content only")
    p_improve.add_argument("--max-retries", type=int, default=2,
                           help="Max validation retries per slide (default: 2)")
    p_improve.add_argument("--no-validate", action="store_true",
                           help="Skip validation pass (revert to single-pass behavior)")

    # models
    p_models = sub.add_parser("models", help="View and manage default model configuration")
    models_sub = p_models.add_subparsers(dest="models_action")
    models_sub.add_parser("init", help="Create models.yaml from models.yaml.example")
    p_models_set = models_sub.add_parser("set", help="Set default model for an operation")
    p_models_set.add_argument("operation", help="Operation name (enhance, feedback, notes, tags, image, improve)")
    p_models_set.add_argument("model_name", help="Model name to set as default")
    models_sub.add_parser("reset", help="(Deprecated) Reset all defaults to built-in values")
    models_sub.add_parser("list-available", help="Show all models in the registry")

    p_write_notes = sub.add_parser("write-notes", help="Write DB notes back to PPTX file")
    p_write_notes.add_argument("deck", help="Path to the PPTX file")
    p_write_notes.add_argument("--db", default="slides.db", help="Path to the SQLite database")

    # db-info
    p_dbinfo = sub.add_parser("db-info", help="Dump database schema, statistics, and content")
    p_dbinfo.add_argument("--db", default="slides.db", help="Path to the SQLite database")
    p_dbinfo.add_argument("--json", action="store_true", help="Output as JSON instead of plain text")
    p_dbinfo.add_argument("--output", help="Write output to a file instead of stdout")
    p_dbinfo.set_defaults(func=cmd_db_info)

    # migrate-paths
    p_migrate = sub.add_parser("migrate-paths", help="Convert absolute DB paths to relative")
    p_migrate.add_argument("--db", default="slides.db", help="Path to the SQLite database")
    p_migrate.add_argument("--base-dir", default=None,
                           help="Base directory for relative paths (default: current directory)")

    # mcp (MCP server management)
    p_mcp = sub.add_parser("mcp", help="Manage MCP server connections")
    p_mcp.add_argument("--config", default="mcp_servers.json", help="MCP servers config file")
    mcp_sub = p_mcp.add_subparsers(dest="mcp_action")

    p_mcp_list = mcp_sub.add_parser("list", help="List configured MCP servers and their tools")
    p_mcp_list.add_argument("--config", default="mcp_servers.json", help="MCP servers config file")
    p_mcp_list.add_argument("--json", action="store_true", help="Output as JSON")

    # merge
    p_merge = sub.add_parser("merge", help="Merge multiple PPTX section files into one deck")
    p_merge.add_argument("chunks", nargs="+", help="PPTX files to merge, in order")
    p_merge.add_argument("-o", "--output", required=True, help="Output file path")
    p_merge.add_argument("--no-renumber", dest="renumber", action="store_false",
                         default=True, help="Skip slide number renumbering")

    # metadata
    p_metadata = sub.add_parser("metadata", help="Extract AI metadata from PPTX speaker notes")
    p_metadata.add_argument("deck", help="Path to PPTX file")
    p_metadata.add_argument("--slide", type=int, default=None, help="Show metadata for specific slide number")

    # merge-template
    p_merge_tpl = sub.add_parser("merge-template", help="Merge generated deck into corporate template")
    p_merge_tpl.add_argument("generated_pptx", help="Path to the generated PPTX deck")
    p_merge_tpl.add_argument("--corp-template", required=True, help="Path to corporate template PPTX")
    p_merge_tpl.add_argument("-o", "--output", required=True, help="Output file path")
    p_merge_tpl.add_argument("--layout-map", default=None, help="JSON file overriding default layout map")
    p_merge_tpl.add_argument("--dry-run", action="store_true", help="Print layout assignments without writing")

    return parser


def main():
    """Main entry point."""
    # Ensure stdout/stderr can handle Unicode (LLM responses often contain
    # arrows, checkmarks, etc. that cp1252 cannot encode on Windows).
    if sys.stdout.encoding and sys.stdout.encoding.lower().replace("-", "") != "utf8":
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

    parser = build_parser()
    args = parser.parse_args()

    if args.debug:
        logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
    else:
        logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

    commands = {
        "create": cmd_create,
        "reverse": cmd_reverse,
        "catalog": cmd_catalog,
        "search": cmd_search,
        "remix": cmd_remix,
        "analyze": cmd_analyze,
        "export": cmd_export,
        "serve": cmd_serve,
        "tags": cmd_tags,
        "tag": cmd_tag,
        "untag": cmd_untag,
        "decks": cmd_decks,
        "ingest": cmd_ingest,
        "export-images": cmd_export_images,
        "improve": cmd_improve,
        "models": cmd_models,
        "write-notes": cmd_write_notes,
        "db-info": cmd_db_info,
        "migrate-paths": cmd_migrate_paths,
        "mcp": cmd_mcp,
        "merge": cmd_merge,
        "metadata": cmd_metadata,
        "merge-template": cmd_merge_template,
    }

    if not args.command:
        parser.print_help()
        return 1

    return commands[args.command](args)


if __name__ == "__main__":
    sys.exit(main() or 0)

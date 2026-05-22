"""Structured metadata blocks in PPTX speaker notes.

Embeds [AIPPT-META] JSON blocks in slide notes to track AI operations
(enhance, improve, image generation) with timestamps and parameters,
as well as source lineage and editing history.
"""

import json
import hashlib
import logging
import re
from datetime import datetime, timezone

logger = logging.getLogger(__name__)

META_START = "[AIPPT-META]"
META_END = "[/AIPPT-META]"
SEPARATOR = "\n\n---\n"

_META_PATTERN = re.compile(
    re.escape(META_START) + r"\n(.*?)\n" + re.escape(META_END),
    re.DOTALL,
)


def format_notes_with_metadata(notes_text: str, metadata_entries: list[dict]) -> str:
    """Combine human-readable notes with a metadata block."""
    if not metadata_entries:
        return notes_text
    meta_block = f"{META_START}\n{json.dumps(metadata_entries, indent=2)}\n{META_END}"
    if notes_text.strip():
        return f"{notes_text}{SEPARATOR}{meta_block}"
    return meta_block


def extract_metadata(slide) -> list[dict]:
    """Extract all metadata entries from a slide's speaker notes."""
    try:
        notes = slide.notes_slide.notes_text_frame.text
    except Exception:
        return []
    if not notes or META_START not in notes:
        return []
    match = _META_PATTERN.search(notes)
    if not match:
        return []
    try:
        return json.loads(match.group(1))
    except (json.JSONDecodeError, TypeError):
        logger.warning("Malformed metadata JSON in slide notes; returning empty list")
        return []


def extract_notes_text(slide) -> str:
    """Extract just the human-readable notes (before the --- separator)."""
    try:
        notes = slide.notes_slide.notes_text_frame.text
    except Exception:
        return ""
    if not notes:
        return ""
    # If notes start directly with the metadata block (no human notes)
    stripped = notes.lstrip()
    if stripped.startswith(META_START):
        return ""
    # Split on the separator that precedes the metadata block
    sep_idx = notes.find(SEPARATOR.lstrip("\n"))
    if sep_idx >= 0 and META_START in notes[sep_idx:]:
        return notes[:sep_idx].rstrip()
    return notes.rstrip()


def append_metadata(slide, operation: str, **kwargs) -> None:
    """Append a metadata entry to the slide's speaker notes."""
    entry = {
        "operation": operation,
        "timestamp": datetime.now(timezone.utc).isoformat(),
        **kwargs,
    }
    existing_entries = extract_metadata(slide)
    existing_entries.append(entry)
    human_notes = extract_notes_text(slide)
    slide.notes_slide.notes_text_frame.text = format_notes_with_metadata(
        human_notes, existing_entries
    )


def content_hash(text: str) -> str:
    """SHA-256 hash of content text (first 12 hex chars)."""
    return hashlib.sha256(text.encode()).hexdigest()[:12]


# ---------------------------------------------------------------------------
# Source lineage helpers
# ---------------------------------------------------------------------------

HISTORY_CAP = 10


def create_lineage_entry(
    source: str,
    layout: str | None = None,
    theme: str | None = None,
    created: str | None = None,
) -> dict:
    """Build the initial 'create' metadata entry with lineage fields.

    Args:
        source: Lineage string, e.g. "outline -> pptxgenjs"
        layout: Slide layout type (e.g. "bullet", "two_column")
        theme: Theme name (e.g. "amd", "default")
        created: ISO date string; defaults to today

    Returns:
        Metadata entry dict suitable for append_metadata().
    """
    entry = {
        "source": source,
        "created": created or datetime.now(timezone.utc).strftime("%Y-%m-%d"),
        "history": [
            f"{datetime.now(timezone.utc).strftime('%Y-%m-%d')}: "
            f"Created from {source}"
            + (f" ({layout} layout)" if layout else ""),
        ],
    }
    if layout:
        entry["layout"] = layout
    if theme:
        entry["theme"] = theme
    return entry


def append_history_entry(slide, entry_text: str, source_tag: str | None = None) -> None:
    """Add a history line to the slide's metadata.

    Finds the most recent metadata entry with a ``history`` list and appends
    the new line.  If no such entry exists, a new metadata entry is created.
    History is capped at HISTORY_CAP entries (oldest are trimmed).

    Args:
        slide: A python-pptx slide object.
        entry_text: Human-readable description of the change.
        source_tag: Optional tag appended in brackets, e.g. "/edit-deck".
    """
    datestamp = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    line = f"{datestamp}: {entry_text}"
    if source_tag:
        line += f" [{source_tag}]"

    entries = extract_metadata(slide)

    # Find last entry with a history list
    target = None
    for e in reversed(entries):
        if "history" in e and isinstance(e["history"], list):
            target = e
            break

    if target is None:
        # No existing lineage entry — create a minimal one
        entries.append({
            "operation": "edit",
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "history": [line],
        })
    else:
        target["history"].append(line)
        # Cap at HISTORY_CAP
        if len(target["history"]) > HISTORY_CAP:
            target["history"] = target["history"][-HISTORY_CAP:]

    human_notes = extract_notes_text(slide)
    slide.notes_slide.notes_text_frame.text = format_notes_with_metadata(
        human_notes, entries
    )


def write_deck_lineage(
    pptx_path: str,
    source: str,
    engine: str | None = None,
    theme: str | None = None,
    generated_at: str | None = None,
) -> None:
    """Open a PPTX file, embed an [AIPPT-META] lineage entry on slide 1, and save.

    This is the canonical call site for recording deck origin after the pipeline
    produces the PPTX file.  It is idempotent: calling it again on a deck that
    already has a lineage entry replaces the existing entry (by appending a new
    one alongside existing metadata entries).

    Args:
        pptx_path: Absolute or CWD-relative path to the .pptx file.
        source: Lineage string, e.g. "outline -> pptxgenjs".
        engine: Engine used ('pptxgenjs', 'python-pptx', etc.).
        theme: Theme name ('amd', 'default', etc.).
        generated_at: ISO timestamp; defaults to now.
    """
    from pptx import Presentation

    generated_at = generated_at or datetime.now(timezone.utc).isoformat()

    try:
        prs = Presentation(pptx_path)
    except Exception as exc:
        logger.warning("write_deck_lineage: could not open %s: %s", pptx_path, exc)
        return

    if not prs.slides:
        logger.warning("write_deck_lineage: presentation has no slides; skipping")
        return

    slide = prs.slides[0]
    lineage_entry = create_lineage_entry(
        source=source,
        theme=theme,
        created=generated_at[:10] if generated_at else None,
    )
    # Add engine and generated_at explicitly so they surface in the JSON block
    lineage_entry["engine"] = engine
    lineage_entry["generated_at"] = generated_at

    existing_entries = extract_metadata(slide)
    existing_entries.append(lineage_entry)
    human_notes = extract_notes_text(slide)
    try:
        slide.notes_slide.notes_text_frame.text = format_notes_with_metadata(
            human_notes, existing_entries
        )
        prs.save(pptx_path)
    except Exception as exc:
        logger.warning("write_deck_lineage: could not save %s: %s", pptx_path, exc)


def get_slide_lineage(slide) -> dict:
    """Extract source lineage information from a slide's metadata.

    Returns a dict with keys: source, created, layout, theme, history.
    Missing fields are returned as None (history defaults to []).
    """
    entries = extract_metadata(slide)
    result = {
        "source": None,
        "created": None,
        "layout": None,
        "theme": None,
        "history": [],
    }
    for entry in entries:
        if "source" in entry and result["source"] is None:
            result["source"] = entry["source"]
        if "created" in entry and result["created"] is None:
            result["created"] = entry["created"]
        # layout/theme: take the latest value (later entries override earlier)
        if "layout" in entry:
            result["layout"] = entry["layout"]
        if "theme" in entry:
            result["theme"] = entry["theme"]
        # history: take the last-seen history list
        if "history" in entry and isinstance(entry["history"], list):
            result["history"] = entry["history"]
    return result

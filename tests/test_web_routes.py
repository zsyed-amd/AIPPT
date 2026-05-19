"""Tests for web API route handlers."""
import os
from unittest.mock import patch

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from aippt.catalog import catalog_deck, get_db, record_edit, add_tags, add_taxonomy_tags
from aippt.web.app import create_app


@pytest.fixture
def deck_path(tmp_path):
    """Create a minimal PPTX with one slide that has notes."""
    prs = Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Test Slide"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Original PPTX notes"
    path = str(tmp_path / "test.pptx")
    prs.save(path)
    return path


@pytest.fixture
def client(tmp_path, deck_path):
    """Create a TestClient with a cataloged deck."""
    db_path = str(tmp_path / "test.db")
    uploads_dir = str(tmp_path / "uploads")
    catalog_deck(deck_path, db_path=db_path)
    app = create_app(db_path=db_path, uploads_dir=uploads_dir)
    return TestClient(app)


@pytest.fixture
def db_path(client):
    """Extract db_path from the test client's app state."""
    return client.app.state.db_path


class TestSaveNotesWithHistory:
    """POST /api/slides/{id}/notes/save should record edit history."""

    def test_save_notes_creates_history_row(self, client, db_path):
        resp = client.post("/api/slides/1/notes/save", json={"notes": "New web notes"})
        assert resp.status_code == 200
        assert resp.json()["status"] == "ok"

        conn = get_db(db_path)
        hist = conn.execute(
            "SELECT * FROM edit_history WHERE slide_id = 1 AND field = 'notes'"
        ).fetchall()
        assert len(hist) == 1
        assert hist[0]["old_value"] == "Original PPTX notes"
        assert hist[0]["new_value"] == "New web notes"
        assert hist[0]["source"] == "web"
        conn.close()

    def test_save_notes_updates_slide_notes(self, client, db_path):
        client.post("/api/slides/1/notes/save", json={"notes": "Updated notes"})
        conn = get_db(db_path)
        row = conn.execute("SELECT notes FROM slides WHERE id = 1").fetchone()
        assert row["notes"] == "Updated notes"
        conn.close()

    def test_save_notes_updates_timestamp(self, client, db_path):
        conn = get_db(db_path)
        before = conn.execute("SELECT updated_at FROM slides WHERE id = 1").fetchone()["updated_at"]
        conn.close()

        client.post("/api/slides/1/notes/save", json={"notes": "Timestamped notes"})

        conn = get_db(db_path)
        after = conn.execute("SELECT updated_at FROM slides WHERE id = 1").fetchone()["updated_at"]
        conn.close()
        assert after >= before


class TestSaveNotesSameValue:
    """Saving identical notes should not create a history row."""

    def test_no_history_for_same_value(self, client, db_path):
        resp = client.post(
            "/api/slides/1/notes/save",
            json={"notes": "Original PPTX notes"},
        )
        assert resp.status_code == 200

        conn = get_db(db_path)
        count = conn.execute("SELECT COUNT(*) as cnt FROM edit_history").fetchone()["cnt"]
        assert count == 0
        conn.close()


class TestSaveNotesValidation:
    """Edge cases for notes save endpoint."""

    def test_empty_notes_rejected(self, client):
        resp = client.post("/api/slides/1/notes/save", json={"notes": ""})
        assert resp.status_code == 400

    def test_missing_slide_404(self, client):
        resp = client.post("/api/slides/9999/notes/save", json={"notes": "test"})
        assert resp.status_code == 404


class TestSaveNotesSource:
    """Verify source field is passed through correctly."""

    def test_ai_source_recorded(self, client, db_path):
        resp = client.post(
            "/api/slides/1/notes/save",
            json={"notes": "AI generated notes", "source": "ai"},
        )
        assert resp.status_code == 200

        conn = get_db(db_path)
        hist = conn.execute(
            "SELECT source FROM edit_history WHERE slide_id = 1"
        ).fetchone()
        assert hist["source"] == "ai"
        conn.close()

    def test_default_source_is_web(self, client, db_path):
        client.post("/api/slides/1/notes/save", json={"notes": "Manual edit"})
        conn = get_db(db_path)
        hist = conn.execute(
            "SELECT source FROM edit_history WHERE slide_id = 1"
        ).fetchone()
        assert hist["source"] == "web"
        conn.close()


class TestNotesHistory:
    """GET /api/slides/{id}/notes/history returns edit history."""

    def test_history_returns_entries_reverse_chronological(self, client, db_path):
        client.post("/api/slides/1/notes/save", json={"notes": "Edit 1"})
        client.post("/api/slides/1/notes/save", json={"notes": "Edit 2"})

        resp = client.get("/api/slides/1/notes/history")
        assert resp.status_code == 200
        entries = resp.json()["history"]
        assert len(entries) == 2
        # Newest first
        assert entries[0]["new_value"] == "Edit 2"
        assert entries[0]["old_value"] == "Edit 1"
        assert entries[1]["new_value"] == "Edit 1"
        assert entries[1]["old_value"] == "Original PPTX notes"

    def test_history_includes_source_and_timestamp(self, client, db_path):
        client.post("/api/slides/1/notes/save", json={"notes": "Web edit"})
        resp = client.get("/api/slides/1/notes/history")
        entry = resp.json()["history"][0]
        assert entry["source"] == "web"
        assert "created_at" in entry


class TestNotesHistoryEmpty:
    """GET /api/slides/{id}/notes/history with no edits."""

    def test_empty_history(self, client):
        resp = client.get("/api/slides/1/notes/history")
        assert resp.status_code == 200
        assert resp.json()["history"] == []

    def test_history_missing_slide_404(self, client):
        resp = client.get("/api/slides/9999/notes/history")
        assert resp.status_code == 404


class TestWriteNotesEndpoint:
    def test_writes_notes_to_pptx(self, client, db_path, deck_path):
        record_edit(1, "notes", "Web write-back test", source="web", db_path=db_path)

        resp = client.post("/api/decks/1/write-notes")
        assert resp.status_code == 200
        data = resp.json()
        assert data["slides_written"] == 1
        assert data["backup_path"] is not None

        prs = Presentation(deck_path)
        assert prs.slides[0].notes_slide.notes_text_frame.text == "Web write-back test"

    def test_deck_not_found(self, client):
        resp = client.post("/api/decks/999/write-notes")
        assert resp.status_code == 404

    def test_slide_count_mismatch(self, client, deck_path, db_path):
        prs = Presentation(deck_path)
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(deck_path)

        resp = client.post("/api/decks/1/write-notes")
        assert resp.status_code == 409
        assert "mismatch" in resp.json()["error"].lower()


class TestDownloadWithNotes:
    def test_download_applies_db_notes(self, client, db_path, deck_path):
        record_edit(1, "notes", "Download test notes", source="web", db_path=db_path)

        resp = client.get("/api/decks/1/download")
        assert resp.status_code == 200

        import tempfile as tf
        with tf.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(resp.content)
            tmp_path = f.name

        prs = Presentation(tmp_path)
        assert prs.slides[0].notes_slide.notes_text_frame.text == "Download test notes"
        os.unlink(tmp_path)

        # Original file should be untouched
        prs_orig = Presentation(deck_path)
        assert prs_orig.slides[0].notes_slide.notes_text_frame.text == "Original PPTX notes"

    def test_download_still_works_without_edits(self, client, deck_path):
        resp = client.get("/api/decks/1/download")
        assert resp.status_code == 200

        import tempfile as tf
        with tf.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(resp.content)
            tmp_path = f.name

        prs = Presentation(tmp_path)
        assert prs.slides[0].notes_slide.notes_text_frame.text == "Original PPTX notes"
        os.unlink(tmp_path)


class TestTemplateEndpoints:
    def test_get_templates_returns_503_when_missing(self, client, tmp_path, monkeypatch):
        monkeypatch.setattr(
            "aippt.config.DEFAULT_TEMPLATE_CONFIG_PATH",
            str(tmp_path / "nonexistent.yaml"),
        )
        resp = client.get("/api/templates")
        assert resp.status_code == 503

    def test_put_templates_updates_config(self, client, tmp_path, monkeypatch):
        config_path = str(tmp_path / "templates.yaml")
        (tmp_path / "templates.yaml").write_text("default_template: old.pptx\n")
        monkeypatch.setattr(
            "aippt.config.DEFAULT_TEMPLATE_CONFIG_PATH", config_path
        )
        resp = client.put(
            "/api/templates",
            json={"default_template": "new/template.pptx"},
        )
        assert resp.status_code == 200
        assert resp.json()["default_template"] == "new/template.pptx"

    def test_put_templates_rejects_empty(self, client):
        resp = client.put("/api/templates", json={"default_template": ""})
        assert resp.status_code == 400


class TestCreateDeckEndpoint:
    @pytest.fixture(autouse=True)
    def _relax_image_requirement(self, monkeypatch):
        """These tests run on Linux but cannot reach SharePoint. The strict
        Graph render path is exercised by TestCreateDeckHardening; here we
        verify the create/save plumbing without requiring real image export."""
        monkeypatch.setattr(
            "aippt.web.routes._require_images_for_render", lambda: False
        )

    def test_create_with_outline_text_no_enhance(self, client, tmp_path, monkeypatch):
        from pptx import Presentation
        template_path = str(tmp_path / "template.pptx")
        prs = Presentation()
        prs.save(template_path)

        config_path = str(tmp_path / "templates.yaml")
        (tmp_path / "templates.yaml").write_text(f"default_template: {template_path}\n")
        monkeypatch.setattr(
            "aippt.config.DEFAULT_TEMPLATE_CONFIG_PATH", config_path
        )

        resp = client.post(
            "/api/decks/create",
            data={
                "outline_text": "# Test\n## Slide 1: Hello\n- World\n",
                "enhance": "false",
            },
            headers={"Authorization": "Bearer tok"},
        )
        assert resp.status_code == 200
        assert resp.headers["content-type"].startswith("text/event-stream")
        body = resp.text
        assert "event: complete" in body

    def test_create_rejects_empty_input(self, client):
        resp = client.post(
            "/api/decks/create",
            data={},
            headers={"Authorization": "Bearer tok"},
        )
        assert resp.status_code == 400

    def test_create_with_md_file_upload(self, client, tmp_path, monkeypatch):
        from pptx import Presentation
        template_path = str(tmp_path / "template.pptx")
        prs = Presentation()
        prs.save(template_path)

        config_path = str(tmp_path / "templates.yaml")
        (tmp_path / "templates.yaml").write_text(f"default_template: {template_path}\n")
        monkeypatch.setattr(
            "aippt.config.DEFAULT_TEMPLATE_CONFIG_PATH", config_path
        )

        md_content = "# File Test\n## Slide 1: From File\n- Content\n"

        resp = client.post(
            "/api/decks/create",
            data={"enhance": "false"},
            files={"outline_file": ("outline.md", md_content, "text/markdown")},
            headers={"Authorization": "Bearer tok"},
        )
        assert resp.status_code == 200
        body = resp.text
        assert "event: complete" in body

    def test_create_saves_outline_to_disk(self, client, tmp_path, monkeypatch):
        """Outline text should be saved to uploads/ for IMAGE: resolution."""
        from pptx import Presentation
        template_path = str(tmp_path / "template.pptx")
        prs = Presentation()
        prs.save(template_path)

        config_path = str(tmp_path / "templates.yaml")
        (tmp_path / "templates.yaml").write_text(f"default_template: {template_path}\n")
        monkeypatch.setattr(
            "aippt.config.DEFAULT_TEMPLATE_CONFIG_PATH", config_path
        )

        md_content = "# Saved\n## Slide 1: Check\n- Content\n"
        resp = client.post(
            "/api/decks/create",
            data={"outline_text": md_content, "enhance": "false"},
            headers={"Authorization": "Bearer tok"},
        )
        assert resp.status_code == 200

        # Verify an outline .md file was saved to uploads/
        uploads_dir = client.app.state.uploads_dir
        md_files = [f for f in os.listdir(uploads_dir) if f.endswith(".md")]
        assert len(md_files) >= 1
        saved_content = open(os.path.join(uploads_dir, md_files[0])).read()
        assert saved_content == md_content

    def test_create_with_image_upload(self, client, tmp_path, monkeypatch):
        """Uploaded images should be saved relative to the outline for IMAGE: resolution."""
        from pptx import Presentation
        template_path = str(tmp_path / "template.pptx")
        prs = Presentation()
        prs.save(template_path)

        config_path = str(tmp_path / "templates.yaml")
        (tmp_path / "templates.yaml").write_text(f"default_template: {template_path}\n")
        monkeypatch.setattr(
            "aippt.config.DEFAULT_TEMPLATE_CONFIG_PATH", config_path
        )

        # Create a tiny 1x1 PNG for the image upload
        import struct, zlib
        def _make_png():
            sig = b'\x89PNG\r\n\x1a\n'
            ihdr_data = struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0)
            ihdr_crc = zlib.crc32(b'IHDR' + ihdr_data) & 0xffffffff
            ihdr = struct.pack('>I', 13) + b'IHDR' + ihdr_data + struct.pack('>I', ihdr_crc)
            raw = b'\x00\x00\x00\x00'
            idat_data = zlib.compress(raw)
            idat_crc = zlib.crc32(b'IDAT' + idat_data) & 0xffffffff
            idat = struct.pack('>I', len(idat_data)) + b'IDAT' + idat_data + struct.pack('>I', idat_crc)
            iend_crc = zlib.crc32(b'IEND') & 0xffffffff
            iend = struct.pack('>I', 0) + b'IEND' + struct.pack('>I', iend_crc)
            return sig + ihdr + idat + iend
        test_png = _make_png()

        md_content = "# Image Test\n## Slide 1: Photo\nIMAGE: imgs/photo.png\n- Caption\n"
        resp = client.post(
            "/api/decks/create",
            data={"outline_text": md_content, "enhance": "false"},
            files=[("image_files", ("photo.png", test_png, "image/png"))],
            headers={"Authorization": "Bearer tok"},
        )
        assert resp.status_code == 200
        assert "event: complete" in resp.text

        # Verify image was saved to the correct subdirectory
        uploads_dir = client.app.state.uploads_dir
        md_files = [f for f in os.listdir(uploads_dir) if f.endswith(".md")]
        assert len(md_files) >= 1
        outline_dir = os.path.dirname(os.path.join(uploads_dir, md_files[0]))
        img_path = os.path.join(outline_dir, "imgs", "photo.png")
        assert os.path.exists(img_path)
        assert open(img_path, 'rb').read() == test_png

    def test_create_passes_outline_path_to_create_deck(self, client, tmp_path, monkeypatch):
        """run_pipeline() should receive outline_path for IMAGE: resolution."""
        from pptx import Presentation
        from unittest.mock import patch
        template_path = str(tmp_path / "template.pptx")
        prs = Presentation()
        prs.save(template_path)

        config_path = str(tmp_path / "templates.yaml")
        (tmp_path / "templates.yaml").write_text(f"default_template: {template_path}\n")
        monkeypatch.setattr(
            "aippt.config.DEFAULT_TEMPLATE_CONFIG_PATH", config_path
        )

        captured_configs = []
        original_run = None

        # Import and capture the real run_pipeline
        from aippt import pipeline as pipeline_module
        original_run = pipeline_module.run_pipeline

        def spy_run_pipeline(config):
            captured_configs.append(config)
            return original_run(config)

        with patch.object(pipeline_module, 'run_pipeline', side_effect=spy_run_pipeline):
            resp = client.post(
                "/api/decks/create",
                data={"outline_text": "# Test\n## Slide 1: Hi\n- OK\n", "enhance": "false"},
                headers={"Authorization": "Bearer tok"},
            )

        assert resp.status_code == 200
        assert len(captured_configs) == 1
        config = captured_configs[0]
        assert config.outline_path is not None
        assert config.outline_path.endswith(".md")


class TestApiTagsEndpoint:
    """GET /api/tags should return tags with counts and categories."""

    def test_returns_empty_when_no_tags(self, client):
        resp = client.get("/api/tags")
        assert resp.status_code == 200
        assert resp.json() == {"tags": []}

    def test_returns_tags_with_counts(self, client, db_path):
        add_tags(1, ["security", "architecture"], db_path=db_path)
        resp = client.get("/api/tags")
        assert resp.status_code == 200
        data = resp.json()
        assert "tags" in data
        names = {t["name"] for t in data["tags"]}
        assert "security" in names
        assert "architecture" in names
        for t in data["tags"]:
            assert t["count"] >= 1

    def test_includes_taxonomy_category(self, client, db_path):
        add_taxonomy_tags([{"name": "security", "category": "topic"}], db_path)
        add_tags(1, ["security"], db_path=db_path)
        resp = client.get("/api/tags")
        data = resp.json()
        sec = next(t for t in data["tags"] if t["name"] == "security")
        assert sec["category"] == "topic"

    def test_uncategorized_tags_have_empty_category(self, client, db_path):
        add_tags(1, ["misc-tag"], db_path=db_path)
        resp = client.get("/api/tags")
        data = resp.json()
        tag = next(t for t in data["tags"] if t["name"] == "misc-tag")
        assert tag["category"] == ""


# ---------------------------------------------------------------------------
# R6: /api/decks/create needs the same R1/R2 hardening as /api/decks/upload.
#
# Without these, the create endpoint silently:
#   - Accepts requests with no Microsoft token (Bearer-less)
#   - Lands a deck with no images when the Graph render path fails
#   - Routes every browser user's renders through the same SP folder
#     (because NTID is never threaded)
# ---------------------------------------------------------------------------


class TestCreateDeckHardening:
    """Same Bearer/NTID/fail-loud contract as /api/decks/upload (R1+R2+R4)."""

    @pytest.fixture
    def template_setup(self, tmp_path, monkeypatch):
        from pptx import Presentation
        template_path = str(tmp_path / "template.pptx")
        prs = Presentation()
        prs.save(template_path)
        config_path = str(tmp_path / "templates.yaml")
        (tmp_path / "templates.yaml").write_text(
            f"default_template: {template_path}\n")
        monkeypatch.setattr(
            "aippt.config.DEFAULT_TEMPLATE_CONFIG_PATH", config_path
        )
        return template_path

    def test_missing_bearer_returns_401(self, client, template_setup):
        resp = client.post(
            "/api/decks/create",
            data={"outline_text": "# T\n## S1\n- x\n", "enhance": "false"},
        )
        assert resp.status_code == 401
        assert "sign-in" in resp.json()["error"].lower() or \
               "microsoft" in resp.json()["error"].lower()

    def test_basic_auth_is_rejected(self, client, template_setup):
        resp = client.post(
            "/api/decks/create",
            data={"outline_text": "# T\n## S1\n- x\n", "enhance": "false"},
            headers={"Authorization": "Basic dXNlcjpwYXNz"},
        )
        assert resp.status_code == 401

    @patch("aippt.web.routes.ingest_deck")
    def test_ntid_header_is_threaded_to_ingest(
        self, mock_ingest, client, template_setup,
    ):
        mock_ingest.return_value = {
            "deck_id": 1, "deck_name": "d", "slide_count": 1,
            "images_exported": True, "tags_generated": False,
            "source_tracked": False,
        }
        resp = client.post(
            "/api/decks/create",
            data={"outline_text": "# T\n## S1\n- x\n", "enhance": "false"},
            headers={
                "Authorization": "Bearer tok",
                "X-AIPPT-NTID": "melliott",
            },
        )
        assert resp.status_code == 200
        _ = resp.text  # drain SSE
        assert mock_ingest.call_args.kwargs.get("ntid") == "melliott"
        assert mock_ingest.call_args.kwargs.get("ms_token") == "tok"

    @patch("aippt.web.routes.ingest_deck")
    def test_create_requires_images_on_linux(
        self, mock_ingest, client, template_setup,
    ):
        import sys
        mock_ingest.return_value = {
            "deck_id": 1, "deck_name": "d", "slide_count": 1,
            "images_exported": True, "tags_generated": False,
            "source_tracked": False,
        }
        resp = client.post(
            "/api/decks/create",
            data={"outline_text": "# T\n## S1\n- x\n", "enhance": "false"},
            headers={"Authorization": "Bearer tok"},
        )
        assert resp.status_code == 200
        _ = resp.text
        if sys.platform.startswith("linux"):
            assert mock_ingest.call_args.kwargs.get("require_images") is True

    @patch("aippt.web.routes.ingest_deck")
    def test_create_graph_error_emits_sse_error_event(
        self, mock_ingest, client, template_setup,
    ):
        from aippt import graph
        mock_ingest.side_effect = graph.GraphError(
            401, "InvalidAuthenticationToken", "Token expired",
        )
        resp = client.post(
            "/api/decks/create",
            data={"outline_text": "# T\n## S1\n- x\n", "enhance": "false"},
            headers={"Authorization": "Bearer tok"},
        )
        assert resp.status_code == 200  # SSE always 200
        body = resp.text
        assert "event: error" in body
        assert "Token expired" in body
        assert "event: complete" not in body
        # The JS sign-out hook (handleCreateEvent) keys off data.status === 401
        # to detect in-band auth failures. If this assertion fails the client
        # silently keeps a dead token in localStorage.
        import json as _json
        error_payload = None
        for block in body.split("\n\n"):
            if "event: error" not in block:
                continue
            for line in block.splitlines():
                if line.startswith("data: "):
                    error_payload = _json.loads(line[len("data: "):])
                    break
            if error_payload is not None:
                break
        assert error_payload is not None, "no SSE error block found"
        assert error_payload.get("status") == 401, (
            "SSE error event missing status: 401 — the JS 401 hook in "
            f"handleCreateEvent won't fire. Got: {error_payload!r}"
        )

    @patch("aippt.web.routes.ingest_deck")
    def test_upload_stream_graph_error_emits_sse_error_event(
        self, mock_ingest, client,
    ):
        """upload-stream must emit ``{status: <code>}`` for in-band Graph
        errors. R9 regression guard for the second SSE endpoint — the JS
        handler in handleUploadEvent keys off the same field.
        """
        from aippt import graph
        mock_ingest.side_effect = graph.GraphError(
            401, "InvalidAuthenticationToken", "Token expired",
        )
        resp = client.post(
            "/api/decks/upload-stream",
            files={"file": ("d.pptx", b"PK\x03\x04not-really",
                            "application/vnd.openxmlformats-officedocument."
                            "presentationml.presentation")},
            headers={"Authorization": "Bearer tok"},
        )
        assert resp.status_code == 200  # SSE always 200
        body = resp.text
        assert "event: error" in body
        import json as _json
        error_payload = None
        for block in body.split("\n\n"):
            if "event: error" not in block:
                continue
            for line in block.splitlines():
                if line.startswith("data: "):
                    error_payload = _json.loads(line[len("data: "):])
                    break
            if error_payload is not None:
                break
        assert error_payload is not None, "no SSE error block found"
        assert error_payload.get("status") == 401, (
            "upload-stream SSE error event missing status: 401 — the JS 401 "
            f"hook in handleUploadEvent won't fire. Got: {error_payload!r}"
        )


# ---------------------------------------------------------------------------
# R10: server-side NTID validation.
#
# X-AIPPT-NTID is interpolated into a SharePoint path. An attacker (or a
# user with a sloppy NTID field) can submit '/', '\\', ':', '..', etc. and
# either escape the per-user folder or produce invalid Graph URLs that 5xx
# deep inside the render pipeline. Validate at the edge with an allowlist
# and reject before any Graph call is attempted.
# ---------------------------------------------------------------------------


class TestNtidValidation:
    """X-AIPPT-NTID must match ^[A-Za-z0-9._-]+$ when present."""

    @pytest.fixture
    def template_setup(self, tmp_path, monkeypatch):
        from pptx import Presentation
        template_path = str(tmp_path / "template.pptx")
        prs = Presentation()
        prs.save(template_path)
        config_path = str(tmp_path / "templates.yaml")
        (tmp_path / "templates.yaml").write_text(
            f"default_template: {template_path}\n")
        monkeypatch.setattr(
            "aippt.config.DEFAULT_TEMPLATE_CONFIG_PATH", config_path
        )
        return template_path

    @pytest.mark.parametrize("bad", [
        "../etc",       # traversal
        "foo/bar",      # path separator
        "foo\\bar",     # windows separator
        "foo:bar",      # SharePoint URL delimiter
        "foo bar",      # whitespace (would 400 deeper in Graph anyway)
        "foo*",         # glob
        "foo?",         # query meta
        "",             # empty after strip — handled separately
    ])
    def test_create_rejects_invalid_ntid(self, client, template_setup, bad):
        # Empty string is "no NTID supplied" — anonymous fallback, not 400.
        if bad == "":
            return
        resp = client.post(
            "/api/decks/create",
            data={"outline_text": "# T\n## S1\n- x\n", "enhance": "false"},
            headers={
                "Authorization": "Bearer tok",
                "X-AIPPT-NTID": bad,
            },
        )
        assert resp.status_code == 400, (
            f"NTID {bad!r} should be rejected with 400, got {resp.status_code}"
        )
        assert "ntid" in resp.json()["error"].lower()

    @patch("aippt.web.routes.ingest_deck")
    def test_create_accepts_typical_ntid(
        self, mock_ingest, client, template_setup,
    ):
        """Standard NTIDs (alnum + . _ -) must still pass."""
        mock_ingest.return_value = {
            "deck_id": 1, "deck_name": "d", "slide_count": 1,
            "images_exported": True, "tags_generated": False,
            "source_tracked": False,
        }
        for good in ["melliott", "user.name", "user_42", "user-42", "U1"]:
            resp = client.post(
                "/api/decks/create",
                data={"outline_text": "# T\n## S1\n- x\n",
                      "enhance": "false"},
                headers={
                    "Authorization": "Bearer tok",
                    "X-AIPPT-NTID": good,
                },
            )
            assert resp.status_code == 200, (
                f"NTID {good!r} should be accepted, got {resp.status_code}"
            )

    def test_upload_stream_rejects_invalid_ntid(self, client, template_setup):
        # Use any non-empty bytes — the request should be rejected before
        # the PPTX parser even runs.
        resp = client.post(
            "/api/decks/upload-stream",
            files={"file": ("d.pptx", b"PK\x03\x04not-really",
                            "application/vnd.openxmlformats-officedocument."
                            "presentationml.presentation")},
            headers={
                "Authorization": "Bearer tok",
                "X-AIPPT-NTID": "../escape",
            },
        )
        assert resp.status_code == 400
        assert "ntid" in resp.json()["error"].lower()

    def test_upload_rejects_invalid_ntid(self, client, template_setup):
        resp = client.post(
            "/api/decks/upload",
            files={"file": ("d.pptx", b"PK\x03\x04not-really",
                            "application/vnd.openxmlformats-officedocument."
                            "presentationml.presentation")},
            headers={
                "Authorization": "Bearer tok",
                "X-AIPPT-NTID": "foo/bar",
            },
        )
        assert resp.status_code == 400
        assert "ntid" in resp.json()["error"].lower()

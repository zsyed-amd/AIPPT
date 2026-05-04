"""Integration tests for web API display_name behavior."""
import os
import uuid

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from aippt.catalog import catalog_deck
from aippt.web.app import create_app


@pytest.fixture
def uuid_deck_path(tmp_path):
    """Create a PPTX saved with a UUID-prefixed filename."""
    prs = Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Slide One"
    prefix = uuid.uuid4().hex  # 32 lowercase hex chars
    filename = f"{prefix}_My Presentation.pptx"
    path = str(tmp_path / filename)
    prs.save(path)
    return path


@pytest.fixture
def client(tmp_path, uuid_deck_path):
    """TestClient with a UUID-prefixed deck cataloged."""
    db_path = str(tmp_path / "test.db")
    uploads_dir = str(tmp_path / "uploads")
    catalog_deck(uuid_deck_path, db_path=db_path)
    app = create_app(db_path=db_path, uploads_dir=uploads_dir)
    return TestClient(app)


class TestHealthzEndpoint:
    """GET /healthz should return status ok."""

    def test_healthz_returns_ok(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        uploads_dir = str(tmp_path / "uploads")
        app = create_app(db_path=db_path, uploads_dir=uploads_dir)
        client = TestClient(app)
        resp = client.get("/healthz")
        assert resp.status_code == 200
        assert resp.json() == {"status": "ok"}


class TestDisplayNameInApi:
    """GET /api/decks should include display_name with UUID stripped."""

    def test_decks_response_has_display_name(self, client):
        resp = client.get("/api/decks")
        assert resp.status_code == 200
        decks = resp.json()
        assert len(decks) == 1
        deck = decks[0]
        assert "display_name" in deck
        assert deck["display_name"] == "My Presentation"
        # Raw name should still contain the UUID prefix
        assert len(deck["name"]) > len(deck["display_name"])

    def test_download_content_disposition_stripped(self, client):
        resp = client.get("/api/decks/1/download")
        assert resp.status_code == 200
        cd = resp.headers.get("content-disposition", "")
        assert "My Presentation.pptx" in cd
        # Should NOT contain the 32-char hex prefix
        assert len(cd.split("_")[0].replace('attachment; filename="', '')) < 32

    def test_search_results_have_display_deck_name(self, client):
        resp = client.get("/api/search?title=Slide")
        assert resp.status_code == 200
        results = resp.json()
        assert len(results) >= 1
        assert results[0]["display_deck_name"] == "My Presentation"


class TestDocsRoute:
    """Tests for the /docs static mount."""

    def test_docs_served_when_built(self, tmp_path):
        """StaticFiles mount serves docs HTML correctly."""
        from fastapi import FastAPI
        from fastapi.staticfiles import StaticFiles

        docs_html_dir = tmp_path / "html"
        docs_html_dir.mkdir()
        (docs_html_dir / "index.html").write_text("<html><body>Docs</body></html>")

        app = FastAPI()
        app.mount("/docs", StaticFiles(directory=str(docs_html_dir), html=True), name="docs")
        client = TestClient(app)
        resp = client.get("/docs/index.html")
        assert resp.status_code == 200
        assert "Docs" in resp.text

    def test_app_starts_without_docs(self, tmp_path):
        """App starts and serves normally regardless of docs build state."""
        db_path = str(tmp_path / "test.db")
        uploads_dir = str(tmp_path / "uploads")
        app = create_app(db_path=db_path, uploads_dir=uploads_dir)
        client = TestClient(app)
        # Root page should always work
        resp = client.get("/")
        assert resp.status_code == 200
        # Static assets should always work
        resp = client.get("/static/index.html")
        assert resp.status_code == 200

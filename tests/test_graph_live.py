"""Live Microsoft Graph integration tests.

Run with:
    MS_ACCESS_TOKEN=$(pbpaste) \\
    AIPPT_SP_SITE_ID=contoso.sharepoint.com,... \\
    AIPPT_SP_DRIVE_ID=b!... \\
    venv/bin/python -m pytest tests/test_graph_live.py -m live -v

Requires:
  - MS_ACCESS_TOKEN: a valid Graph access token (Bearer, Files.ReadWrite.All
    + Sites.ReadWrite.All). Obtain via the web UI's device-code flow and
    copy the access_token out of the localStorage entry.
  - Network egress to graph.microsoft.com.
  - For the render test: AIPPT_SP_SITE_ID + AIPPT_SP_DRIVE_ID pointing at a
    library you can write to. The test uploads, downloads as PDF, then
    deletes the staged file.

These tests are excluded from the default pytest run (`-m 'not e2e and not
live'` in pyproject). They exist so we can prove against a real tenant that
the stdlib HTTP client agrees with Graph's wire format.
"""
from __future__ import annotations

import os
import tempfile
from pathlib import Path

import pytest

from aippt import graph, render


pytestmark = pytest.mark.live


MS_TOKEN = os.environ.get("MS_ACCESS_TOKEN", "").strip()
SP_SITE_ID = os.environ.get("AIPPT_SP_SITE_ID", "").strip()
SP_DRIVE_ID = os.environ.get("AIPPT_SP_DRIVE_ID", "").strip()

SKIP_NO_TOKEN = pytest.mark.skipif(
    not MS_TOKEN, reason="MS_ACCESS_TOKEN not set"
)
SKIP_NO_SP = pytest.mark.skipif(
    not (MS_TOKEN and SP_SITE_ID and SP_DRIVE_ID),
    reason="MS_ACCESS_TOKEN + AIPPT_SP_SITE_ID + AIPPT_SP_DRIVE_ID required",
)


@SKIP_NO_TOKEN
def test_can_call_me_endpoint():
    """Smoke test: token is valid and Graph round-trips JSON we can parse."""
    me = graph.get_json("/me", token=graph.get_token_from_env() or MS_TOKEN)
    # Graph guarantees `id` on /me responses; everything else is optional.
    assert "id" in me, f"Unexpected /me response shape: {me!r}"


@SKIP_NO_SP
def test_render_small_fixture_produces_pngs(tmp_path: Path):
    """End-to-end: tiny PPTX -> Graph -> PDF -> pdftoppm PNGs.

    Builds the fixture in-process (python-pptx) so we don't depend on
    example assets that may move.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    pptx_path = tmp_path / "live-fixture.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only
    slide.shapes.title.text = "AIPPT live render check"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    box.text_frame.text = "If you can see this, Graph rendered the PDF."
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    prs.save(str(pptx_path))

    # NTID just has to be a valid SharePoint folder name; pick something
    # obvious so leftover folders are easy to spot if cleanup ever fails.
    ntid = "aippt-livetest"

    out_dir = tmp_path / "pngs"
    pngs = render.render_pptx_to_pngs(
        pptx_path=str(pptx_path),
        out_dir=str(out_dir),
        token=MS_TOKEN,
        ntid=ntid,
        site_id=SP_SITE_ID,
        drive_id=SP_DRIVE_ID,
    )

    assert pngs, "render returned no PNGs"
    for png in pngs:
        assert png.exists(), f"missing PNG: {png}"
        assert png.stat().st_size > 0, f"empty PNG: {png}"
        # PNG signature: 0x89 P N G
        assert png.read_bytes()[:4] == b"\x89PNG", (
            f"not a PNG: {png}"
        )

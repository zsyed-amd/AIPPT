"""Tests for aippt/template_merge.py — corporate template merge."""

import json
import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt


def _make_slide_with_metadata(prs, layout_selected: str | None = None, notes_text: str = ""):
    """Helper: create a slide with optional [AIPPT-META] layout_selected."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.7))
    txBox.text_frame.text = "Test Slide"

    if layout_selected is not None:
        entries = [{"operation": "enhance", "layout_selected": layout_selected}]
        meta_block = f"[AIPPT-META]\n{json.dumps(entries, indent=2)}\n[/AIPPT-META]"
        if notes_text:
            slide.notes_slide.notes_text_frame.text = f"{notes_text}\n\n---\n{meta_block}"
        else:
            slide.notes_slide.notes_text_frame.text = meta_block
    elif notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    return slide


class TestFindLayoutByName:
    """Test _find_layout_by_name() against the real corporate template."""

    def test_find_existing_layout(self):
        from aippt.template_merge import _find_layout_by_name

        prs = Presentation("templates/corp.pptx")
        layout = _find_layout_by_name(prs, "Title and Content")
        assert layout is not None
        assert layout.name == "Title and Content"

    def test_find_blank_layout(self):
        from aippt.template_merge import _find_layout_by_name

        prs = Presentation("templates/corp.pptx")
        layout = _find_layout_by_name(prs, "Blank")
        assert layout is not None
        assert layout.name == "Blank"

    def test_find_nonexistent_layout(self):
        from aippt.template_merge import _find_layout_by_name

        prs = Presentation("templates/corp.pptx")
        layout = _find_layout_by_name(prs, "Does Not Exist")
        assert layout is None

    def test_all_corp_layout_map_entries_exist(self):
        from aippt.template_merge import _find_layout_by_name, CORP_LAYOUT_MAP

        prs = Presentation("templates/corp.pptx")
        for layout_type, layout_name in CORP_LAYOUT_MAP.items():
            layout = _find_layout_by_name(prs, layout_name)
            assert layout is not None, f"Layout '{layout_name}' (for type '{layout_type}') not found in corp template"


class TestGetLayoutForSlide:
    """Test _get_layout_for_slide() metadata-to-layout resolution."""

    def test_mapped_layout_type(self):
        from aippt.template_merge import _get_layout_for_slide, CORP_LAYOUT_MAP

        src_prs = Presentation()
        slide = _make_slide_with_metadata(src_prs, layout_selected="bullet")

        target_prs = Presentation("templates/corp.pptx")
        layout, source_type, target_name = _get_layout_for_slide(
            slide, CORP_LAYOUT_MAP, "Blank", target_prs
        )
        assert source_type == "bullet"
        assert target_name == "Title and Content"
        assert layout.name == "Title and Content"

    def test_unmapped_layout_type_falls_back(self):
        from aippt.template_merge import _get_layout_for_slide, CORP_LAYOUT_MAP

        src_prs = Presentation()
        slide = _make_slide_with_metadata(src_prs, layout_selected="iconRows")

        target_prs = Presentation("templates/corp.pptx")
        layout, source_type, target_name = _get_layout_for_slide(
            slide, CORP_LAYOUT_MAP, "Blank", target_prs
        )
        assert source_type == "iconRows"
        assert target_name == "Blank"

    def test_no_metadata_falls_back(self):
        from aippt.template_merge import _get_layout_for_slide, CORP_LAYOUT_MAP

        src_prs = Presentation()
        slide = _make_slide_with_metadata(src_prs, layout_selected=None, notes_text="Just notes")

        target_prs = Presentation("templates/corp.pptx")
        layout, source_type, target_name = _get_layout_for_slide(
            slide, CORP_LAYOUT_MAP, "Blank", target_prs
        )
        assert source_type is None
        assert target_name == "Blank"

    def test_no_notes_at_all_falls_back(self):
        from aippt.template_merge import _get_layout_for_slide, CORP_LAYOUT_MAP

        src_prs = Presentation()
        slide = src_prs.slides.add_slide(src_prs.slide_layouts[0])

        target_prs = Presentation("templates/corp.pptx")
        layout, source_type, target_name = _get_layout_for_slide(
            slide, CORP_LAYOUT_MAP, "Blank", target_prs
        )
        assert source_type is None
        assert target_name == "Blank"

    def test_each_corp_layout_map_entry(self):
        from aippt.template_merge import _get_layout_for_slide, CORP_LAYOUT_MAP

        target_prs = Presentation("templates/corp.pptx")
        for layout_type, expected_name in CORP_LAYOUT_MAP.items():
            src_prs = Presentation()
            slide = _make_slide_with_metadata(src_prs, layout_selected=layout_type)
            layout, source_type, target_name = _get_layout_for_slide(
                slide, CORP_LAYOUT_MAP, "Blank", target_prs
            )
            assert target_name == expected_name, f"Type '{layout_type}' should map to '{expected_name}', got '{target_name}'"


class TestCopySlideToTemplate:
    """Test _copy_slide_to_template() preserves shapes and background."""

    def test_shapes_preserved(self):
        from aippt.template_merge import _copy_slide_to_template, _find_layout_by_name

        src_prs = Presentation()
        slide = src_prs.slides.add_slide(src_prs.slide_layouts[0])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        slide.shapes[-1].text_frame.text = "Hello World"
        slide.shapes.add_shape(1, Inches(2), Inches(3), Inches(2), Inches(1))
        original_shape_texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]

        target_prs = Presentation("templates/corp.pptx")
        blank = _find_layout_by_name(target_prs, "Blank")
        dest = _copy_slide_to_template(slide, target_prs, blank)

        dest_texts = [s.text_frame.text for s in dest.shapes if s.has_text_frame]
        for text in original_shape_texts:
            assert text in dest_texts, f"Text '{text}' not found in copied slide"

    def test_placeholder_shapes_removed(self):
        from aippt.template_merge import _copy_slide_to_template, _find_layout_by_name

        src_prs = Presentation()
        slide = src_prs.slides.add_slide(src_prs.slide_layouts[0])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        slide.shapes[-1].text_frame.text = "Content"

        target_prs = Presentation("templates/corp.pptx")
        content_layout = _find_layout_by_name(target_prs, "Title and Content")
        dest = _copy_slide_to_template(slide, target_prs, content_layout)

        placeholders = [s for s in dest.shapes if s.is_placeholder]
        assert len(placeholders) == 0, f"Found {len(placeholders)} placeholder shapes that should have been removed"


class TestTransferNotes:
    """Test _transfer_notes() preserves notes and appends history."""

    def test_notes_with_metadata_transferred(self):
        from aippt.template_merge import _transfer_notes

        src_prs = Presentation()
        slide = _make_slide_with_metadata(src_prs, layout_selected="bullet", notes_text="Speaker notes here")

        dest_prs = Presentation()
        dest_slide = dest_prs.slides.add_slide(dest_prs.slide_layouts[0])

        _transfer_notes(slide, dest_slide, "Title and Content")

        dest_notes = dest_slide.notes_slide.notes_text_frame.text
        assert "Speaker notes here" in dest_notes
        assert "[AIPPT-META]" in dest_notes
        assert "layout_selected" in dest_notes
        assert "/template-merge" in dest_notes

    def test_notes_without_metadata_transferred(self):
        from aippt.template_merge import _transfer_notes

        src_prs = Presentation()
        slide = _make_slide_with_metadata(src_prs, layout_selected=None, notes_text="Plain notes")

        dest_prs = Presentation()
        dest_slide = dest_prs.slides.add_slide(dest_prs.slide_layouts[0])

        _transfer_notes(slide, dest_slide, "Blank")

        dest_notes = dest_slide.notes_slide.notes_text_frame.text
        assert "Plain notes" in dest_notes

    def test_empty_notes_no_crash(self):
        from aippt.template_merge import _transfer_notes

        src_prs = Presentation()
        slide = src_prs.slides.add_slide(src_prs.slide_layouts[0])

        dest_prs = Presentation()
        dest_slide = dest_prs.slides.add_slide(dest_prs.slide_layouts[0])

        _transfer_notes(slide, dest_slide, "Blank")


class TestMergeWithTemplate:
    """Test the public merge_with_template() API."""

    def _make_generated_deck(self, path, slide_specs):
        """Create a generated deck with slides that have metadata.

        slide_specs: list of (title, layout_selected_or_None) tuples.
        """
        prs = Presentation()
        for title, layout_type in slide_specs:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.7))
            txBox.text_frame.text = title

            if layout_type is not None:
                entries = [{"operation": "enhance", "layout_selected": layout_type}]
                meta_block = f"[AIPPT-META]\n{json.dumps(entries, indent=2)}\n[/AIPPT-META]"
                slide.notes_slide.notes_text_frame.text = f"Notes for {title}\n\n---\n{meta_block}"
            else:
                slide.notes_slide.notes_text_frame.text = f"Notes for {title}"

        prs.save(path)

    def test_merge_basic(self, tmp_path):
        from aippt.template_merge import merge_with_template

        gen_path = str(tmp_path / "generated.pptx")
        out_path = str(tmp_path / "merged.pptx")
        self._make_generated_deck(gen_path, [
            ("Title Slide", "title"),
            ("Content Slide", "bullet"),
            ("Code Slide", "code"),
        ])

        result = merge_with_template(gen_path, "templates/corp.pptx", out_path)

        assert result["slide_count"] == 3
        assert result["output_path"] == out_path
        assert len(result["layout_assignments"]) == 3
        assert os.path.isfile(out_path)

        merged_prs = Presentation(out_path)
        assert len(merged_prs.slides) == 3

    def test_layout_assignments_correct(self, tmp_path):
        from aippt.template_merge import merge_with_template

        gen_path = str(tmp_path / "generated.pptx")
        out_path = str(tmp_path / "merged.pptx")
        self._make_generated_deck(gen_path, [
            ("Title", "title"),
            ("Bullets", "bullet"),
            ("Two Col", "two_column"),
            ("Code", "code"),
            ("Divider", "section_divider"),
            ("Close", "closing"),
        ])

        result = merge_with_template(gen_path, "templates/corp.pptx", out_path)

        expected = [
            ("Title Slide - No Image", "title"),
            ("Title and Content", "bullet"),
            ("Two Content", "two_column"),
            ("Developer Code Layout", "code"),
            ("Divider slide", "section_divider"),
            ("Closing logo slide", "closing"),
        ]
        for assignment, (exp_layout, exp_source) in zip(result["layout_assignments"], expected):
            assert assignment["target_layout"] == exp_layout
            assert assignment["source_layout"] == exp_source

    def test_fallback_layout(self, tmp_path):
        from aippt.template_merge import merge_with_template

        gen_path = str(tmp_path / "generated.pptx")
        out_path = str(tmp_path / "merged.pptx")
        self._make_generated_deck(gen_path, [
            ("Unknown Layout", "iconRows"),
            ("No Metadata", None),
        ])

        result = merge_with_template(gen_path, "templates/corp.pptx", out_path)

        assert result["layout_assignments"][0]["target_layout"] == "Blank"
        assert result["layout_assignments"][1]["target_layout"] == "Blank"

    def test_metadata_preserved(self, tmp_path):
        from aippt.template_merge import merge_with_template
        from aippt.metadata import extract_metadata

        gen_path = str(tmp_path / "generated.pptx")
        out_path = str(tmp_path / "merged.pptx")
        self._make_generated_deck(gen_path, [("Test", "bullet")])

        merge_with_template(gen_path, "templates/corp.pptx", out_path)

        merged_prs = Presentation(out_path)
        entries = extract_metadata(merged_prs.slides[0])
        layout_types = [e.get("layout_selected") for e in entries if "layout_selected" in e]
        assert "bullet" in layout_types

        notes = merged_prs.slides[0].notes_slide.notes_text_frame.text
        assert "/template-merge" in notes

    def test_shapes_preserved(self, tmp_path):
        from aippt.template_merge import merge_with_template

        gen_path = str(tmp_path / "generated.pptx")
        out_path = str(tmp_path / "merged.pptx")
        self._make_generated_deck(gen_path, [("My Title", "bullet")])

        merge_with_template(gen_path, "templates/corp.pptx", out_path)

        merged_prs = Presentation(out_path)
        texts = [s.text_frame.text for s in merged_prs.slides[0].shapes if s.has_text_frame]
        assert "My Title" in texts

    def test_empty_deck(self, tmp_path):
        from aippt.template_merge import merge_with_template

        gen_path = str(tmp_path / "generated.pptx")
        out_path = str(tmp_path / "merged.pptx")
        prs = Presentation()
        prs.save(gen_path)

        result = merge_with_template(gen_path, "templates/corp.pptx", out_path)

        assert result["slide_count"] == 0
        assert result["layout_assignments"] == []

    def test_custom_layout_map(self, tmp_path):
        from aippt.template_merge import merge_with_template

        gen_path = str(tmp_path / "generated.pptx")
        out_path = str(tmp_path / "merged.pptx")
        self._make_generated_deck(gen_path, [("Test", "bullet")])

        custom_map = {"bullet": "Blank"}
        result = merge_with_template(gen_path, "templates/corp.pptx", out_path, layout_map=custom_map)

        assert result["layout_assignments"][0]["target_layout"] == "Blank"

    def test_missing_template_file(self, tmp_path):
        from aippt.template_merge import merge_with_template

        gen_path = str(tmp_path / "generated.pptx")
        prs = Presentation()
        prs.save(gen_path)

        with pytest.raises(FileNotFoundError):
            merge_with_template(gen_path, "nonexistent.pptx", str(tmp_path / "out.pptx"))

    def test_missing_generated_file(self, tmp_path):
        from aippt.template_merge import merge_with_template

        with pytest.raises(FileNotFoundError):
            merge_with_template("nonexistent.pptx", "templates/corp.pptx", str(tmp_path / "out.pptx"))

    def test_corporate_master_present(self, tmp_path):
        from aippt.template_merge import merge_with_template

        gen_path = str(tmp_path / "generated.pptx")
        out_path = str(tmp_path / "merged.pptx")
        self._make_generated_deck(gen_path, [("Test", "bullet")])

        merge_with_template(gen_path, "templates/corp.pptx", out_path)

        merged_prs = Presentation(out_path)
        layout_names = [l.name for l in merged_prs.slide_layouts]
        assert "Title Slide - No Image" in layout_names
        assert "Developer Code Layout" in layout_names
        assert "Closing logo slide" in layout_names
        assert len(layout_names) >= 31


class TestMergeTemplateCLI:
    """Test the merge-template CLI subcommand."""

    def _make_generated_deck(self, path, slide_specs):
        """Create a generated deck with slides that have metadata."""
        prs = Presentation()
        for title, layout_type in slide_specs:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.7))
            txBox.text_frame.text = title
            if layout_type is not None:
                entries = [{"operation": "enhance", "layout_selected": layout_type}]
                meta_block = f"[AIPPT-META]\n{json.dumps(entries, indent=2)}\n[/AIPPT-META]"
                slide.notes_slide.notes_text_frame.text = meta_block
        prs.save(path)

    def test_merge_template_subcommand(self, tmp_path):
        from aippt.cli import build_parser

        gen_path = str(tmp_path / "generated.pptx")
        out_path = str(tmp_path / "merged.pptx")
        self._make_generated_deck(gen_path, [("Test", "bullet")])

        parser = build_parser()
        args = parser.parse_args([
            "merge-template", gen_path,
            "--corp-template", "templates/corp.pptx",
            "-o", out_path,
        ])
        assert args.command == "merge-template"
        assert args.generated_pptx == gen_path
        assert args.corp_template == "templates/corp.pptx"
        assert args.output == out_path

    def test_merge_template_dry_run_flag(self, tmp_path):
        from aippt.cli import build_parser

        gen_path = str(tmp_path / "generated.pptx")
        self._make_generated_deck(gen_path, [("Test", "bullet")])

        parser = build_parser()
        args = parser.parse_args([
            "merge-template", gen_path,
            "--corp-template", "templates/corp.pptx",
            "-o", str(tmp_path / "out.pptx"),
            "--dry-run",
        ])
        assert args.dry_run is True

    def test_cmd_merge_template_runs(self, tmp_path):
        from aippt.cli import cmd_merge_template

        gen_path = str(tmp_path / "generated.pptx")
        out_path = str(tmp_path / "merged.pptx")
        self._make_generated_deck(gen_path, [("Test", "bullet")])

        class Args:
            generated_pptx = gen_path
            corp_template = "templates/corp.pptx"
            output = out_path
            layout_map = None
            dry_run = False

        result = cmd_merge_template(Args())
        assert result == 0
        assert os.path.isfile(out_path)

    def test_cmd_merge_template_dry_run(self, tmp_path, capsys):
        from aippt.cli import cmd_merge_template

        gen_path = str(tmp_path / "generated.pptx")
        self._make_generated_deck(gen_path, [("Test", "bullet")])

        class Args:
            generated_pptx = gen_path
            corp_template = "templates/corp.pptx"
            output = str(tmp_path / "out.pptx")
            layout_map = None
            dry_run = True

        result = cmd_merge_template(Args())
        assert result == 0
        assert not os.path.isfile(str(tmp_path / "out.pptx"))
        captured = capsys.readouterr()
        assert "bullet" in captured.out
        assert "Title and Content" in captured.out


class TestPipelineIntegration:
    """Test that PipelineConfig accepts corp_template field."""

    def test_pipeline_config_has_corp_template(self):
        from aippt.pipeline import PipelineConfig

        config = PipelineConfig(
            outline_text="## Test\n- bullet",
            template_path="templates/corp.pptx",
            output_path="/tmp/test.pptx",
            corp_template="templates/corp.pptx",
        )
        assert config.corp_template == "templates/corp.pptx"

    def test_pipeline_config_corp_template_default_none(self):
        from aippt.pipeline import PipelineConfig

        config = PipelineConfig(
            outline_text="## Test\n- bullet",
            template_path="templates/corp.pptx",
            output_path="/tmp/test.pptx",
        )
        assert config.corp_template is None

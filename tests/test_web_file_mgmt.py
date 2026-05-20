"""Tests for upload and download API endpoints."""
import io
import os
from unittest.mock import patch

import pytest
from pptx import Presentation
from starlette.testclient import TestClient

from aippt.web.app import create_app
from aippt.catalog import catalog_deck


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def make_pptx(num_slides: int = 1) -> io.BytesIO:
    """Return an in-memory PPTX with ``num_slides`` slides."""
    prs = Presentation()
    layout = prs.slide_layouts[0]
    for _ in range(num_slides):
        prs.slides.add_slide(layout)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def app(tmp_path):
    """Create a fresh app with isolated db and uploads_dir."""
    db_path = str(tmp_path / "test.db")
    uploads_dir = str(tmp_path / "uploads")
    # Force view_only=False so the upload endpoint isn't 403-gated by the
    # PRD-A view-only block. A test Bearer token is supplied in _upload.
    return create_app(db_path=db_path, uploads_dir=uploads_dir, view_only=False)


@pytest.fixture
def client(app):
    """Starlette TestClient wrapping the app."""
    return TestClient(app)


# Default Bearer header for every upload — PRD A requires Microsoft sign-in
# for the Linux Graph render path. Tests don't actually hit Graph (export is
# mocked) but the upload endpoint requires the header to clear its 401 gate.
_TEST_AUTH = {"Authorization": "Bearer test-ms-token"}


def _upload(client, filename="my_deck.pptx", num_slides=3, generate_tags=False):
    """Helper: upload a PPTX, optionally with generate_tags."""
    buf = make_pptx(num_slides=num_slides)
    data = {}
    if generate_tags:
        data["generate_tags"] = "true"
    return client.post(
        "/api/decks/upload",
        files={"file": (filename, buf, PPTX_MIME)},
        data=data,
        headers=_TEST_AUTH,
    )


# ---------------------------------------------------------------------------
# TestUploadDeck
# ---------------------------------------------------------------------------


class TestUploadDeck:

    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_upload_valid_pptx(self, _mock_export, client):
        """Upload a valid PPTX and verify the response fields and deck list."""
        response = _upload(client, num_slides=3)

        assert response.status_code == 200, response.text
        data = response.json()

        # Required fields present
        assert "id" in data
        assert "name" in data
        assert "slide_count" in data
        assert "message" in data
        assert "images_exported" in data
        assert "tags_generated" in data

        # Slide count matches what we put in
        assert data["slide_count"] == 3

        # Image export succeeded (mocked)
        assert data["images_exported"] is True
        assert data["tags_generated"] is False

        # Deck appears in the list endpoint
        list_response = client.get("/api/decks")
        assert list_response.status_code == 200
        decks = list_response.json()
        deck_ids = [d["id"] for d in decks]
        assert data["id"] in deck_ids

    def test_upload_invalid_file_type(self, client):
        """Uploading a non-.pptx file must return 400 with an error message."""
        buf = io.BytesIO(b"this is not a presentation")
        response = client.post(
            "/api/decks/upload",
            files={"file": ("notes.txt", buf, "text/plain")},
        )

        assert response.status_code == 400, response.text
        data = response.json()
        assert "error" in data

    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_upload_duplicate(self, _mock_export, client):
        """Uploading the same PPTX bytes twice should succeed both times."""
        buf = make_pptx(num_slides=2)
        pptx_bytes = buf.read()

        # First upload
        resp1 = client.post(
            "/api/decks/upload",
            files={"file": ("deck.pptx", io.BytesIO(pptx_bytes), PPTX_MIME)},
            headers=_TEST_AUTH,
        )
        assert resp1.status_code == 200, resp1.text
        id1 = resp1.json()["id"]

        # Second upload (identical content)
        resp2 = client.post(
            "/api/decks/upload",
            files={"file": ("deck.pptx", io.BytesIO(pptx_bytes), PPTX_MIME)},
            headers=_TEST_AUTH,
        )
        assert resp2.status_code == 200, resp2.text
        id2 = resp2.json()["id"]

        # Both should have returned valid deck IDs (existing deck returned)
        assert id1 > 0
        assert id2 > 0

    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_upload_with_images_exported(self, _mock_export, client):
        """When image export succeeds, response reflects images_exported=True."""
        response = _upload(client, num_slides=2)

        assert response.status_code == 200, response.text
        data = response.json()
        assert data["images_exported"] is True
        assert data["tags_generated"] is False

    @patch("aippt.ingest.cmd_analyze", return_value=0)
    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_upload_with_generate_tags(self, _mock_export, _mock_analyze, client):
        """When generate_tags is true, tags are generated."""
        response = _upload(client, num_slides=2, generate_tags=True)

        assert response.status_code == 200, response.text
        data = response.json()
        assert data["images_exported"] is True
        assert data["tags_generated"] is True
        assert "tags generated" in data["message"]

    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_upload_without_generate_tags_flag(self, _mock_export, client):
        """Default upload without generate_tags should not generate tags."""
        response = _upload(client, num_slides=1)

        assert response.status_code == 200, response.text
        data = response.json()
        assert data["tags_generated"] is False


# ---------------------------------------------------------------------------
# TestDownloadDeck
# ---------------------------------------------------------------------------


class TestDownloadDeck:

    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def _upload_pptx(self, client, _mock_export, filename: str = "test_deck.pptx", num_slides: int = 1):
        """Helper: upload a PPTX and return the parsed response JSON."""
        buf = make_pptx(num_slides=num_slides)
        resp = client.post(
            "/api/decks/upload",
            files={"file": (filename, buf, PPTX_MIME)},
            headers=_TEST_AUTH,
        )
        assert resp.status_code == 200, resp.text
        return resp.json()

    def test_download_existing_deck(self, client):
        """Upload then download; verify status, Content-Type, and non-empty body."""
        deck = self._upload_pptx(client)
        deck_id = deck["id"]

        response = client.get(f"/api/decks/{deck_id}/download")

        assert response.status_code == 200, response.text
        assert response.headers["content-type"] == PPTX_MIME
        assert len(response.content) > 0

    def test_download_nonexistent_deck(self, client):
        """Requesting a deck ID that does not exist must return 404."""
        response = client.get("/api/decks/99999/download")

        assert response.status_code == 404, response.text

    def test_download_missing_file(self, client, app):
        """Catalog a deck, delete its file from disk, then attempt download -> 404."""
        deck = self._upload_pptx(client, filename="vanishing.pptx")
        deck_id = deck["id"]

        # Retrieve the file path from the db directly
        from aippt.catalog import get_deck_by_id
        db_path = app.state.db_path
        deck_record = get_deck_by_id(deck_id, db_path)
        assert deck_record is not None
        file_path = deck_record["file_path"]
        assert os.path.exists(file_path)

        # Remove the file to simulate a missing source file
        os.remove(file_path)
        assert not os.path.exists(file_path)

        # Download should now return 404 with appropriate message
        response = client.get(f"/api/decks/{deck_id}/download")
        assert response.status_code == 404, response.text
        data = response.json()
        assert "error" in data
        assert "Source file not found" in data["error"]


# ---------------------------------------------------------------------------
# SSE helper
# ---------------------------------------------------------------------------


def _parse_sse(text: str) -> list:
    """Parse SSE text into a list of {"event": ..., "data": ...} dicts.

    Each SSE message is separated by a blank line.  Lines beginning with
    ``event:`` set the event name; lines beginning with ``data:`` carry the
    payload (parsed as JSON when possible).
    """
    import json

    messages = []
    for block in text.strip().split("\n\n"):
        if not block.strip():
            continue
        event = None
        data = None
        for line in block.splitlines():
            if line.startswith("event:"):
                event = line[len("event:"):].strip()
            elif line.startswith("data:"):
                raw = line[len("data:"):].strip()
                try:
                    data = json.loads(raw)
                except (ValueError, TypeError):
                    data = raw
        if event is not None or data is not None:
            messages.append({"event": event, "data": data})
    return messages


# ---------------------------------------------------------------------------
# TestUploadStream
# ---------------------------------------------------------------------------


class TestUploadStream:
    """Tests for POST /api/decks/upload-stream (SSE endpoint)."""

    def _stream_upload(self, client, filename="stream_deck.pptx", num_slides=3, generate_tags=False):
        """Helper: upload via the SSE endpoint and return the raw response."""
        buf = make_pptx(num_slides=num_slides)
        data = {}
        if generate_tags:
            data["generate_tags"] = "true"
        return client.post(
            "/api/decks/upload-stream",
            files={"file": (filename, buf, PPTX_MIME)},
            data=data,
            headers=_TEST_AUTH,
        )

    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_stream_returns_sse_events(self, _mock_export, client):
        """Upload via SSE endpoint; verify content-type and presence of key events."""
        response = self._stream_upload(client)

        assert response.status_code == 200, response.text
        assert "text/event-stream" in response.headers.get("content-type", "")

        events = _parse_sse(response.text)
        event_names = [e["event"] for e in events]

        # export_images progress events should be present
        assert "progress" in event_names
        # complete event must arrive at the end
        assert "complete" in event_names
        assert event_names[-1] == "complete"

        # Verify export_images step appears (running or skipped)
        progress_events = [e for e in events if e["event"] == "progress"]
        steps = [e["data"].get("step") for e in progress_events if isinstance(e["data"], dict)]
        assert "export_images" in steps

        # Catalog step must also appear
        assert "catalog" in steps

    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_stream_complete_has_deck_info(self, _mock_export, client):
        """The complete SSE event must carry deck_id, deck_name, and slide_count."""
        response = self._stream_upload(client, num_slides=4)

        assert response.status_code == 200, response.text

        events = _parse_sse(response.text)
        complete_events = [e for e in events if e["event"] == "complete"]
        assert len(complete_events) == 1, "Expected exactly one complete event"

        data = complete_events[0]["data"]
        assert isinstance(data, dict)
        assert "deck_id" in data
        assert "deck_name" in data
        assert "slide_count" in data
        assert data["deck_id"] > 0
        assert data["slide_count"] == 4

    @patch("aippt.ingest.cmd_analyze", return_value=0)
    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_stream_with_tags(self, _mock_export, _mock_analyze, client):
        """When generate_tags=true, a tags progress step must appear in the SSE stream."""
        response = self._stream_upload(client, num_slides=2, generate_tags=True)

        assert response.status_code == 200, response.text

        events = _parse_sse(response.text)
        progress_events = [e for e in events if e["event"] == "progress"]
        steps = [e["data"].get("step") for e in progress_events if isinstance(e["data"], dict)]
        assert "tags" in steps

        complete_events = [e for e in events if e["event"] == "complete"]
        assert len(complete_events) == 1
        complete_data = complete_events[0]["data"]
        assert complete_data.get("tags_generated") is True
        assert complete_data.get("images_exported") is True

    def test_stream_rejects_non_pptx(self, client):
        """Non-.pptx upload must return 400 JSON error (not SSE)."""
        buf = io.BytesIO(b"not a presentation")
        response = client.post(
            "/api/decks/upload-stream",
            files={"file": ("notes.txt", buf, "text/plain")},
        )

        assert response.status_code == 400, response.text
        data = response.json()
        assert "error" in data

    def test_stream_without_bearer_returns_401(self, client):
        """Missing Bearer token must return 401 JSON (not SSE)."""
        buf = make_pptx(num_slides=1)
        response = client.post(
            "/api/decks/upload-stream",
            files={"file": ("deck.pptx", buf, PPTX_MIME)},
        )

        assert response.status_code == 401, response.text
        data = response.json()
        assert "sign-in" in data["error"].lower() or "microsoft" in data["error"].lower()

    def test_stream_in_view_only_returns_403(self, tmp_path):
        """View-only deployments must block the SSE upload endpoint with 403."""
        db_path = str(tmp_path / "test.db")
        uploads_dir = str(tmp_path / "uploads")
        app = create_app(db_path=db_path, uploads_dir=uploads_dir, view_only=True)
        view_only_client = TestClient(app)

        buf = make_pptx(num_slides=1)
        response = view_only_client.post(
            "/api/decks/upload-stream",
            files={"file": ("deck.pptx", buf, PPTX_MIME)},
            headers=_TEST_AUTH,
        )

        assert response.status_code == 403, response.text
        data = response.json()
        assert "view-only" in data["error"].lower()

"""Unit tests for the in-memory log ring buffer + /api/logs route."""
from __future__ import annotations

import logging
import os

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from aippt.catalog import catalog_deck
from aippt.web.app import create_app
from aippt.web.log_buffer import RingBufferLogHandler, install_ring_buffer


# ---------------------------------------------------------------------------
# RingBufferLogHandler — pure unit behavior, no FastAPI
# ---------------------------------------------------------------------------


class TestRingBufferHandler:
    def test_emit_captures_message_and_metadata(self):
        h = RingBufferLogHandler(capacity=10)
        rec = logging.LogRecord(
            name="aippt.test", level=logging.INFO, pathname=__file__,
            lineno=1, msg="hello %s", args=("world",), exc_info=None,
        )
        h.emit(rec)
        snap = h.snapshot()
        assert len(snap) == 1
        assert snap[0]["message"] == "hello world"
        assert snap[0]["level"] == "INFO"
        assert snap[0]["logger"] == "aippt.test"
        assert isinstance(snap[0]["id"], int)
        assert isinstance(snap[0]["ts"], float)

    def test_ring_drops_oldest_when_capacity_exceeded(self):
        h = RingBufferLogHandler(capacity=3)
        for i in range(5):
            h.emit(logging.LogRecord(
                "x", logging.INFO, __file__, 1, "msg-%d" % i, None, None,
            ))
        snap = h.snapshot()
        # Only the most recent 3 survive.
        assert [r["message"] for r in snap] == ["msg-2", "msg-3", "msg-4"]

    def test_snapshot_filters_by_level(self):
        h = RingBufferLogHandler(capacity=10)
        for lvl, msg in [
            (logging.DEBUG, "d"),
            (logging.INFO, "i"),
            (logging.WARNING, "w"),
            (logging.ERROR, "e"),
        ]:
            h.emit(logging.LogRecord(
                "x", lvl, __file__, 1, msg, None, None,
            ))
        assert [r["message"] for r in h.snapshot(level="WARNING")] == ["w", "e"]
        assert [r["message"] for r in h.snapshot(level="ERROR")] == ["e"]

    def test_snapshot_filters_by_logger_prefix(self):
        h = RingBufferLogHandler(capacity=10)
        for name in ("aippt.render", "uvicorn.access", "aippt.web"):
            h.emit(logging.LogRecord(
                name, logging.INFO, __file__, 1, "x", None, None,
            ))
        loggers = [r["logger"] for r in h.snapshot(logger_prefix="aippt")]
        assert loggers == ["aippt.render", "aippt.web"]

    def test_snapshot_since_returns_only_new(self):
        h = RingBufferLogHandler(capacity=10)
        for i in range(3):
            h.emit(logging.LogRecord(
                "x", logging.INFO, __file__, 1, "m%d" % i, None, None,
            ))
        first = h.snapshot()
        cursor = first[-1]["id"]
        h.emit(logging.LogRecord(
            "x", logging.INFO, __file__, 1, "after-cursor", None, None,
        ))
        new = h.snapshot(since=cursor)
        assert [r["message"] for r in new] == ["after-cursor"]

    def test_bad_format_args_does_not_crash_emit(self):
        h = RingBufferLogHandler(capacity=5)
        # %d format with a string arg would raise inside getMessage().
        h.emit(logging.LogRecord(
            "x", logging.INFO, __file__, 1, "n=%d", ("not-an-int",), None,
        ))
        snap = h.snapshot()
        assert len(snap) == 1
        # The handler swallowed the format error; some representation made it
        # through rather than crashing the producer.
        assert snap[0]["message"]


class TestInstallRingBuffer:
    def teardown_method(self):
        # Clean the root logger so tests don't leak handlers into each other.
        root = logging.getLogger()
        root.handlers = [
            h for h in root.handlers if not isinstance(h, RingBufferLogHandler)
        ]

    def test_idempotent_install_does_not_stack_handlers(self):
        first = install_ring_buffer(capacity=50)
        second = install_ring_buffer(capacity=50)
        assert first is second, (
            "calling install_ring_buffer twice in one process must reuse "
            "the existing handler; otherwise every log record gets "
            "double-counted"
        )
        ring_handlers = [
            h for h in logging.getLogger().handlers
            if isinstance(h, RingBufferLogHandler)
        ]
        assert len(ring_handlers) == 1

    def test_install_captures_app_logger_emissions(self):
        h = install_ring_buffer(capacity=50)
        h.clear()
        logging.getLogger("aippt.test.installed").info("captured-line")
        msgs = [r["message"] for r in h.snapshot(logger_prefix="aippt")]
        assert "captured-line" in msgs

    def test_install_attaches_to_uvicorn_loggers(self):
        # uvicorn.access and uvicorn.error have propagate=False by default,
        # so attaching only to root would silently drop the HTTP access log
        # — the most useful thing a developer would expect to see.
        h = install_ring_buffer()
        for name in ("uvicorn.access", "uvicorn.error"):
            assert h in logging.getLogger(name).handlers, (
                f"ring buffer must be attached to {name} or HTTP traffic "
                "will not appear in /api/logs"
            )

    def test_install_raises_root_level_above_warning_default(self):
        # uvicorn boot sets root to WARNING which would drop INFO; the
        # installer must lower it so the ring buffer is actually useful.
        root = logging.getLogger()
        original_level = root.level
        try:
            root.setLevel(logging.NOTSET)
            install_ring_buffer()
            assert root.level <= logging.INFO
        finally:
            root.setLevel(original_level)


# ---------------------------------------------------------------------------
# /api/logs route — integration through FastAPI TestClient
# ---------------------------------------------------------------------------


@pytest.fixture
def deck_path(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "T"
    p = str(tmp_path / "t.pptx")
    prs.save(p)
    return p


@pytest.fixture
def client(tmp_path, deck_path):
    db_path = str(tmp_path / "logs.db")
    catalog_deck(deck_path, db_path=db_path)
    app = create_app(
        db_path=db_path, uploads_dir=str(tmp_path / "u"),
        images_dir=str(tmp_path / "img"),
    )
    return TestClient(app)


class TestLogsRoute:
    def test_requires_bearer_token(self, client):
        resp = client.get("/api/logs")
        assert resp.status_code == 401
        assert "sign-in" in resp.json()["error"].lower()

    def test_returns_recent_records_with_cursor(self, client):
        # Emit a couple of unique lines we can search for in the response.
        logging.getLogger("aippt.test.routes").warning(
            "test-marker-AAA")
        logging.getLogger("aippt.test.routes").error(
            "test-marker-BBB")

        resp = client.get(
            "/api/logs?logger_prefix=aippt.test.routes",
            headers={"Authorization": "Bearer tok"},
        )
        assert resp.status_code == 200
        body = resp.json()
        msgs = [r["message"] for r in body["records"]]
        assert "test-marker-AAA" in msgs
        assert "test-marker-BBB" in msgs
        assert body["next_cursor"] is not None
        # Polling with the cursor returns nothing new yet.
        poll = client.get(
            f"/api/logs?since={body['next_cursor']}"
            f"&logger_prefix=aippt.test.routes",
            headers={"Authorization": "Bearer tok"},
        ).json()
        assert poll["count"] == 0

        # After another emit, polling picks up only the new line.
        logging.getLogger("aippt.test.routes").info("test-marker-CCC")
        poll2 = client.get(
            f"/api/logs?since={body['next_cursor']}"
            f"&logger_prefix=aippt.test.routes",
            headers={"Authorization": "Bearer tok"},
        ).json()
        assert [r["message"] for r in poll2["records"]] == ["test-marker-CCC"]

    def test_level_filter_drops_lower_records(self, client):
        logging.getLogger("aippt.test.level").info("LOW")
        logging.getLogger("aippt.test.level").error("HIGH")
        resp = client.get(
            "/api/logs?level=ERROR&logger_prefix=aippt.test.level",
            headers={"Authorization": "Bearer tok"},
        )
        msgs = [r["message"] for r in resp.json()["records"]]
        assert "HIGH" in msgs
        assert "LOW" not in msgs

    def test_limit_clamped_to_capacity(self, client):
        resp = client.get(
            "/api/logs?limit=999999",
            headers={"Authorization": "Bearer tok"},
        )
        body = resp.json()
        # Limit clamped to handler capacity; never errors.
        assert body["count"] <= body["capacity"]

    def test_bearer_token_is_scrubbed_in_captured_logs(self, client):
        # If any handler logs the Authorization header verbatim, the
        # pre-existing scrub filter on the root logger must redact it
        # before the ring buffer sees the line.
        secret = "Bearer ey-very-real-token-xxx"
        logging.getLogger("aippt.test.scrub").warning(
            "request headers={'Authorization': '%s'}", secret,
        )
        resp = client.get(
            "/api/logs?logger_prefix=aippt.test.scrub",
            headers={"Authorization": "Bearer tok"},
        )
        all_text = " ".join(r["message"] for r in resp.json()["records"])
        assert "ey-very-real-token-xxx" not in all_text, (
            "ring buffer captured an unscrubbed Bearer token — the scrub "
            "filter on the root logger must run before this handler"
        )
        assert "Bearer <redacted>" in all_text

    def test_works_in_view_only_mode(self, tmp_path, deck_path):
        # Logs are read-only triage; view-only deployments still need them.
        db_path = str(tmp_path / "vo.db")
        catalog_deck(deck_path, db_path=db_path)
        app = create_app(
            db_path=db_path, uploads_dir=str(tmp_path / "u"),
            images_dir=str(tmp_path / "img"), view_only=True,
        )
        vo_client = TestClient(app)
        resp = vo_client.get(
            "/api/logs",
            headers={"Authorization": "Bearer tok"},
        )
        assert resp.status_code == 200

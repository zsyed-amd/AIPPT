"""Test slide master support in pptxgenjs-generated PPTX files.

Generates decks with and without masters and validates the PPTX
structure using python-pptx.
"""

import os
import subprocess
import tempfile

import pytest
from pptx import Presentation

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
NODE_PATH = subprocess.check_output(
    ["npm", "root", "-g"], cwd=REPO_ROOT, text=True
).strip()


def _generate_deck(use_master: bool, output_path: str) -> str:
    """Generate a minimal pptxgenjs deck with or without masters."""
    master_flag = "true" if use_master else "false"
    helpers_path = os.path.join(REPO_ROOT, "lib", "pptxgenjs-helpers.mjs")
    theme_path = os.path.join(REPO_ROOT, "themes", "default.yaml")
    script = f"""\
import {{ createDeck, addTitleSlide, addBulletSlide, addSectionDivider,
         addClosingSlide, addProcessFlow, addCardGrid, SW, SH
       }} from '{helpers_path}';

const deck = createDeck('{theme_path}', {{ useSlideMaster: {master_flag} }});

addTitleSlide(deck, 'Test Title', 'Subtitle', 1);
addBulletSlide(deck, 'Bullets', ['Item 1', 'Item 2', 'Item 3'], 2, 'notes');
addSectionDivider(deck, 1, 'Section One', 3);
addProcessFlow(deck, 'Process', ['Step 1', 'Step 2', 'Step 3'], 4, 'notes');
addCardGrid(deck, 'Cards', [
  {{ title: 'A', body: 'Desc A' }},
  {{ title: 'B', body: 'Desc B' }},
], 5, 'notes');
addClosingSlide(deck, 6, 'Thank you');

await deck.save('{output_path}');
"""
    script_path = os.path.join(tempfile.gettempdir(), "test_master_gen.mjs")
    with open(script_path, "w") as f:
        f.write(script)

    env = {**os.environ, "NODE_PATH": NODE_PATH}
    result = subprocess.run(
        ["node", script_path],
        cwd=REPO_ROOT,
        env=env,
        capture_output=True,
        text=True,
        timeout=30,
    )
    assert result.returncode == 0, f"Generation failed:\n{result.stderr}"
    assert os.path.exists(output_path), f"Output not created: {output_path}"
    return output_path


class TestMasterDefinition:
    """Generated PPTX has expected masters when useSlideMaster is true."""

    def test_master_deck_has_named_layouts(self, tmp_path):
        out = str(tmp_path / "master.pptx")
        _generate_deck(use_master=True, output_path=out)
        p = Presentation(out)
        layout_names = [sl.name for sm in p.slide_masters for sl in sm.slide_layouts]
        assert "TITLE_MASTER" in layout_names, f"Missing TITLE_MASTER in {layout_names}"
        assert "CONTENT_MASTER" in layout_names, f"Missing CONTENT_MASTER in {layout_names}"
        assert "SECTION_DIVIDER_MASTER" in layout_names, f"Missing SECTION_DIVIDER_MASTER in {layout_names}"

    def test_no_master_deck_has_default_master_only(self, tmp_path):
        out = str(tmp_path / "no_master.pptx")
        _generate_deck(use_master=False, output_path=out)
        p = Presentation(out)
        assert len(p.slide_masters) == 1

    def test_both_decks_have_same_slide_count(self, tmp_path):
        out_m = str(tmp_path / "master.pptx")
        out_n = str(tmp_path / "no_master.pptx")
        _generate_deck(use_master=True, output_path=out_m)
        _generate_deck(use_master=False, output_path=out_n)
        p_m = Presentation(out_m)
        p_n = Presentation(out_n)
        assert len(p_m.slides) == len(p_n.slides) == 6


class TestBackwardCompat:
    """Default useSlideMaster=false produces identical structure to today."""

    def test_default_creates_deck_without_masters(self, tmp_path):
        out = str(tmp_path / "default.pptx")
        _generate_deck(use_master=False, output_path=out)
        p = Presentation(out)
        assert len(p.slide_masters) == 1
        assert len(p.slides) == 6

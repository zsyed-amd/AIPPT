"""Tests for view-only mode (library mode)."""
import os
import pytest
from unittest.mock import patch
from fastapi.testclient import TestClient
from pptx import Presentation

from aippt.catalog import catalog_deck
from aippt.web.app import create_app, detect_view_only


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def deck_path(tmp_path):
    """Create a minimal PPTX with one slide."""
    prs = Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Test Slide"
    path = str(tmp_path / "test.pptx")
    prs.save(path)
    return path


@pytest.fixture
def view_only_client(tmp_path, deck_path):
    """TestClient with view_only=True."""
    db_path = str(tmp_path / "test.db")
    uploads_dir = str(tmp_path / "uploads")
    catalog_deck(deck_path, db_path=db_path)
    app = create_app(db_path=db_path, uploads_dir=uploads_dir, view_only=True)
    return TestClient(app)


@pytest.fixture
def full_client(tmp_path, deck_path):
    """TestClient with view_only=False."""
    db_path = str(tmp_path / "test.db")
    uploads_dir = str(tmp_path / "uploads")
    catalog_deck(deck_path, db_path=db_path)
    app = create_app(db_path=db_path, uploads_dir=uploads_dir, view_only=False)
    return TestClient(app)


# ---------------------------------------------------------------------------
# Config endpoint
# ---------------------------------------------------------------------------

class TestConfigEndpoint:
    """GET /api/config returns view_only flag."""

    def test_config_endpoint_full_mode(self, full_client):
        resp = full_client.get("/api/config")
        assert resp.status_code == 200
        assert resp.json() == {"view_only": False}

    def test_config_endpoint_view_only(self, view_only_client):
        resp = view_only_client.get("/api/config")
        assert resp.status_code == 200
        assert resp.json() == {"view_only": True}


# ---------------------------------------------------------------------------
# 403 guards on LLM endpoints
# ---------------------------------------------------------------------------

class TestViewOnly403Guards:
    """LLM endpoints return 403 in view-only mode."""

    def test_analyze_blocked_view_only(self, view_only_client):
        resp = view_only_client.post("/api/slides/1/analyze", json={})
        assert resp.status_code == 403
        assert "view-only" in resp.json()["error"]

    def test_notes_blocked_view_only(self, view_only_client):
        resp = view_only_client.post("/api/slides/1/notes", json={})
        assert resp.status_code == 403
        assert "view-only" in resp.json()["error"]

    def test_improvements_blocked_view_only(self, view_only_client):
        resp = view_only_client.post("/api/slides/1/improvements", json={})
        assert resp.status_code == 403
        assert "view-only" in resp.json()["error"]

    def test_create_blocked_view_only(self, view_only_client):
        resp = view_only_client.post(
            "/api/decks/create",
            data={"outline_text": "# Test\n- bullet"},
        )
        assert resp.status_code == 403
        assert "view-only" in resp.json()["error"]


# ---------------------------------------------------------------------------
# Non-LLM endpoints still work in view-only mode
# ---------------------------------------------------------------------------

class TestViewOnlyAllowedEndpoints:
    """Browse, search, and upload still work in view-only mode."""

    def test_search_works_view_only(self, view_only_client):
        resp = view_only_client.get("/api/search?tags=&title=")
        assert resp.status_code == 200

    def test_decks_list_works_view_only(self, view_only_client):
        resp = view_only_client.get("/api/decks")
        assert resp.status_code == 200
        assert len(resp.json()) >= 1

    def test_upload_blocked_view_only(self, view_only_client, deck_path):
        # Per PRD A: deck ingest now requires Microsoft sign-in for the
        # Linux Graph render path, so view-only deployments cannot ingest.
        with open(deck_path, "rb") as f:
            resp = view_only_client.post(
                "/api/decks/upload",
                files={"file": ("upload.pptx", f, "application/octet-stream")},
                data={"generate_tags": "true"},
                headers={"Authorization": "Bearer tok"},
            )
        assert resp.status_code == 403
        assert "view-only" in resp.json()["error"].lower()


# ---------------------------------------------------------------------------
# Auto-detection logic
# ---------------------------------------------------------------------------

class TestAutoDetection:
    """detect_view_only() and auto-detection in create_app()."""

    def test_auto_detect_no_config(self, tmp_path):
        """No gateway.yaml, no env vars → view-only."""
        with patch.dict(os.environ, {}, clear=True):
            assert detect_view_only(str(tmp_path / "nonexistent.yaml")) is True

    def test_auto_detect_with_gateway(self, tmp_path):
        """Existing gateway.yaml → not view-only."""
        gw = tmp_path / "gateway.yaml"
        gw.write_text("gateway:\n  base_url: http://example.com\n")
        with patch.dict(os.environ, {}, clear=True):
            assert detect_view_only(str(gw)) is False

    def test_auto_detect_with_anthropic_key(self):
        """ANTHROPIC_API_KEY set → not view-only."""
        with patch.dict(os.environ, {"ANTHROPIC_API_KEY": "sk-test"}, clear=True):
            assert detect_view_only(None) is False

    def test_auto_detect_with_openai_key(self):
        """OPENAI_API_KEY set → not view-only."""
        with patch.dict(os.environ, {"OPENAI_API_KEY": "sk-test"}, clear=True):
            assert detect_view_only(None) is False

    def test_flag_overrides_detection(self, tmp_path):
        """--view-only flag forces True even with gateway present."""
        gw = tmp_path / "gateway.yaml"
        gw.write_text("gateway:\n  base_url: http://example.com\n")
        app = create_app(
            db_path=str(tmp_path / "test.db"),
            gateway_config=str(gw),
            view_only=True,
        )
        assert app.state.view_only is True

    def test_auto_detect_in_create_app(self, tmp_path):
        """create_app with view_only=None auto-detects."""
        with patch.dict(os.environ, {}, clear=True):
            app = create_app(
                db_path=str(tmp_path / "test.db"),
                gateway_config=str(tmp_path / "nonexistent.yaml"),
                view_only=None,
            )
            assert app.state.view_only is True

    def test_env_var_forces_view_only(self):
        """AIPPT_VIEW_ONLY=1 forces view-only even with API keys."""
        with patch.dict(os.environ, {"AIPPT_VIEW_ONLY": "1", "ANTHROPIC_API_KEY": "sk-test"}, clear=True):
            assert detect_view_only(None) is True

    def test_env_var_true_forces_view_only(self):
        """AIPPT_VIEW_ONLY=true forces view-only."""
        with patch.dict(os.environ, {"AIPPT_VIEW_ONLY": "true"}, clear=True):
            assert detect_view_only(None) is True

    def test_env_var_zero_does_not_force_view_only(self):
        """AIPPT_VIEW_ONLY=0 does not force view-only (falls to auto-detect)."""
        with patch.dict(os.environ, {"AIPPT_VIEW_ONLY": "0", "ANTHROPIC_API_KEY": "sk-test"}, clear=True):
            assert detect_view_only(None) is False

    def test_env_var_zero_without_keys(self):
        """AIPPT_VIEW_ONLY=0 explicitly disables view-only even without keys."""
        with patch.dict(os.environ, {"AIPPT_VIEW_ONLY": "0"}, clear=True):
            assert detect_view_only(None) is False

    def test_flag_overrides_env_var(self, tmp_path):
        """--view-only flag (view_only=True) overrides AIPPT_VIEW_ONLY=0."""
        with patch.dict(os.environ, {"AIPPT_VIEW_ONLY": "0"}, clear=True):
            app = create_app(
                db_path=str(tmp_path / "test.db"),
                view_only=True,
            )
            assert app.state.view_only is True

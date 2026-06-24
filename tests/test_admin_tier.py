"""Tests for the v1 admin tier: NTID allowlist, admin-gated endpoints, whoami."""
from __future__ import annotations

import base64
import json
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from aippt.catalog import catalog_deck
from aippt.config import load_admin_ntids
from aippt.web.app import create_app


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def deck_path(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "T"
    p = str(tmp_path / "t.pptx")
    prs.save(p)
    return p


@pytest.fixture
def app_with_admin(tmp_path, deck_path):
    db_path = str(tmp_path / "admin.db")
    catalog_deck(deck_path, db_path=db_path)
    app = create_app(
        db_path=db_path, uploads_dir=str(tmp_path / "u"),
        images_dir=str(tmp_path / "img"),
    )
    # Lowercase allowlist (as load_admin_ntids produces). Includes the six
    # NTIDs added alongside the case-insensitivity fix so they're exercised
    # end-to-end through the gate.
    app.state.admin_ntids = {
        "melliott", "jdoe",
        "ansgputa", "edtian", "egroenke", "miroy", "yrajesh", "zsyed",
    }
    return app


@pytest.fixture
def client(app_with_admin):
    return TestClient(app_with_admin)


def _fake_jwt(payload: dict) -> str:
    """Build a minimal unsigned JWT (header.payload.signature) for the
    Bearer-identity extraction tests. The signature is junk; the extractor
    is explicitly unverified."""
    def b64(d: dict) -> str:
        raw = json.dumps(d).encode()
        return base64.urlsafe_b64encode(raw).rstrip(b"=").decode()
    return f"{b64({'typ':'JWT'})}.{b64(payload)}.signature-not-verified"


# ---------------------------------------------------------------------------
# load_admin_ntids — config loader
# ---------------------------------------------------------------------------


class TestLoadAdminNtids:
    def test_empty_when_no_config(self):
        assert load_admin_ntids(None) == set()

    def test_empty_when_file_missing(self, tmp_path):
        assert load_admin_ntids(str(tmp_path / "nope.yaml")) == set()

    def test_empty_when_key_absent(self, tmp_path):
        cfg = tmp_path / "gw.yaml"
        cfg.write_text("gateway:\n  base_url: x\n", encoding="utf-8")
        assert load_admin_ntids(str(cfg)) == set()

    def test_reads_list(self, tmp_path):
        cfg = tmp_path / "gw.yaml"
        cfg.write_text(
            "admin_ntids:\n  - melliott\n  - jdoe\n  - alpha.beta\n",
            encoding="utf-8",
        )
        assert load_admin_ntids(str(cfg)) == {"melliott", "jdoe", "alpha.beta"}

    def test_drops_malformed_entries(self, tmp_path):
        cfg = tmp_path / "gw.yaml"
        cfg.write_text(
            "admin_ntids:\n  - melliott\n  - 'has space'\n  - bad@char\n  - 42\n",
            encoding="utf-8",
        )
        # Only the well-formed entry survives.
        assert load_admin_ntids(str(cfg)) == {"melliott"}

    def test_empty_list_is_no_admins(self, tmp_path):
        cfg = tmp_path / "gw.yaml"
        cfg.write_text("admin_ntids: []\n", encoding="utf-8")
        assert load_admin_ntids(str(cfg)) == set()

    def test_lowercases_entries(self, tmp_path):
        # Matching is case-insensitive: entries are lowercased at load.
        cfg = tmp_path / "gw.yaml"
        cfg.write_text(
            "admin_ntids:\n  - MElliott\n  - ZSYED\n", encoding="utf-8",
        )
        assert load_admin_ntids(str(cfg)) == {"melliott", "zsyed"}


# ---------------------------------------------------------------------------
# /api/auth/whoami
# ---------------------------------------------------------------------------


class TestWhoami:
    def test_signed_out_returns_no_identity(self, client):
        resp = client.get("/api/auth/whoami")
        assert resp.status_code == 200
        body = resp.json()
        assert body == {
            "signed_in": False, "ntid": "", "is_admin": False,
            "suggested_ntid": "",
        }

    def test_signed_in_non_admin(self, client):
        resp = client.get(
            "/api/auth/whoami",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "rando"},
        )
        body = resp.json()
        assert body["signed_in"] is True
        assert body["ntid"] == "rando"
        assert body["is_admin"] is False

    def test_signed_in_admin(self, client):
        resp = client.get(
            "/api/auth/whoami",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "melliott"},
        )
        body = resp.json()
        assert body["signed_in"] is True
        assert body["ntid"] == "melliott"
        assert body["is_admin"] is True

    def test_malformed_ntid_returns_blank_not_admin(self, client):
        # Server allowlist rejects path-traversal / spaces / etc.
        resp = client.get(
            "/api/auth/whoami",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "has space"},
        )
        body = resp.json()
        assert body["signed_in"] is True
        # Malformed → treated as empty; is_admin must be False.
        assert body["ntid"] == ""
        assert body["is_admin"] is False

    def test_mixed_case_ntid_is_admin(self, client):
        # Header arrives mixed-case; allowlist is lowercase. Case-insensitive
        # match → admin True. ntid echoes the original case the client sent
        # (audit/SharePoint fidelity), but the gate still recognizes it.
        resp = client.get(
            "/api/auth/whoami",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "MElliott"},
        )
        body = resp.json()
        assert body["ntid"] == "MElliott"
        assert body["is_admin"] is True

    def test_newly_added_ntid_is_admin(self, client):
        # One of the six NTIDs added with this change.
        resp = client.get(
            "/api/auth/whoami",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "zsyed"},
        )
        assert resp.json()["is_admin"] is True

    def test_near_miss_ntid_is_not_admin(self, client):
        # A near-miss of a real admin NTID must not match.
        resp = client.get(
            "/api/auth/whoami",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "zsyed1"},
        )
        assert resp.json()["is_admin"] is False


# ---------------------------------------------------------------------------
# Admin gate on DELETE /api/decks/{id}
# ---------------------------------------------------------------------------


class TestDeleteDeckAdminGate:
    def test_admin_can_delete(self, client):
        resp = client.delete(
            "/api/decks/1",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "melliott"},
        )
        assert resp.status_code == 200

    def test_admin_can_delete_mixed_case(self, client):
        # The gate (not just whoami) honors case-insensitive matching.
        resp = client.delete(
            "/api/decks/1",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "MElliott"},
        )
        assert resp.status_code == 200

    def test_non_admin_gets_403(self, client):
        resp = client.delete(
            "/api/decks/1",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "rando"},
        )
        assert resp.status_code == 403
        assert "admin" in resp.json()["error"].lower()

    def test_missing_ntid_gets_403(self, client):
        resp = client.delete(
            "/api/decks/1",
            headers={"Authorization": "Bearer tok"},
        )
        assert resp.status_code == 403

    def test_bearer_check_runs_first(self, client):
        # No Bearer → 401, not 403. The auth ladder is bearer → admin.
        resp = client.delete("/api/decks/1")
        assert resp.status_code == 401

    def test_view_only_check_runs_first(self, tmp_path, deck_path):
        db_path = str(tmp_path / "vo.db")
        catalog_deck(deck_path, db_path=db_path)
        app = create_app(
            db_path=db_path, uploads_dir=str(tmp_path / "u"),
            images_dir=str(tmp_path / "img"), view_only=True,
        )
        app.state.admin_ntids = {"melliott"}
        c = TestClient(app)
        resp = c.delete(
            "/api/decks/1",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "melliott"},
        )
        # Even for an admin, view_only wins — admin-tier doesn't bypass
        # deployment posture.
        assert resp.status_code == 403
        assert "view-only" in resp.json()["error"].lower()


# ---------------------------------------------------------------------------
# Admin gate on GET /api/logs
# ---------------------------------------------------------------------------


class TestLogsAdminGate:
    def test_admin_can_read(self, client):
        resp = client.get(
            "/api/logs",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "melliott"},
        )
        assert resp.status_code == 200

    def test_non_admin_gets_403(self, client):
        resp = client.get(
            "/api/logs",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "rando"},
        )
        assert resp.status_code == 403

    def test_bearer_check_runs_first(self, client):
        # No Bearer → 401 (sign-in required) before admin check kicks in.
        resp = client.get("/api/logs")
        assert resp.status_code == 401

    def test_view_only_allows_admin(self, tmp_path, deck_path):
        # /api/logs is read-only; view-only deployments still expose it
        # to admins.
        db_path = str(tmp_path / "vo.db")
        catalog_deck(deck_path, db_path=db_path)
        app = create_app(
            db_path=db_path, uploads_dir=str(tmp_path / "u"),
            images_dir=str(tmp_path / "img"), view_only=True,
        )
        app.state.admin_ntids = {"melliott"}
        c = TestClient(app)
        resp = c.get(
            "/api/logs",
            headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "melliott"},
        )
        assert resp.status_code == 200


# ---------------------------------------------------------------------------
# Audit logging for admin actions (impersonation paper trail)
# ---------------------------------------------------------------------------


class TestAdminAuditLog:
    def test_admin_action_logs_both_identities(self, client, caplog):
        # Bearer that decodes to upn=jdoe@amd.com; X-AIPPT-NTID claims melliott.
        # In v1 the gate trusts the header (melliott is in the allowlist) so the
        # action goes through, but the audit line records both so any later
        # investigation sees the discrepancy.
        jwt = _fake_jwt({"upn": "jdoe@amd.com", "oid": "guid"})
        with caplog.at_level("INFO", logger="aippt.web.routes"):
            resp = client.delete(
                "/api/decks/1",
                headers={"Authorization": f"Bearer {jwt}", "X-AIPPT-NTID": "melliott"},
            )
        assert resp.status_code == 200
        joined = " ".join(r.getMessage() for r in caplog.records)
        assert "admin_action" in joined
        assert "ntid=melliott" in joined
        assert "upn:jdoe@amd.com" in joined

    def test_denied_admin_action_logs_warning(self, client, caplog):
        with caplog.at_level("WARNING", logger="aippt.web.routes"):
            resp = client.delete(
                "/api/decks/1",
                headers={"Authorization": "Bearer tok", "X-AIPPT-NTID": "rando"},
            )
        assert resp.status_code == 403
        joined = " ".join(r.getMessage() for r in caplog.records)
        assert "admin_denied" in joined
        assert "ntid=rando" in joined


# ---------------------------------------------------------------------------
# Bearer-identity-unverified helper (audit only; never used to gate)
# ---------------------------------------------------------------------------


class TestBearerIdentityUnverified:
    def test_returns_absent_when_no_bearer(self, client):
        # Pull the helper through whoami's behavior: with no Bearer, the
        # audit machinery sees "absent". Easier to test via an admin-denial
        # log line than to import the helper directly.
        # Direct import is still the cleanest assertion:
        from aippt.web.routes import _bearer_identity_unverified
        from starlette.requests import Request as _Req

        scope = {"type": "http", "headers": [], "method": "GET", "path": "/"}
        req = _Req(scope)
        assert _bearer_identity_unverified(req) == "absent"

    def test_returns_unparseable_for_non_jwt(self):
        from aippt.web.routes import _bearer_identity_unverified
        from starlette.requests import Request as _Req

        scope = {
            "type": "http",
            "method": "GET",
            "path": "/",
            "headers": [(b"authorization", b"Bearer not-a-jwt")],
        }
        req = _Req(scope)
        assert _bearer_identity_unverified(req) == "unparseable"

    def test_prefers_upn_over_oid(self):
        from aippt.web.routes import _bearer_identity_unverified
        from starlette.requests import Request as _Req

        jwt = _fake_jwt({"oid": "guid-here", "upn": "jdoe@amd.com"})
        scope = {
            "type": "http",
            "method": "GET",
            "path": "/",
            "headers": [(b"authorization", f"Bearer {jwt}".encode())],
        }
        req = _Req(scope)
        assert _bearer_identity_unverified(req) == "upn:jdoe@amd.com"

    def test_falls_through_to_sub_when_no_upn(self):
        from aippt.web.routes import _bearer_identity_unverified
        from starlette.requests import Request as _Req

        jwt = _fake_jwt({"sub": "subject-id-only"})
        scope = {
            "type": "http",
            "method": "GET",
            "path": "/",
            "headers": [(b"authorization", f"Bearer {jwt}".encode())],
        }
        req = _Req(scope)
        assert _bearer_identity_unverified(req) == "sub:subject-id-only"


# ---------------------------------------------------------------------------
# suggested_ntid — Bearer-derived NTID hint (UX only; never used to gate)
# ---------------------------------------------------------------------------


class TestSuggestedNtid:
    def test_whoami_suggests_local_part_from_upn(self, client):
        # upn=MElliott@amd.com → local-part lowercased → "melliott".
        jwt = _fake_jwt({"upn": "MElliott@amd.com"})
        resp = client.get(
            "/api/auth/whoami",
            headers={"Authorization": f"Bearer {jwt}"},
        )
        assert resp.json()["suggested_ntid"] == "melliott"

    def test_whoami_prefers_preferred_username(self, client):
        # preferred_username wins over upn for the suggestion.
        jwt = _fake_jwt({
            "preferred_username": "ZSYED@amd.com",
            "upn": "someoneelse@amd.com",
        })
        resp = client.get(
            "/api/auth/whoami",
            headers={"Authorization": f"Bearer {jwt}"},
        )
        assert resp.json()["suggested_ntid"] == "zsyed"

    def test_whoami_suggestion_empty_without_claim(self, client):
        # A JWT with no usable identity claim → empty suggestion.
        jwt = _fake_jwt({"oid": "guid-only"})
        resp = client.get(
            "/api/auth/whoami",
            headers={"Authorization": f"Bearer {jwt}"},
        )
        assert resp.json()["suggested_ntid"] == ""

    def test_whoami_suggestion_empty_without_bearer(self, client):
        resp = client.get("/api/auth/whoami")
        assert resp.json()["suggested_ntid"] == ""

    def test_suggestion_does_not_grant_admin(self, client):
        # A Bearer whose identity is an admin, but no X-AIPPT-NTID header:
        # the suggestion is offered, yet is_admin stays False because the gate
        # trusts the explicit header, not the token-derived hint.
        jwt = _fake_jwt({"upn": "melliott@amd.com"})
        resp = client.get(
            "/api/auth/whoami",
            headers={"Authorization": f"Bearer {jwt}"},
        )
        body = resp.json()
        assert body["suggested_ntid"] == "melliott"
        assert body["ntid"] == ""
        assert body["is_admin"] is False

    def test_helper_rejects_malformed_local_part(self):
        # Local-part with a disallowed char (after split/lower) → "".
        from aippt.web.routes import _suggested_ntid_from_bearer
        from starlette.requests import Request as _Req

        jwt = _fake_jwt({"upn": "bad name@amd.com"})
        scope = {
            "type": "http",
            "method": "GET",
            "path": "/",
            "headers": [(b"authorization", f"Bearer {jwt}".encode())],
        }
        assert _suggested_ntid_from_bearer(_Req(scope)) == ""

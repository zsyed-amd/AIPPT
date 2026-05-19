"""Tests for the ``--images-dir`` plumbing through ``cmd_serve`` and the
web layer's ingest_deck calls.

Without this wiring, container deployments lose rendered PNGs on pod
restart: ``aippt serve`` previously had no way to override the default
``images/`` cwd-relative directory, so a Kubernetes PVC mounted at
``/app/data`` only captured uploads + the SQLite DB — not images.
"""
from __future__ import annotations

import argparse
import os
from unittest.mock import patch, ANY

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from aippt import cli
from aippt.catalog import catalog_deck
from aippt.web.app import create_app


@pytest.fixture
def deck_path(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test"
    p = str(tmp_path / "test.pptx")
    prs.save(p)
    return p


class TestServeArgParser:
    """``aippt serve --images-dir <path>`` must parse without error and
    the value must reach ``create_app(images_dir=...)``."""

    def test_images_dir_flows_from_cli_to_create_app(self, tmp_path):
        from aippt import cli as _cli
        images_target = str(tmp_path / "container-images")
        args = argparse.Namespace(
            host="127.0.0.1",
            port=8765,
            db="slides.db",
            gateway_config=None,
            uploads_dir=str(tmp_path / "up"),
            images_dir=images_target,
            view_only=False,
        )
        with patch("uvicorn.run") as _run, \
             patch("aippt.web.app.create_app") as mock_create_app:
            mock_create_app.return_value.state.view_only = False
            _cli.cmd_serve(args)
            kwargs = mock_create_app.call_args.kwargs
            assert kwargs.get("images_dir") == images_target, (
                f"--images-dir must flow to create_app(images_dir=...), got {kwargs!r}"
            )

    def test_images_dir_falls_back_to_dirs_yaml_when_unset(self, tmp_path):
        from aippt import cli as _cli
        args = argparse.Namespace(
            host="127.0.0.1",
            port=8765,
            db="slides.db",
            gateway_config=None,
            uploads_dir=str(tmp_path / "up"),
            images_dir=None,  # not provided on CLI
            view_only=False,
        )
        with patch("uvicorn.run") as _run, \
             patch("aippt.web.app.create_app") as mock_create_app:
            mock_create_app.return_value.state.view_only = False
            _cli.cmd_serve(args)
            kwargs = mock_create_app.call_args.kwargs
            # Falls back to dirs.yaml's images entry, resolved against base_dir.
            # We don't assert the exact value (depends on test cwd's dirs.yaml),
            # but it must be set and absolute (resolve_path returns absolute).
            assert kwargs.get("images_dir") is not None
            assert os.path.isabs(kwargs["images_dir"])


class TestCreateAppImagesDir:
    """``create_app`` must accept an ``images_dir`` arg and stash it on state."""

    def test_images_dir_lands_on_app_state(self, tmp_path):
        db = str(tmp_path / "x.db")
        uploads = str(tmp_path / "up")
        images = str(tmp_path / "img-base")
        app = create_app(
            db_path=db, uploads_dir=uploads, images_dir=images,
        )
        assert app.state.images_dir == images

    def test_images_dir_defaults_when_omitted(self, tmp_path):
        """Backward compat: existing tests construct create_app without
        images_dir. State should fall back to ``"images"`` (cwd-relative,
        the historical default)."""
        app = create_app(
            db_path=str(tmp_path / "x.db"),
            uploads_dir=str(tmp_path / "up"),
        )
        # The exact fallback must match what ingest_deck would have used.
        assert app.state.images_dir == "images"


class TestRoutesPassImagesDirToIngest:
    """All three ``ingest_deck`` call sites in routes.py must pass the
    app-state ``images_dir`` so the configured dir actually gets used."""

    @pytest.fixture
    def template_setup(self, tmp_path, monkeypatch):
        template_path = str(tmp_path / "template.pptx")
        Presentation().save(template_path)
        config_path = str(tmp_path / "templates.yaml")
        (tmp_path / "templates.yaml").write_text(
            f"default_template: {template_path}\n"
        )
        monkeypatch.setattr(
            "aippt.config.DEFAULT_TEMPLATE_CONFIG_PATH", config_path
        )
        return template_path

    @pytest.fixture
    def client(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        uploads_dir = str(tmp_path / "uploads")
        images_dir = str(tmp_path / "custom-images")
        # Catalog an empty deck so the app starts cleanly
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = str(tmp_path / "seed.pptx")
        prs.save(path)
        catalog_deck(path, db_path=db_path)
        app = create_app(
            db_path=db_path, uploads_dir=uploads_dir, images_dir=images_dir,
        )
        return TestClient(app), images_dir

    @patch("aippt.web.routes.ingest_deck")
    def test_upload_passes_app_state_images_dir(
        self, mock_ingest, client, deck_path,
    ):
        tc, images_dir = client
        mock_ingest.return_value = {
            "deck_id": 1, "deck_name": "d", "slide_count": 1,
            "images_exported": True, "tags_generated": False,
            "source_tracked": False, "images_dir": images_dir,
        }
        with open(deck_path, "rb") as fh:
            tc.post(
                "/api/decks/upload",
                files={"file": ("d.pptx", fh.read(),
                                "application/vnd.openxmlformats-officedocument."
                                "presentationml.presentation")},
                headers={"Authorization": "Bearer tok"},
            )
        assert mock_ingest.called, "ingest_deck was not invoked"
        kwargs = mock_ingest.call_args.kwargs
        passed = kwargs.get("images_dir")
        assert passed is not None, (
            "upload endpoint must thread app.state.images_dir to ingest_deck "
            "or rendered PNGs land in cwd (K8s data-loss bug)"
        )
        assert passed.startswith(images_dir), (
            f"images_dir should be rooted at app.state.images_dir, got {passed!r}"
        )

    @patch("aippt.web.routes.ingest_deck")
    def test_upload_stream_passes_app_state_images_dir(
        self, mock_ingest, client, deck_path,
    ):
        tc, images_dir = client
        mock_ingest.return_value = {
            "deck_id": 1, "deck_name": "d", "slide_count": 1,
            "images_exported": True, "tags_generated": False,
            "source_tracked": False, "images_dir": images_dir,
        }
        with open(deck_path, "rb") as fh:
            resp = tc.post(
                "/api/decks/upload-stream",
                files={"file": ("d.pptx", fh.read(),
                                "application/vnd.openxmlformats-officedocument."
                                "presentationml.presentation")},
                headers={"Authorization": "Bearer tok"},
            )
            _ = resp.text  # drain SSE
        assert mock_ingest.called
        passed = mock_ingest.call_args.kwargs.get("images_dir")
        assert passed and passed.startswith(images_dir), (
            f"upload-stream must use app.state.images_dir, got {passed!r}"
        )

    @patch("aippt.web.routes.ingest_deck")
    def test_create_passes_app_state_images_dir(
        self, mock_ingest, client, template_setup,
    ):
        tc, images_dir = client
        mock_ingest.return_value = {
            "deck_id": 1, "deck_name": "d", "slide_count": 1,
            "images_exported": True, "tags_generated": False,
            "source_tracked": False, "images_dir": images_dir,
        }
        resp = tc.post(
            "/api/decks/create",
            data={"outline_text": "# T\n## S1\n- x\n", "enhance": "false"},
            headers={"Authorization": "Bearer tok"},
        )
        _ = resp.text  # drain SSE
        assert mock_ingest.called
        passed = mock_ingest.call_args.kwargs.get("images_dir")
        assert passed and passed.startswith(images_dir), (
            f"create must use app.state.images_dir, got {passed!r}"
        )
